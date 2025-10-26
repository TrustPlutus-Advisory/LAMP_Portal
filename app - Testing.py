import dash
from dash import dcc, html, Input, Output, State, callback_context, ALL
import dash_bootstrap_components as dbc
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import os
import random
import string
import hashlib
from datetime import datetime, timedelta
import base64
import io
import json
import numpy as np
import io
from datetime import datetime
import base64
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import plotly.graph_objects as go
import kaleido

# Initialize the Dash app
app = dash.Dash(__name__, suppress_callback_exceptions=True)
server = app.server
app.title = "LAMP - Large Account Management Process"

# Premium CSS with TRUSTPLUTUS Branding
premium_css = """
@import url("https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap");
:root {
  --maroon: #b04d59;
  --maroon-light: #e07094;
  --maroon-dark: #cc6270;
  --lamp-light: rgba(255, 215, 0, 0.7);
  --gold: #f7ce7c;
  --gold-light: #e6daa1;
  --gold-dark: #e0d077;

  --gray-light: #F5F5F5;
  --gray-dark: #3A3A3A;
  --white: #FFFFFF;
  --platinum-gray: #E5E4E2;
  --charcoal-gray: #36454F;
  --slate-gray: #708090;
  --light-gray: #F8F9FA;
  --medium-gray: #6C757D;
  --dark-gray: #343A40;
  --pure-white: #FFFFFF;
  --text-dark: #2C3E50;
  --success-green: #28A745;
  --danger-red: #d95965;
  --warning-amber: #f2d683;
  --info-blue: #a2aab3;
  --shadow-light: rgba(204, 98, 112, 0.1);
  --shadow-medium: rgba(204, 98, 112, 0.2);
  --shadow-heavy: rgba(204, 98, 112,O.3);
  --shadow-sm: 0 4px 10px rgba(0,0,0,0.05);
  --shadow-md: 0 6px 20px rgba(0,0,0,0.08);
}
# :root {
#     --maroon: #800020;
#     --maroon-light: #a0203a;
#     --maroon-dark: #5c0015;
#     --gold: #FFD700;
#     --gold-light: #FFED4E;
#     --gold-dark: #B8860B;
#     --lamp-light: rgba(255, 215, 0, 0.7);
#     --platinum-gray: #E5E4E2;
#     --charcoal-gray: #36454F;
#     --slate-gray: #708090;
#     --light-gray: #F8F9FA;
#     --medium-gray: #6C757D;
#     --dark-gray: #343A40;
#     --pure-white: #FFFFFF;
#     --text-dark: #2C3E50;
#     --success-green: #28A745;
#     --danger-red: #DC3545;
#     --warning-amber: #FFC107;
#     --info-blue: #007BFF;
#     --shadow-light: rgba(128, 0, 32, 0.1);
#     --shadow-medium: rgba(128, 0, 32, 0.2);
#     --shadow-heavy: rgba(128, 0, 32, 0.3);
# }

* { margin: 0; padding: 0; box-sizing: border-box; }

body {
    font-family: 'Inter', -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif !important;
    background: linear-gradient(135deg, var(--light-gray) 0%, var(--platinum-gray) 100%) !important;
    color: var(--text-dark) !important;
    line-height: 1.6;
    min-height: 100vh;
    font-weight: 400;
    overflow-x: hidden;
}

.navbar-brand {
    font-family: 'Inter', sans-serif !important;
    font-size: 1.8rem !important;
    font-weight: 700 !important;
    letter-spacing: -0.02em;
    display: flex !important;
    align-items: center !important;
    justify-content: center !important;
}

.trust-text { color: var(--maroon) !important; }
.plutus-text { color: var(--gold-dark) !important; }
.lamp-text {
    color: var(--lamp-light) !important;
    font-size: 1.8rem !important;
    font-weight: 700 !important;
    margin-left: 1rem !important;
}

.main-navbar {
    background: var(--pure-white) !important;
    border-bottom: 1px solid var(--platinum-gray) !important;
    box-shadow: 0 2px 20px var(--shadow-light) !important;
    backdrop-filter: blur(10px);
    padding: 1rem 0 !important;
}

.nav-btn {
    background: transparent !important;
    border: 2px solid var(--pure-white) !important;
    color: var(--text-dark) !important;
    padding: 0.5rem 1.2rem !important;  /* Reduced from 0.6rem 1.5rem */
    border-radius: 50px !important;
    font-weight: 600 !important;
    transition: all 0.3s ease !important;
    margin: 0 0.5rem !important;
    font-size: 0.813rem !important;  /* Reduced from 0.9rem (13px instead of 14.4px) */
}

.nav-btn:hover {
    background: var(--pure-white) !important;
    color: var(--text-dark) !important;
    transform: translateY(-2px) !important;
    box-shadow: 0 8px 25px var(--pure-white) !important;
}

.nav-btn-primary {
    background: linear-gradient(135deg, var(--pure-white) 0%, var(--pure-white) 100%) !important;
    color: var(--text-dark) !important;
    border: none !important;
}

.hero-section {
    background: linear-gradient(135deg, 
    var(--maroon-dark) 0%, 
    var(--maroon) 70%, 
    var(--gold-light) 50%),
        url("data:image/svg+xml,%3Csvg width='60' height='60' viewBox='0 0 60 60' xmlns='http://www.w3.org/2000/svg'%3E%3Cg fill='none' fill-rule='evenodd'%3E%3Cg fill='%23ffffff' fill-opacity='0.05'%3E%3Ccircle cx='30' cy='30' r='2'/%3E%3C/g%3E%3C/g%3E%3C/svg%3E") !important;
    min-height: 70vh;
    display: flex;
    align-items: center;
    position: relative;
    overflow: hidden;
}

.hero-content { position: relative; z-index: 2; }

.hero-title {
    font-family: 'Inter', sans-serif !important;
    font-size: 3.2rem !important;
    font-weight: 700 !important;
    color: var(--pure-white) !important;
    margin-bottom: 1.5rem !important;
    letter-spacing: -0.02em;
    line-height: 1.1;
}

.hero-subtitle {
    font-family: 'Inter', sans-serif !important;
    font-size: 1.1rem !important;
    color: rgba(255, 255, 255, 0.9) !important;
    margin-bottom: 2rem !important;
    font-weight: 400;
    max-width: 600px;
    line-height: 1.5;
}

.hero-lamp {
    background: linear-gradient(135deg, var(--lamp-light) 0%, var(--gold-light) 100%);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
    font-weight: 700 !important;
    font-size: 3.2rem !important;
}

.premium-container { background: transparent; min-height: 100vh; padding: 0; }
.content-section { padding: 4rem 0; background: var(--pure-white); }

h1, h2, h3, h4, h5, h6 {
    font-family: 'Inter', sans-serif !important;
    font-weight: 700 !important;
    color: var(--text-dark) !important;
    letter-spacing: -0.025em;
    margin-bottom: 1.5rem;
}

h2 {
    font-size: 2.2rem;
    background: linear-gradient(135deg, var(--maroon) 0%, var(--gold-dark) 100%);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
}

.card {
    background: var(--pure-white) !important;
    border: none !important;
    border-radius: 20px !important;
    box-shadow: 0 20px 40px var(--shadow-light) !important;
    backdrop-filter: blur(20px);
    transition: all 0.4s ease !important;
    overflow: hidden;
    position: relative;
}

.card::before {
    content: '';
    position: absolute;
    top: 0;
    left: 0;
    right: 0;
    height: 4px;
    background: linear-gradient(90deg, var(--maroon) 0%, var(--gold) 100%);
}

.card:hover {
    transform: translateY(-8px) !important;
    box-shadow: 0 30px 60px var(--shadow-medium) !important;
}

.card-body { padding: 2.5rem !important; }

.btn {
    background: linear-gradient(135deg, var(--pure-white) 0%, var(--pure-white) 100%) !important;
    color: var(--text-dark) !important;
    box-shadow: 0 8px 25px var(--pure-white) !important;
    font-family: 'Inter', sans-serif !important;
    font-weight: 600 !important;
    border-radius: 50px !important;
    padding: 0.7rem 1.8rem !important;  /* Reduced from 1rem 2.5rem */
    border: none !important;
    transition: all 0.4s ease !important;
    text-transform: none !important;
    font-size: 0.875rem !important;  /* Reduced from 1rem (14px instead of 16px) */
    letter-spacing: 0.025em;
    position: relative;
    overflow: hidden;
}

.btn-primary {
    background: linear-gradient(135deg, var(--pure-white) 0%, var(--pure-white) 100%) !important;
    color: var(--text-dark) !important;
    box-shadow: 0 8px 25px var(--pure-white) !important;
}

.btn-primary:hover {
    background: linear-gradient(135deg, var(--pure-white) 0%, var(--pure-white) 100%) !important;
    transform: translateY(-3px) !important;
    box-shadow: 0 15px 35px var(--pure-white) !important;
    color: var(--text-dark) !important;
}

.btn-secondary {
    background: linear-gradient(135deg, var(--pure-white) 0%, var(--pure-white) 100%) !important;
    color: var(--text-dark) !important;
    box-shadow: 0 8px 25px rgba(54, 69, 79, 0.3) !important;
}

.btn-gold {
    background: linear-gradient(135deg, var(--pure-white) 0%, var(--pure-white) 100%) !important;
    color: var(--text-dark) !important;
    box-shadow: 0 8px 25px rgba(255, 215, 0, 0.4) !important;
    font-weight: 700 !important;
}

.btn-link {
    background: linear-gradient(135deg, var(--pure-white) 0%, var(--pure-white) 100%) !important;
    color: var(--text-dark) !important;
    box-shadow: 0 2px 8px rgba(217, 199, 121, 0.4) !important;
    font-weight: 700 !important;
    color: var(--text-dark) !important;
    text-decoration: none !important;
    border-radius: 50px !important;
    padding: 0.8rem 2rem !important;
    transition: all 0.3s ease !important;
}

.form-control, .form-select {
    font-family: 'Inter', sans-serif !important;
    border: 2px solid var(--platinum-gray) !important;
    border-radius: 15px !important;
    padding: 1rem 1.5rem !important;
    background: var(--pure-white) !important;
    color: var(--text-dark) !important;
    font-size: 1rem !important;
    transition: all 0.3s ease !important;
    box-shadow: 0 2px 10px rgba(217, 199, 121, 0.05) !important;
}

.form-control:focus, .form-select:focus {
    box-shadow: 0 2px 10px rgba(217, 199, 121, 0.05) !important;
    transform: translateY(-2px) !important;
    outline: none !important;
}

.alert {
    border: none !important;
    border-radius: 15px !important;
    padding: 1.5rem 2rem !important;
    font-weight: 500 !important;
    border-left: 5px solid transparent !important;
    backdrop-filter: blur(10px);
    box-shadow: 0 8px 25px rgba(0, 0, 0, 0.1) !important;
}

.alert-success {
    background: linear-gradient(135deg, rgba(40, 167, 69, 0.1) 0%, rgba(32, 201, 151, 0.1) 100%) !important;
    color: var(--success-green) !important;
    border-left-color: var(--success-green) !important;
}

.upload-area {
    border: 3px dashed var(--gold) !important;
    border-radius: 20px !important;
    background: linear-gradient(135deg, rgba(255, 215, 0, 0.05) 0%, rgba(255, 237, 78, 0.05) 100%) !important;
    transition: all 0.4s ease !important;
    cursor: pointer !important;
    position: relative;
    overflow: hidden;
}

.dashboard-card {
    background: var(--pure-white) !important;
    border-radius: 25px !important;
    padding: 2.5rem !important;
    text-align: center;
    transition: all 0.4s ease !important;
    border: 1px solid var(--platinum-gray) !important;
    position: relative;
    overflow: hidden;
    height: 380px !important;  /* Add this line */
    display: flex !important;  /* Add this line */
    flex-direction: column !important;  /* Add this line */
    justify-content: space-between !important;  /* Add this line */
}

.dashboard-card::before {
    content: '';
    position: absolute;
    top: 0;
    left: 0;
    right: 0;
    height: 5px;
    background: linear-gradient(90deg, var(--maroon) 0%, var(--gold) 100%);
}

.dashboard-icon {
    font-size: 4rem;
    margin-bottom: 1.5rem;
    background: linear-gradient(135deg, var(--maroon) 0%, var(--gold) 100%);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
}

.table {
    border-radius: 15px !important;
    overflow: hidden !important;
    box-shadow: 0 10px 30px rgba(128, 0, 32, 0.1) !important;
}

.table th {
    background: linear-gradient(135deg, var(--maroon) 0%, var(--maroon-dark) 100%) !important;
    color: var(--pure-white) !important;
    font-weight: 600 !important;
    border: none !important;
    padding: 1.5rem !important;
    font-size: 1rem !important;
}

.fade-in-up { animation: fadeInUp 0.8s ease-out; }

@keyframes fadeInUp {
    from { opacity: 0; transform: translateY(40px); }
    to { opacity: 1; transform: translateY(0); }
}

@media (max-width: 768px) {
    .hero-title { font-size: 2.2rem !important; }
    .hero-subtitle { font-size: 1rem !important; }
    .navbar-brand { font-size: 1.4rem !important; }
    .lamp-text { font-size: 1.4rem !important; margin-left: 0.5rem !important; }
}
.Select-menu-outer {
    max-height: 300px !important;
    z-index: 9999 !important;
}

._dash-dropdown {
    z-index: 100 !important;
}

.card {
    overflow: visible !important;
}
.glass-card {
    background: rgba(255, 255, 255, 0.7) !important;
    backdrop-filter: blur(10px) !important;
    -webkit-backdrop-filter: blur(10px) !important;
    border: 1px solid rgba(255, 255, 255, 0.3) !important;
    box-shadow: 0 8px 32px 0 rgba(31, 38, 135, 0.15) !important;
}

/* Animated Gradient Background */
.animated-gradient {
    background: linear-gradient(-45deg, #f8f9fa, #e9ecef, #dee2e6, #f8f9fa);
    background-size: 400% 400%;
    animation: gradientShift 15s ease infinite;
}

@keyframes gradientShift {
    0% { background-position: 0% 50%; }
    50% { background-position: 100% 50%; }
    100% { background-position: 0% 50%; }
}

/* Pulse Animation for KPI Cards */
@keyframes pulse {
    0%, 100% { transform: scale(1); }
    50% { transform: scale(1.02); }
}

.kpi-card-hover:hover {
    animation: pulse 2s infinite;
    box-shadow: 0 8px 24px rgba(196, 30, 58, 0.2) !important;
    transition: all 0.3s ease;
}

/* Modern Filter Sidebar */
.filter-sidebar {
    background: linear-gradient(135deg, #ffffff 0%, #f8f9fa 100%);
    border-right: 1px solid rgba(196, 30, 58, 0.1);
}

/* Sleek Dropdown Styling */
.Select-control {
    background: rgba(255, 255, 255, 0.9) !important;
    border: 2px solid transparent !important;
    transition: all 0.3s ease !important;
}

.Select-control:hover {
    border-color: rgba(196, 30, 58, 0.3) !important;
    box-shadow: 0 4px 12px rgba(196, 30, 58, 0.1) !important;
}

/* Smooth Transitions */
* {
    transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
}

/* Loading Skeleton Animation */
@keyframes shimmer {
    0% { background-position: -1000px 0; }
    100% { background-position: 1000px 0; }
}

.skeleton {
    animation: shimmer 2s infinite;
    background: linear-gradient(to right, #f6f7f8 0%, #edeef1 20%, #f6f7f8 40%, #f6f7f8 100%);
    background-size: 1000px 100%;
}

/* Status Badges */
.status-badge {
    display: inline-block;
    padding: 4px 12px;
    border-radius: 20px;
    font-size: 12px;
    font-weight: 600;
    text-transform: uppercase;
    letter-spacing: 0.5px;
}

.status-positive {
    background: linear-gradient(135deg, #10b981 0%, #059669 100%);
    color: white;
}

.status-negative {
    background: linear-gradient(135deg, #ef4444 0%, #dc2626 100%);
    color: white;
}

/* Modern Scrollbar */
::-webkit-scrollbar {
    width: 8px;
    height: 8px;
}

::-webkit-scrollbar-track {
    background: #f1f1f1;
    border-radius: 10px;
}

::-webkit-scrollbar-thumb {
    background: linear-gradient(135deg, #C41E3A 0%, #8b1528 100%);
    border-radius: 10px;
}

::-webkit-scrollbar-thumb:hover {
    background: #8b1528;
}

/* Data Visualization Enhancement */
.metric-trend-up::after {
    content: "↑";
    color: #10b981;
    margin-left: 5px;
    font-size: 1.2em;
}

.metric-trend-down::after {
    content: "↓";
    color: #ef4444;
    margin-left: 5px;
    font-size: 1.2em;
}

/* Micro-interactions */
.interactive-element {
    cursor: pointer;
    position: relative;
    overflow: hidden;
}

.interactive-element::before {
    content: '';
    position: absolute;
    top: 50%;
    left: 50%;
    width: 0;
    height: 0;
    border-radius: 50%;
    background: rgba(196, 30, 58, 0.1);
    transform: translate(-50%, -50%);
    transition: width 0.6s, height 0.6s;
}

.interactive-element:active::before {
    width: 300px;
    height: 300px;
}

/* Modern Card Shadows with Depth */
.depth-card {
    box-shadow: 
        0 1px 2px rgba(0,0,0,0.07),
        0 2px 4px rgba(0,0,0,0.07),
        0 4px 8px rgba(0,0,0,0.07),
        0 8px 16px rgba(0,0,0,0.07),
        0 16px 32px rgba(0,0,0,0.07),
        0 32px 64px rgba(0,0,0,0.07) !important;
}
.modern-tab-btn {
    background: white;
    border: 2px solid #e5e7eb;
    padding: 12px 24px;
    fontSize: 14px;
    fontWeight: 600;
    color: #374151;
    borderRadius: 10px;
    cursor: pointer;
    transition: all 0.3s ease;
    display: flex;
    align-items: center;
    gap: 8px;
}

.modern-tab-btn:hover {
    background: linear-gradient(135deg, #C41E3A 0%, #8b1528 100%);
    color: white;
    border-color: #C41E3A;
    transform: translateY(-2px);
    box-shadow: 0 4px 12px rgba(196, 30, 58, 0.3);
}

.modern-tab-btn:active {
    transform: translateY(0);
}
"""

app.index_string = f'''
<!DOCTYPE html>
<html>
    <head>
        {{%metas%}}
        <title>{{%title%}}</title>
        {{%favicon%}}
        <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.1.3/dist/css/bootstrap.min.css" rel="stylesheet">
        <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.0.0/css/all.min.css">
        <style>{premium_css}</style>
        {{%css%}}
    </head>
    <body>
        {{%app_entry%}}
        <footer>{{%config%}}{{%scripts%}}{{%renderer%}}</footer>
    </body>
</html>
'''

# Global variables
current_user = {}
otp_storage = {}
VALID_DOMAINS = ['@trustgroup.in', '@trustplutus.com', '@prajana-alternatives.com']

# Ensure data directories exist
os.makedirs('data', exist_ok=True)
os.makedirs('data/user_data', exist_ok=True)


def save_plotly_as_image(fig, width=800, height=600):
    """Convert Plotly figure to image bytes for embedding in PDF/PPT"""
    try:
        import kaleido
        img_bytes = fig.to_image(format="png", width=width, height=height, engine="kaleido")
        return img_bytes
    except ImportError:
        print("ERROR: kaleido not installed. Run: pip install kaleido==0.2.1")
        return None
    except Exception as e:
        print(f"ERROR converting plot to image: {e}")
        import traceback
        traceback.print_exc()
        return None


def create_title_slide_ppt(prs, title, subtitle, filter_desc):
    """Create title slide for PowerPoint with TRUSTPLUTUS branding"""
    MAROON = RGBColor(176, 77, 89)
    GOLD = RGBColor(247, 206, 124)

    title_slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(title_slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "TRUSTPLUTUS"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = MAROON
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = GOLD
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Filter info
    filter_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.5))
    filter_frame = filter_box.text_frame
    filter_frame.text = f"Filter: {filter_desc}"
    filter_para = filter_frame.paragraphs[0]
    filter_para.font.size = Pt(16)
    filter_para.font.color.rgb = RGBColor(102, 102, 102)
    filter_para.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = datetime.now().strftime('%B %d, %Y')
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(14)
    date_para.font.color.rgb = RGBColor(102, 102, 102)
    date_para.alignment = PP_ALIGN.CENTER

    return prs


def add_chart_slide_ppt(prs, title, fig, notes=""):
    """Add a slide with a Plotly chart to PowerPoint"""
    MAROON = RGBColor(176, 77, 89)

    slide_layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.color.rgb = MAROON

    # Convert Plotly to image
    img_bytes = save_plotly_as_image(fig, width=900, height=500)

    if img_bytes:
        img_stream = io.BytesIO(img_bytes)
        pic = slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.5), width=Inches(9))

    # Add notes if provided
    if notes:
        notes_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
        notes_frame = notes_box.text_frame
        notes_frame.text = notes
        notes_para = notes_frame.paragraphs[0]
        notes_para.font.size = Pt(12)
        notes_para.font.color.rgb = RGBColor(102, 102, 102)

    return prs


# Utility functions
def generate_user_id(name, role):
    role_mapping = {'Investment Strategist': 'IS', 'Technical Analyst': 'TA', 'Banker': 'BA', 'Management Team': 'MT',
                    'Others': 'O'}
    clean_name = ''.join(e for e in name if e.isalnum())
    random_num = ''.join(random.choices(string.digits, k=4))
    role_initial = role_mapping.get(role, 'O')
    return f"{role_initial}_{clean_name}_{random_num}"


def generate_otp():
    return ''.join(random.choices(string.digits, k=6))


def hash_password(password):
    return hashlib.sha256(password.encode()).hexdigest()


# NEW DATA LOADING FUNCTIONS
def load_key_matrix_data():
    """Load key matrix quarterly data"""
    try:
        return pd.read_excel('data/key_matrix_data.xlsx')
    except FileNotFoundError:
        return pd.DataFrame()


def load_crosssell_detailed_data():
    """Load detailed cross-sell data with hierarchy"""
    try:
        return pd.read_excel('data/crosssell_data.xlsx')
    except FileNotFoundError:
        return pd.DataFrame()


def load_mandate_detailed_data():
    """Load detailed mandate data with hierarchy"""
    try:
        return pd.read_excel('data/mandate_data.xlsx')
    except FileNotFoundError:
        return pd.DataFrame()


def filter_data_by_hierarchy(df, region, banker, strategist):
    """Filter dataframe based on hierarchical selection
    Investment Strategist filter is INDEPENDENT"""
    if strategist and strategist != 'All Strategists':
        filtered_df = df[df['Investment_Strategist'] == strategist].copy()
    else:
        if region == 'All Regions':
            filtered_df = df.copy()
        else:
            filtered_df = df[df['Region'] == region].copy()
            if banker != 'All Bankers':
                filtered_df = filtered_df[filtered_df['Banker'] == banker].copy()
    return filtered_df


def get_comparison_data(df, compare_type, compare_values, metric_columns):
    """Get data for comparison across multiple entities"""
    comparison_results = []
    for value in compare_values:
        if compare_type == 'region':
            filtered = df[df['Region'] == value]
        elif compare_type == 'banker':
            filtered = df[df['Banker'] == value]
        elif compare_type == 'strategist':
            filtered = df[df['Investment_Strategist'] == value]
        else:
            continue

        if not filtered.empty:
            metrics = {'Name': value}
            for col in metric_columns:
                if col in filtered.columns:
                    metrics[col] = filtered[col].sum() if pd.api.types.is_numeric_dtype(filtered[col]) else \
                    filtered[col].iloc[0]
            comparison_results.append(metrics)
    return pd.DataFrame(comparison_results)


def create_comparison_chart(comparison_df, metric_name, chart_title, chart_type='bar'):
    """Create a comparison chart (bar or line)"""
    if comparison_df.empty:
        return go.Figure()

    fig = go.Figure()
    if chart_type == 'bar':
        fig.add_trace(go.Bar(
            x=comparison_df['Name'],
            y=comparison_df[metric_name],
            marker_color='#C41E3A',
            text=comparison_df[metric_name],
            textposition='outside',
            texttemplate='%{text:.1f}'
        ))
    else:
        fig.add_trace(go.Scatter(
            x=comparison_df['Name'],
            y=comparison_df[metric_name],
            mode='lines+markers+text',
            marker=dict(size=10, color='#C41E3A'),
            line=dict(color='#C41E3A', width=2),
            text=comparison_df[metric_name],
            textposition='top center',
            texttemplate='%{text:.1f}'
        ))

    fig.update_layout(
        title=chart_title,
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=12),
        title_font=dict(size=16, color='#C41E3A'),
        height=400,
        showlegend=False,
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=True, gridcolor='#f0f0f0')
    )
    return fig


def send_otp_email(email, otp):
    print(f"OTP for {email}: {otp}")
    return True


def validate_email_domain(email):
    return any(email.endswith(domain) for domain in VALID_DOMAINS)


def load_users():
    users_file = 'data/users.xlsx'
    if os.path.exists(users_file):
        return pd.read_excel(users_file)
    else:
        return pd.DataFrame(columns=['Name', 'Phone', 'Email', 'Role', 'UserID', 'Password'])


def save_user(name, phone, email, role, user_id, password):
    users_df = load_users()
    new_user = pd.DataFrame({'Name': [name], 'Phone': [phone], 'Email': [email], 'Role': [role], 'UserID': [user_id],
                             'Password': [hash_password(password)]})
    users_df = pd.concat([users_df, new_user], ignore_index=True)
    users_df.to_excel('data/users.xlsx', index=False)


def authenticate_user(user_id, password):
    users_df = load_users()
    user_row = users_df[users_df['UserID'] == user_id]
    if not user_row.empty:
        stored_password = user_row.iloc[0]['Password']
        if stored_password == hash_password(password):
            return user_row.iloc[0].to_dict()
    return None


def parse_excel_upload(contents, filename):
    try:
        content_type, content_string = contents.split(',')
        decoded = base64.b64decode(content_string)
        df = pd.read_excel(io.BytesIO(decoded))
        required_columns = ['Family_Name', 'Product', 'Stage', 'New_Liquidity']
        if all(col in df.columns for col in required_columns):
            df['Timestamp'] = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            return df[['Timestamp'] + required_columns], None
        else:
            missing_cols = [col for col in required_columns if col not in df.columns]
            return None, f"Missing required columns: {', '.join(missing_cols)}"
    except Exception as e:
        return None, f"Error processing file: {str(e)}"


def save_bulk_cross_sell_data(user_id, df):
    filename = f'data/user_data/{user_id}_crosssell.xlsx'
    if os.path.exists(filename):
        existing_df = pd.read_excel(filename)
        combined_df = pd.concat([existing_df, df], ignore_index=True)
    else:
        combined_df = df
    combined_df.to_excel(filename, index=False)
    return len(df)
def save_proposal_data(user_id, proposal_data):
    """Save proposal data to Excel file"""
    filename = f'data/user_data/{user_id}_proposals.xlsx'
    if os.path.exists(filename):
        existing_df = pd.read_excel(filename)
        combined_df = pd.concat([existing_df, proposal_data], ignore_index=True)
    else:
        combined_df = proposal_data
    combined_df.to_excel(filename, index=False)
    return True

def save_pdf_file(user_id, pdf_data, original_filename):
    """Save uploaded PDF file"""
    user_folder = f'data/user_data/{user_id}_proposals'
    os.makedirs(user_folder, exist_ok=True)
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    safe_filename = f"{timestamp}_{original_filename}"
    filepath = os.path.join(user_folder, safe_filename)
    with open(filepath, 'wb') as f:
        f.write(pdf_data)
    return safe_filename

def search_existing_clients(family_name_search):
    """Search for existing clients with similar names"""
    user_id = current_user.get('UserID', '')
    filename = f'data/user_data/{user_id}_proposals.xlsx'

    if not os.path.exists(filename):
        return []

    try:
        df = pd.read_excel(filename)
        if 'Family_Name' in df.columns:
            # Search for partial matches (case insensitive)
            matches = df[df['Family_Name'].str.contains(family_name_search, case=False, na=False)]
            # Remove duplicates based on Family_Name and return latest entry for each
            latest_matches = matches.groupby('Family_Name').last().reset_index()
            return latest_matches.to_dict('records')
        return []
    except Exception as e:
        print(f"Error searching clients: {e}")
        return []
def load_client_data(family_name):
    """Load existing client data for editing"""
    user_id = current_user.get('UserID', '')
    filename = f'data/user_data/{user_id}_proposals.xlsx'

    if not os.path.exists(filename):
        return None

    try:
        df = pd.read_excel(filename)
        client_data = df[df['Family_Name'] == family_name]
        if not client_data.empty:
            return client_data.iloc[-1].to_dict()  # Get the latest entry
        return None
    except Exception as e:
        print(f"Error loading client data: {e}")
        return None

def update_proposal_data(user_id, proposal_data, family_name):
    """Update existing proposal data instead of creating new entry"""
    filename = f'data/user_data/{user_id}_proposals.xlsx'

    if os.path.exists(filename):
        existing_df = pd.read_excel(filename)
        # Remove existing entries for this client
        existing_df = existing_df[existing_df['Family_Name'] != family_name]
        # Add the updated entry
        proposal_df = pd.DataFrame(proposal_data)
        combined_df = pd.concat([existing_df, proposal_df], ignore_index=True)
    else:
        combined_df = pd.DataFrame(proposal_data)

    combined_df.to_excel(filename, index=False)
    return True
def load_dashboard_data():
    try:
        dashboard_df = pd.read_excel('data/key_matrix_data.xlsx',sheet_name = "Dashboard")
        crosssell_df = pd.read_excel('data/crosssell_data.xlsx')
        mandate_df = pd.read_excel('data/mandate_data.xlsx')
        return dashboard_df, crosssell_df, mandate_df
    except FileNotFoundError:
        return pd.DataFrame(), pd.DataFrame(), pd.DataFrame()


def get_dashboard_metrics(dashboard_df, selected_year=None):
    if selected_year and not dashboard_df.empty:
        filtered_df = dashboard_df[dashboard_df['Year'] == selected_year]
    else:
        if not dashboard_df.empty:
            latest_year = dashboard_df['Year'].max()
            filtered_df = dashboard_df[dashboard_df['Year'] == latest_year]
        else:
            return {}

    if filtered_df.empty:
        return {}

    metrics = {
        'families': len(filtered_df['Family_Name'].unique()),
        'aum_cr': round(filtered_df['AUM_Cr'].sum(), 1),
        'proprietary_products': len(filtered_df[filtered_df['Status'] == 'Proprietary Products']),
        'billed_revenue': round(filtered_df['Revenue_Cr'].sum(), 1),
        'total_revenue': round(filtered_df['Revenue_Cr'].sum() * 1.2, 1),
        'year': filtered_df['Year'].iloc[0] if len(filtered_df) > 0 else 2024
    }
    return metrics


def create_key_matrix_charts(dashboard_df):
    if dashboard_df.empty:
        return {}, {}, {}

    dashboard_df['Period'] = dashboard_df['Quarter'] + ' ' + dashboard_df['Year'].astype(str)
    quarterly_data = dashboard_df.groupby('Period').agg({
        'Family_Name': 'nunique',
        'AUM_Cr': 'sum',
        'Revenue_Cr': 'sum'
    }).reset_index()

    trust_maroon = '#800020'
    plutus_gold = '#B8860B'

    fig_families = go.Figure()
    fig_families.add_trace(
        go.Bar(x=quarterly_data['Period'], y=quarterly_data['Family_Name'], marker_color=trust_maroon, name='Families'))
    fig_families.update_layout(title='No. of Families', plot_bgcolor='white', paper_bgcolor='white',
                               font=dict(family='Inter', size=12), title_font=dict(size=14, color=trust_maroon),
                               height=300)

    fig_aum = go.Figure()
    fig_aum.add_trace(
        go.Bar(x=quarterly_data['Period'], y=quarterly_data['AUM_Cr'], marker_color=plutus_gold, name='AUM (Cr)'))
    fig_aum.update_layout(title='AUM in Cr', plot_bgcolor='white', paper_bgcolor='white',
                          font=dict(family='Inter', size=12), title_font=dict(size=14, color=trust_maroon), height=300)

    fig_revenue = go.Figure()
    fig_revenue.add_trace(
        go.Bar(x=quarterly_data['Period'], y=quarterly_data['Revenue_Cr'], marker_color='#6C757D', name='Revenue (Cr)'))
    fig_revenue.update_layout(title='Total Revenue', plot_bgcolor='white', paper_bgcolor='white',
                              font=dict(family='Inter', size=12), title_font=dict(size=14, color=trust_maroon),
                              height=300)

    return fig_families, fig_aum, fig_revenue


def create_mandate_pipeline_chart(mandate_df):
    if mandate_df.empty:
        return {}

    mandate_df['Period'] = mandate_df['Quarter'] + ' ' + mandate_df['Year'].astype(str)
    pivot_data = mandate_df.pivot_table(index='Period', columns='Stage', values='Count', fill_value=0).reset_index()

    trust_maroon = '#800020'
    colors = {'Pitched': trust_maroon, 'Offered': '#A0203A', 'Converted': '#B8860B', 'Declined': '#6C757D'}

    fig = go.Figure()
    for stage in ['Pitched', 'Offered', 'Converted', 'Declined']:
        if stage in pivot_data.columns:
            fig.add_trace(go.Bar(name=stage, x=pivot_data['Period'], y=pivot_data[stage],
                                 marker_color=colors.get(stage, trust_maroon)))

    fig.update_layout(title='New Mandates Pipeline', barmode='group', plot_bgcolor='white', paper_bgcolor='white',
                      font=dict(family='Inter', size=12), title_font=dict(size=16, color=trust_maroon), height=400,
                      legend=dict(orientation="h", y=-0.2))
    return fig


def get_navigation(is_logged_in=False, user_name=""):
    """Navigation bar that changes based on login status"""

    if is_logged_in and user_name:
        # Show logout button when logged in
        auth_buttons = html.Div([
            html.Span(f"Welcome, {user_name}",
                      className="me-3",
                      style={"fontSize": "13px", "color": "#666", "fontWeight": "500"}),
            dbc.Button("Logout", id="nav-logout-btn", className="btn-link",
                       style={"fontSize": "13px", "padding": "8px 20px"})
        ], className="d-flex align-items-center")
    else:
        # Show login/register buttons when not logged in
        auth_buttons = html.Div([
            dbc.Button("Login", id="nav-login-btn", className="btn-link me-2",
                       style={"fontSize": "13px", "padding": "8px 20px"}),
            dbc.Button("New to LAMP", id="nav-register-btn", className="btn-link",
                       style={"fontSize": "13px", "padding": "8px 20px"})
        ], className="d-flex")

    return dbc.Navbar([
        dbc.Container([
            dbc.Row([
                dbc.Col([
                    html.Span(
                        [html.Span("TRUST", className="trust-text"),
                         html.Span("PLUTUS", className="plutus-text")],
                        className="navbar-brand")
                ], width="auto")
            ], align="center", className="g-0"),
            dbc.Row([
                dbc.Col([auth_buttons], width="auto")
            ], align="center", className="g-0")
        ], fluid=True, className="d-flex justify-content-between align-items-center")
    ], className="main-navbar fixed-top", expand="lg")

def get_hero_section():
    return html.Div([
        dbc.Container([
            dbc.Row([
                dbc.Col([
                    html.Div([
                        html.H1([html.Span("LAMP", className="hero-lamp me-3"), html.Br(),
                                 "Large Account Management Process"], className="hero-title"),
                        html.P(
                            "Revolutionizing wealth management through intelligent automation and strategic insights. Experience the future of account management with TRUSTPLUTUS.",
                            className="hero-subtitle"),
                        html.Div([
                            dbc.Button("Get Started", id="hero-get-started", className="btn-link",
                                      style={"fontSize": "14px", "padding": "10px 24px"}),
                            dbc.Button("Learn More", id="hero-learn-more", className="btn-link",
                                      style={"fontSize": "14px", "padding": "10px 24px"})
                        ])
                    ], className="hero-content")
                ], lg=8)
            ], justify="start", align="center", style={"minHeight": "70vh"})
        ], fluid=True)
    ], className="hero-section")

# Layout functions
def get_home_layout():
    return html.Div([
        get_hero_section(),
        html.Div([
            dbc.Container([
                dbc.Row([
                    dbc.Col([
                        html.H2("Why Choose LAMP?", className="text-center lamp-text mb-5"),
                        dbc.Row([
                            dbc.Col([
                                html.Div([
                                    html.I(className="fas fa-chart-line dashboard-icon"),
                                    html.H5("Advanced Analytics", className="mb-3"),
                                    html.P(
                                        "Leverage cutting-edge analytics to make informed investment decisions and optimize portfolio performance.",
                                        className="text-muted")
                                ], className="text-center")
                            ], lg=4, className="mb-4"),
                            dbc.Col([
                                html.Div([
                                    html.I(className="fas fa-shield-alt dashboard-icon"),
                                    html.H5("Secure Platform", className="mb-3"),
                                    html.P(
                                        "Enterprise-grade security with multi-factor authentication and encrypted data storage.",
                                        className="text-muted")
                                ], className="text-center")
                            ], lg=4, className="mb-4"),
                            dbc.Col([
                                html.Div([
                                    html.I(className="fas fa-users dashboard-icon"),
                                    html.H5("Client Management", className="mb-3"),
                                    html.P(
                                        "Streamlined client relationship management with automated workflows and proposal tracking.",
                                        className="text-muted")
                                ], className="text-center")
                            ], lg=4, className="mb-4")
                        ])
                    ], lg=10)
                ], justify="center")
            ], fluid=True)
        ], className="content-section")
    ], className="fade-in-up")


def get_registration_layout():
    return html.Div([
        html.Div([
            dbc.Container([
                dbc.Row([
                    dbc.Col([
                        html.Div([
                            html.H2("Join LAMP", className="text-center mb-2"),
                            html.P("Create your LAMP account and unlock powerful wealth management tools",
                                   className="text-center text-muted mb-5",
                                   style={"fontSize": "1.1rem", "fontWeight": "400"})
                        ], className="mb-5"),
                        dbc.Card([
                            dbc.CardHeader([html.H5("Account Registration", className="mb-0 text-center")]),
                            dbc.CardBody([
                                dbc.Form([
                                    dbc.Row([
                                        dbc.Label("Full Name", width=3, className="form-label"),
                                        dbc.Col([dbc.Input(id="reg-name", type="text",
                                                           placeholder="Enter your full name",
                                                           className="form-control")], width=9)
                                    ], className="mb-4"),
                                    dbc.Row([
                                        dbc.Label("Phone Number", width=3, className="form-label"),
                                        dbc.Col([dbc.Input(id="reg-phone", type="text",
                                                           placeholder="Enter phone number", className="form-control")],
                                                width=9)
                                    ], className="mb-4"),
                                    dbc.Row([
                                        dbc.Label("Email ID", width=3, className="form-label"),
                                        dbc.Col([
                                            dbc.Input(id="reg-email", type="email", placeholder="Enter office email",
                                                      className="form-control mb-2"),
                                            html.Small(
                                                "Must end with @trustgroup.in, @trustplutus.com, or @prajana-alternatives.com",
                                                className="text-muted")
                                        ], width=9)
                                    ], className="mb-4"),
                                    dbc.Row([
                                        dbc.Label("Role", width=3, className="form-label"),
                                        dbc.Col([
                                            dcc.Dropdown(
                                                id="reg-role",
                                                options=[
                                                    {"label": "Investment Strategist",
                                                     "value": "Investment Strategist"},
                                                    {"label": "Technical Analyst", "value": "Technical Analyst"},
                                                    {"label": "Banker", "value": "Banker"},
                                                    {"label": "Management Team", "value": "Management Team"},
                                                    {"label": "Others", "value": "Others"}
                                                ],
                                                placeholder="Select your role"
                                            )
                                        ], width=9)
                                    ], className="mb-4"),
                                    dbc.Row([
                                        dbc.Col(
                                            [dbc.Button("Send OTP", id="send-otp-btn",
                                               className="btn-link btn-primary w-100",
                                                style={"fontSize": "13px", "padding": "10px 20px"})])
                                    ], className="mb-4"),
                                    html.Div(id="otp-section", style={"display": "none"}, children=[
                                        dbc.Alert([
                                            html.H6("OTP Sent Successfully!", className="alert-heading"),
                                            html.P("Check your console for the OTP code", className="mb-0")
                                        ], color="success", className="mb-4"),
                                        dbc.Row([
                                            dbc.Label("Enter OTP", width=3, className="form-label"),
                                            dbc.Col([dbc.Input(id="otp-input", type="text",
                                                               placeholder="Enter 6-digit OTP",
                                                               className="form-control")], width=9)
                                        ], className="mb-4"),
                                        dbc.Row([
                                            dbc.Label("Create Password", width=3, className="form-label"),
                                            dbc.Col([dbc.Input(id="reg-password", type="password",
                                                               placeholder="Create a secure password",
                                                               className="form-control")], width=9)
                                        ], className="mb-4"),
                                        dbc.Button("Complete Registration", id="register-btn",
                                                   className="btn-success w-100",
                                                   style={"fontSize": "13px", "padding": "10px 20px"})
                                    ])
                                ])
                            ])
                        ], className="card shadow-lg")
                    ], width=10, lg=8)
                ], justify="center"),
                html.Hr(className="my-5"),
                html.Div([
                    html.P("Already have an account?", className="text-center mb-3", style={"fontSize": "1.1rem"}),
                    dbc.Button("Sign In", id="login-btn",
                               className="btn-link btn-primary w-100",
                               style={"fontSize": "13px", "padding": "10px 20px"})
                ])
            ], fluid=True)
        ], className="content-section")
    ], className="fade-in-up", style={"paddingTop": "2rem"})


def get_login_layout():
    return html.Div([
        html.Div([
            dbc.Container([
                dbc.Row([
                    dbc.Col([
                        # Welcome Header Section
                        html.Div([
                            html.H2("Welcome Back!", className="text-center mb-2",
                                    style={"color": "var(--maroon)", "fontSize": "1.8rem"}),
                            html.P("Sign in to your LAMP account to continue",
                                   className="text-center text-muted mb-4",
                                   style={"fontSize": "1rem", "fontWeight": "400"})
                        ], className="mb-4"),

                        # Login Card
                        dbc.Card([
                            dbc.CardBody([
                                dbc.Form([
                                    # User ID Field
                                    dbc.Row([
                                        dbc.Col([
                                            html.Label([
                                                html.I(className="fas fa-user",
                                                       style={"marginRight": "8px", "color": "var(--maroon)"}),
                                                "User ID"
                                            ], className="form-label",
                                                style={"fontWeight": "600", "fontSize": "0.95rem"}),
                                            dbc.Input(
                                                id="login-userid",
                                                type="text",
                                                placeholder="Enter your User ID",
                                                className="form-control",
                                                style={"fontSize": "0.95rem"}
                                            )
                                        ], width=12)
                                    ], className="mb-4"),

                                    # Password Field
                                    dbc.Row([
                                        dbc.Col([
                                            html.Label([
                                                html.I(className="fas fa-lock",
                                                       style={"marginRight": "8px", "color": "var(--maroon)"}),
                                                "Password"
                                            ], className="form-label",
                                                style={"fontWeight": "600", "fontSize": "0.95rem"}),
                                            dbc.Input(
                                                id="login-password",
                                                type="password",
                                                placeholder="Enter your password",
                                                className="form-control",
                                                style={"fontSize": "0.95rem"}
                                            )
                                        ], width=12)
                                    ], className="mb-4"),

                                    # Login Button
                                    dbc.Button([
                                        html.I(className="fas fa-sign-in-alt",
                                               style={"marginRight": "8px"}),
                                        "Sign In"
                                    ], id="login-btn",
                                        className="btn-primary w-100",
                                        style={
                                            "fontSize": "1rem",
                                            "padding": "12px 20px",
                                            "fontWeight": "600"
                                        })
                                ])
                            ], style={"padding": "2rem"})
                        ], className="card shadow-lg", style={"border": "none"}),

                        # Divider
                        html.Div([
                            html.Div(style={
                                "height": "1px",
                                "background": "linear-gradient(90deg, transparent, #ddd, transparent)",
                                "margin": "2rem 0"
                            }),

                            # New User Section
                            html.Div([
                                html.P("Don't have an account?",
                                       className="text-center mb-3",
                                       style={"fontSize": "0.95rem", "color": "#666"}),
                                dbc.Button([
                                    html.I(className="fas fa-user-plus",
                                           style={"marginRight": "8px"}),
                                    "Create New Account"
                                ], id="go-to-register",
                                    className="btn-link d-block mx-auto",
                                    style={
                                        "fontSize": "0.95rem",
                                        "padding": "10px 24px"
                                    })
                            ])
                        ])

                    ], width=12, lg=5, md=7, sm=10)
                ], justify="center", style={"minHeight": "calc(100vh - 160px)", "alignItems": "center"})
            ], fluid=True)
        ], className="content-section", style={"paddingTop": "2rem", "paddingBottom": "2rem"})
    ], className="fade-in-up")

def get_main_dashboard_layout():
    user_name = current_user.get('Name', 'User')
    user_role = current_user.get('Role', 'User')
    return html.Div([
        html.Div([
            dbc.Container([
                dbc.Row([
                    dbc.Col([
                        dbc.Card([
                            dbc.CardHeader([html.H5("Quick Actions", className="mb-0 text-center")]),
                            dbc.CardBody([
                                html.P("Choose your workflow to get started", className="text-center text-muted mb-4"),
                                dbc.Row([
                                    dbc.Col([
                                        html.Div([
                                            html.I(className="fas fa-chart-bar dashboard-icon"),
                                            html.H6("Dashboard", className="text-center mb-3"),
                                            html.P("View comprehensive analytics and performance metrics",
                                                   className="text-center text-muted small mb-3"),
                                            dbc.Button("Launch Dashboard", id="dashboard-btn",
                                                       className="btn-link btn-primary w-100",
                                                       style={"fontSize": "13px", "padding": "10px 20px"})
                                        ], className="dashboard-card")
                                    ], lg=4, className="mb-4"),
                                    dbc.Col([
                                        html.Div([
                                            html.I(className="fas fa-chart-line dashboard-icon"),
                                            html.H6("Cross Sell", className="text-center mb-3"),
                                            html.P("Manage cross-selling opportunities and track client interactions",
                                                   className="text-center text-muted small mb-3"),
                                            dbc.Button("Launch Cross Sell", id="cross-sell-btn",
                                                       className="btn-link btn-primary w-100",
                                                       style={"fontSize": "13px", "padding": "10px 20px"})
                                        ], className="dashboard-card")
                                    ], lg=4, className="mb-4"),
                                    dbc.Col([
                                        html.Div([
                                            html.I(className="fas fa-file-alt dashboard-icon"),
                                            html.H6("Proposals", className="text-center mb-3"),
                                            html.P("Create and manage client proposals with advanced tracking",
                                                   className="text-center text-muted small mb-3"),
                                            dbc.Button("Launch Proposals", id="proposal-btn",
                                                       className="btn-link btn-primary w-100",
                                                       style={"fontSize": "13px", "padding": "10px 20px"})
                                        ], className="dashboard-card")
                                    ], lg=4, className="mb-4")
                                ])
                            ])
                        ], className="card shadow-lg")
                    ], width=12, lg=10)
                ], justify="center"),
                html.Hr(className="my-5"),
                html.Div([
                    dbc.Button("Sign Out", id="logout-btn", className="btn-link btn-outline-primary")
                ], className="text-center")
            ], fluid=True)
        ], className="content-section")
    ], className="fade-in-up", style={"paddingTop": "2rem"})



def get_dashboard_layout():
    """
    Dashboard layout with filters in left sidebar and content on right
    Optimized for professional look and better screen fitting
    """
    # Load all data to get unique values
    try:
        full_df = pd.read_excel('data/key_matrix_data.xlsx', sheet_name='Dashboard')
    except Exception as e:
        print(f"Error reading key_matrix_data.xlsx: {e}")
        full_df = pd.DataFrame()

    # Extract unique values
    regions = ['All Regions']
    region_banker_map = {'All Regions': ['All Bankers']}
    all_strategists = ['All Strategists']

    if not full_df.empty:
        # Get all unique strategists (independent filter)
        if 'Investment_Strategist' in full_df.columns:
            unique_strategists = sorted(full_df[full_df['Investment_Strategist'].notna()]['Investment_Strategist'].unique().tolist())
            all_strategists.extend(unique_strategists)

        # Get regions and bankers
        if 'Region' in full_df.columns:
            unique_regions = sorted(full_df[full_df['Region'].notna()]['Region'].unique().tolist())
            regions.extend(unique_regions)

            # Map bankers to each region
            for region in unique_regions:
                if 'Banker' in full_df.columns:
                    bankers = sorted(full_df[full_df['Region'] == region]['Banker'].dropna().unique().tolist())
                    region_banker_map[region] = ['All Bankers'] + bankers
                else:
                    region_banker_map[region] = ['All Bankers']

    return html.Div([
        dbc.Container([
            # Modern Header with Gradient - Reduced spacing
            dbc.Row([
                dbc.Col([
                    html.Div([
                        html.H1("Dashboard FY23", style={
                            "background": "linear-gradient(135deg, #C41E3A 0%, #8b1528 100%)",
                            "-webkit-background-clip": "text",
                            "-webkit-text-fill-color": "transparent",
                            "fontWeight": "800",
                            "fontSize": "28px",
                            "margin": "0px 0 5px 0",
                            "letterSpacing": "-0.5px"
                        }),
                        html.Div([
                            html.Span("●", style={"color": "#10b981", "marginRight": "6px", "fontSize": "10px"}),
                            html.Span("Live Data", style={"fontSize": "12px", "color": "#6b7280", "fontWeight": "500"})
                        ])
                    ])
                ], width=12)
            ]),

            # Modern Navigation Tabs with Icons - Reduced spacing
            html.Div([
                html.Button([
                    html.I(className="fas fa-chart-line", style={"marginRight": "6px"}),
                    "Key Matrix"
                ], id="key-matrix-btn", className="modern-tab-btn", style={
                    "padding": "8px 16px",
                    "fontSize": "13px"
                }),
                html.Button([
                    html.I(className="fas fa-exchange-alt", style={"marginRight": "6px"}),
                    "X-Sell"
                ], id="x-sell-dash-btn", className="modern-tab-btn", style={
                    "padding": "8px 16px",
                    "fontSize": "13px"
                }),
                html.Button([
                    html.I(className="fas fa-briefcase", style={"marginRight": "6px"}),
                    "New Mandates"
                ], id="new-mandates-btn", className="modern-tab-btn", style={
                    "padding": "8px 16px",
                    "fontSize": "13px"
                })
            ], style={
                "marginBottom": "15px",
                "display": "flex",
                "gap": "8px"
            }),

            # Main Content Row
            dbc.Row([
                # LEFT SIDEBAR - Enhanced Filters - More compact
                dbc.Col([
                    html.Div([
                        # Filter Header with Icon
                        html.Div([
                            html.I(className="fas fa-filter", style={
                                "color": "#C41E3A",
                                "marginRight": "8px",
                                "fontSize": "16px"
                            }),
                            html.H5("Filters", style={
                                "display": "inline",
                                "color": "#1f2937",
                                "fontWeight": "700",
                                "margin": "0",
                                "fontSize": "16px"
                            })
                        ], style={
                            "marginBottom": "15px",
                            "paddingBottom": "10px",
                            "borderBottom": "2px solid #e5e7eb"
                        }),

                        # Region Filter with Icon
                        html.Div([
                            html.Label([
                                html.I(className="fas fa-map-marker-alt", style={
                                    "marginRight": "6px",
                                    "color": "#6b7280",
                                    "fontSize": "12px"
                                }),
                                "Region"
                            ], style={
                                "fontSize": "13px",
                                "fontWeight": "600",
                                "color": "#374151",
                                "marginBottom": "8px",
                                "display": "block"
                            }),
                            dcc.Dropdown(
                                id='region-filter',
                                options=[{'label': r, 'value': r} for r in regions],
                                value='All Regions',
                                clearable=False,
                                style={"fontSize": "13px"},
                                optionHeight=40,
                                className="modern-dropdown"
                            )
                        ], style={"marginBottom": "15px"}),

                        # Banker Filter with Icon
                        html.Div([
                            html.Label([
                                html.I(className="fas fa-user-tie", style={
                                    "marginRight": "6px",
                                    "color": "#6b7280",
                                    "fontSize": "12px"
                                }),
                                "Banker"
                            ], style={
                                "fontSize": "13px",
                                "fontWeight": "600",
                                "color": "#374151",
                                "marginBottom": "8px",
                                "display": "block"
                            }),
                            dcc.Dropdown(
                                id='banker-filter',
                                options=[{'label': 'All Bankers', 'value': 'All Bankers'}],
                                value='All Bankers',
                                clearable=False,
                                style={"fontSize": "13px"},
                                optionHeight=40,
                                className="modern-dropdown"
                            )
                        ], style={"marginBottom": "15px"}),

                        # Elegant Divider
                        html.Div(style={
                            "height": "1px",
                            "background": "linear-gradient(90deg, transparent, #e5e7eb, transparent)",
                            "margin": "20px 0"
                        }),

                        # Investment Strategist with Badge
                        html.Div([
                            html.Label([
                                html.I(className="fas fa-user-chart", style={
                                    "marginRight": "6px",
                                    "color": "#6b7280",
                                    "fontSize": "12px"
                                }),
                                "Investment Strategist",
                                html.Span("INDEPENDENT", style={
                                    "marginLeft": "8px",
                                    "fontSize": "9px",
                                    "fontWeight": "700",
                                    "color": "#10b981",
                                    "background": "rgba(16, 185, 129, 0.1)",
                                    "padding": "2px 6px",
                                    "borderRadius": "8px",
                                    "letterSpacing": "0.5px"
                                })
                            ], style={
                                "fontSize": "13px",
                                "fontWeight": "600",
                                "color": "#374151",
                                "marginBottom": "8px",
                                "display": "block"
                            }),
                            dcc.Dropdown(
                                id='strategist-filter',
                                options=[{'label': s, 'value': s} for s in all_strategists],
                                value='All Strategists',
                                clearable=False,
                                style={"fontSize": "13px"},
                                optionHeight=40,
                                className="modern-dropdown"
                            )
                        ], style={"marginBottom": "20px"}),

                        # Modern Back Button - More compact
                        dbc.Button([
                            html.I(className="fas fa-arrow-left", style={"marginRight": "6px"}),
                            "Back to Main"
                        ], id="back-to-main-dash",
                        color="light",
                        style={
                            "width": "100%",
                            "borderRadius": "10px",
                            "fontWeight": "600",
                            "border": "2px solid #e5e7eb",
                            "color": "#374151",
                            "transition": "all 0.3s ease",
                            "fontSize": "13px",
                            "padding": "8px 16px"
                        })
                    ], className="glass-card", style={
                        "padding": "20px",
                        "borderRadius": "12px",
                        "position": "sticky",
                        "top": "20px"
                    })
                ], width=3),

                # RIGHT SIDE - Enhanced Dashboard
                dbc.Col([
                    html.Div(id='dashboard-kpi-content', className="depth-card", style={
                        "padding": "20px",
                        "borderRadius": "12px",
                        "background": "white"
                    })
                ], width=9)
            ], style={"marginTop": "10px"})

        ], fluid=True, style={"padding": "20px"}),

        dcc.Store(id='region-banker-map', data=region_banker_map)
    ], className="animated-gradient", style={"minHeight": "100vh", "paddingTop": "5px", "paddingBottom": "5px"})


# ============================================================================
# UPDATED KEY MATRIX LAYOUT - Replace get_key_matrix_layout function
# ============================================================================
# Find this function (around line 1100) and replace it entirely

def get_key_matrix_layout():
    """Enhanced Key Matrix with Export functionality"""
    df = load_key_matrix_data()

    # Extract unique values for filters
    regions = ['All Regions']
    region_banker_map = {'All Regions': ['All Bankers']}
    all_strategists = ['All Strategists']

    if not df.empty:
        if 'Investment_Strategist' in df.columns:
            unique_strategists = sorted(
                df[df['Investment_Strategist'].notna()]['Investment_Strategist'].unique().tolist())
            all_strategists.extend([s for s in unique_strategists if s != 'All Strategists'])

        if 'Region' in df.columns:
            unique_regions = sorted(df[df['Region'].notna()]['Region'].unique().tolist())
            regions.extend([r for r in unique_regions if r != 'All Regions'])

            for region in unique_regions:
                if 'Banker' in df.columns:
                    bankers = sorted(df[df['Region'] == region]['Banker'].dropna().unique().tolist())
                    region_banker_map[region] = ['All Bankers'] + [b for b in bankers if b != 'All Bankers']

    return html.Div([
        dbc.Container([
            # Enhanced Header with Export Dropdown
            dbc.Row([
                dbc.Col([
                    html.Div([
                        html.H1("Key Matrix", style={
                            "background": "linear-gradient(135deg, #C41E3A 0%, #8b1528 100%)",
                            "-webkit-background-clip": "text",
                            "-webkit-text-fill-color": "transparent",
                            "fontWeight": "800",
                            "fontSize": "26px",
                            "margin": "15px 0 8px 0",
                            "letterSpacing": "-0.5px",
                            "display": "inline-block"
                        }),
                        html.Span("📊", style={"fontSize": "24px", "marginLeft": "10px"})
                    ])
                ], width=6),
                dbc.Col([
                    html.Div([
                        # Export Dropdown Button
                        dbc.DropdownMenu(
                            [
                                dbc.DropdownMenuItem([
                                    html.I(className="fas fa-file-excel", style={"marginRight": "8px"}),
                                    "Export as Excel"
                                ], id="km-export-excel-btn", className="dropdown-item-custom"),
                                dbc.DropdownMenuItem([
                                    html.I(className="fas fa-file-pdf", style={"marginRight": "8px"}),
                                    "Export as PDF"
                                ], id="km-export-pdf-btn", className="dropdown-item-custom"),
                                dbc.DropdownMenuItem([
                                    html.I(className="fas fa-file-powerpoint", style={"marginRight": "8px"}),
                                    "Export as PowerPoint"
                                ], id="km-export-ppt-btn", className="dropdown-item-custom"),
                            ],
                            label=[
                                html.I(className="fas fa-download", style={"marginRight": "6px"}),
                                "Export"
                            ],
                            color="light",
                            size="sm",
                            className="me-2",
                            style={"display": "inline-block", "fontSize": "12px"}
                        ),

                        # # Refresh Button
                        # dbc.Button([
                        #     html.I(className="fas fa-sync-alt", style={"marginRight": "6px"}),
                        #     "Refresh"
                        # ], id="km-refresh-btn", color="light", size="sm",
                        #     className="me-2", style={"fontSize": "12px"}),

                        # Back Button
                        dbc.Button([
                            html.I(className="fas fa-arrow-left", style={"marginRight": "6px"}),
                            "Back"
                        ], id="back-to-dashboard-matrix", color="secondary", size="sm",
                            style={"fontSize": "12px"})
                    ], style={"textAlign": "right"})
                ], width=6, style={"paddingTop": "18px"})
            ]),

            html.Hr(style={"margin": "8px 0", "borderColor": "#e0e0e0"}),

            # Main content with filters
            dbc.Row([
                # Left Sidebar - Enhanced with collapsible sections
                dbc.Col([
                    html.Div([
                        # Filter Header with collapse button
                        html.Div([
                            html.Div([
                                html.I(className="fas fa-filter", style={
                                    "color": "#C41E3A",
                                    "marginRight": "8px",
                                    "fontSize": "14px"
                                }),
                                html.H5("Filters", style={
                                    "display": "inline",
                                    "fontWeight": "700",
                                    "fontSize": "15px",
                                    "margin": "0"
                                })
                            ], style={"display": "inline-block"}),
                            html.Button([
                                html.I(id="km-filter-collapse-icon", className="fas fa-chevron-up")
                            ], id="km-filter-collapse-btn", style={
                                "border": "none",
                                "background": "transparent",
                                "float": "right",
                                "cursor": "pointer",
                                "color": "#666"
                            })
                        ], style={
                            "marginBottom": "12px",
                            "paddingBottom": "8px",
                            "borderBottom": "2px solid #e5e7eb"
                        }),

                        # Collapsible Filter Content
                        dbc.Collapse([
                            # Region Filter
                            html.Div([
                                html.Label([
                                    html.I(className="fas fa-map-marker-alt", style={
                                        "marginRight": "6px",
                                        "color": "#6b7280",
                                        "fontSize": "11px"
                                    }),
                                    "Region"
                                ], style={
                                    "fontSize": "12px",
                                    "fontWeight": "600",
                                    "marginBottom": "6px",
                                    "display": "block",
                                    "color": "#374151"
                                }),
                                dcc.Dropdown(
                                    id='km-region-filter',
                                    options=[{'label': r, 'value': r} for r in regions],
                                    value='All Regions',
                                    clearable=False,
                                    style={"fontSize": "12px", "minHeight": "36px"},
                                    optionHeight=35
                                )
                            ], style={"marginBottom": "12px"}),

                            # Banker Filter
                            html.Div([
                                html.Label([
                                    html.I(className="fas fa-user-tie", style={
                                        "marginRight": "6px",
                                        "color": "#6b7280",
                                        "fontSize": "11px"
                                    }),
                                    "Banker"
                                ], style={
                                    "fontSize": "12px",
                                    "fontWeight": "600",
                                    "marginBottom": "6px",
                                    "display": "block",
                                    "color": "#374151"
                                }),
                                dcc.Dropdown(
                                    id='km-banker-filter',
                                    options=[{'label': 'All Bankers', 'value': 'All Bankers'}],
                                    value='All Bankers',
                                    clearable=False,
                                    style={"fontSize": "12px", "minHeight": "36px"},
                                    optionHeight=35
                                )
                            ], style={"marginBottom": "12px"}),

                            html.Div(style={
                                "height": "1px",
                                "background": "linear-gradient(90deg, transparent, #e5e7eb, transparent)",
                                "margin": "15px 0"
                            }),

                            # Strategist Filter
                            html.Div([
                                html.Label([
                                    html.I(className="fas fa-chart-line", style={
                                        "marginRight": "6px",
                                        "color": "#6b7280",
                                        "fontSize": "11px"
                                    }),
                                    "Investment Strategist ",
                                    html.Span("INDEPENDENT", style={
                                        "fontSize": "9px",
                                        "fontWeight": "700",
                                        "color": "#10b981",
                                        "background": "rgba(16, 185, 129, 0.1)",
                                        "padding": "2px 6px",
                                        "borderRadius": "8px",
                                        "letterSpacing": "0.5px"
                                    })
                                ], style={
                                    "fontSize": "12px",
                                    "fontWeight": "600",
                                    "marginBottom": "6px",
                                    "display": "block",
                                    "color": "#374151"
                                }),
                                dcc.Dropdown(
                                    id='km-strategist-filter',
                                    options=[{'label': s, 'value': s} for s in all_strategists],
                                    value='All Strategists',
                                    clearable=False,
                                    style={"fontSize": "12px", "minHeight": "36px"},
                                    optionHeight=35
                                )
                            ], style={"marginBottom": "15px"}),
                        ], id="km-filter-collapse", is_open=True),

                        html.Div(style={
                            "height": "1px",
                            "background": "linear-gradient(90deg, transparent, #e5e7eb, transparent)",
                            "margin": "15px 0"
                        }),

                        # Comparison Section - Collapsible
                        html.Div([
                            html.Div([
                                html.I(className="fas fa-balance-scale", style={
                                    "color": "#C41E3A",
                                    "marginRight": "8px",
                                    "fontSize": "14px"
                                }),
                                html.H5("Compare", style={
                                    "display": "inline",
                                    "fontWeight": "700",
                                    "fontSize": "15px",
                                    "margin": "0"
                                }),
                                html.Button([
                                    html.I(id="km-compare-collapse-icon", className="fas fa-chevron-down")
                                ], id="km-compare-collapse-btn", style={
                                    "border": "none",
                                    "background": "transparent",
                                    "float": "right",
                                    "cursor": "pointer",
                                    "color": "#666"
                                })
                            ], style={"marginBottom": "12px"}),

                            dbc.Collapse([
                                html.Label("Compare Multiple", style={
                                    "fontSize": "12px",
                                    "fontWeight": "600",
                                    "color": "#374151",
                                    "marginBottom": "6px",
                                    "display": "block"
                                }),
                                dcc.Dropdown(
                                    id='km-compare-type',
                                    options=[
                                        {'label': '🗺️ Compare Regions', 'value': 'region'},
                                        {'label': '👔 Compare Bankers', 'value': 'banker'},
                                        {'label': '📈 Compare Strategists', 'value': 'strategist'}
                                    ],
                                    placeholder='Select comparison type',
                                    clearable=True,
                                    style={"fontSize": "12px", "marginBottom": "8px", "minHeight": "36px"},
                                    optionHeight=35
                                ),
                                dcc.Dropdown(
                                    id='km-compare-values',
                                    options=[],
                                    placeholder='Select 2-5 items',
                                    multi=True,
                                    clearable=True,
                                    style={"fontSize": "12px", "minHeight": "36px"},
                                    optionHeight=35
                                )
                            ], id="km-compare-collapse", is_open=False)
                        ], style={"marginBottom": "15px"})
                    ], className="glass-card", style={
                        "padding": "18px",
                        "borderRadius": "12px",
                        "position": "sticky",
                        "top": "20px",
                        "maxHeight": "calc(100vh - 150px)",
                        "overflowY": "auto"
                    })
                ], width=3),

                # Right Side - Enhanced Content with Tabs
                dbc.Col([
                    # View Toggle Tabs
                    dbc.Tabs([
                        dbc.Tab(label="📊 Dashboard View", tab_id="km-tab-dashboard",
                                label_style={"fontSize": "12px", "fontWeight": "600"}),
                        dbc.Tab(label="📈 Trend Analysis", tab_id="km-tab-trends",
                                label_style={"fontSize": "12px", "fontWeight": "600"}),
                        dbc.Tab(label="📋 Data Table", tab_id="km-tab-table",
                                label_style={"fontSize": "12px", "fontWeight": "600"}),
                    ], id="km-tabs", active_tab="km-tab-dashboard", style={"marginBottom": "15px"}),

                    html.Div(id='km-content', style={
                        "padding": "15px",
                        "minHeight": "calc(100vh - 250px)"
                    })
                ], width=9)
            ], style={"marginTop": "10px", "alignItems": "stretch"})
        ], fluid=True),

        # Store components
        dcc.Store(id='km-region-banker-map', data=region_banker_map),
        dcc.Store(id='km-last-update', data=datetime.now().strftime('%Y-%m-%d %H:%M:%S')),

        # Download component for all exports
        dcc.Download(id="km-download-dataframe")

    ], className="animated-gradient", style={"minHeight": "100vh", "padding": "20px"})

def get_x_sell_dashboard_layout():
    """Enhanced X-Sell - Removed search box and non-functional elements"""
    df = load_crosssell_detailed_data()

    regions = ['All Regions']
    region_banker_map = {'All Regions': ['All Bankers']}
    all_strategists = ['All Strategists']
    all_products = ['All Products']
    all_stages = ['All Stages']

    if not df.empty:
        if 'Investment_Strategist' in df.columns:
            unique_strategists = sorted(
                df[df['Investment_Strategist'].notna()]['Investment_Strategist'].unique().tolist())
            all_strategists.extend([s for s in unique_strategists if s != 'All Strategists'])

        if 'Region' in df.columns:
            unique_regions = sorted(df[df['Region'].notna()]['Region'].unique().tolist())
            regions.extend([r for r in unique_regions if r != 'All Regions'])

            for region in unique_regions:
                if 'Banker' in df.columns:
                    bankers = sorted(df[df['Region'] == region]['Banker'].dropna().unique().tolist())
                    region_banker_map[region] = ['All Bankers'] + [b for b in bankers if b != 'All Bankers']

        if 'Product' in df.columns:
            unique_products = sorted(df['Product'].dropna().unique().tolist())
            all_products.extend(unique_products)

        if 'Stage' in df.columns:
            unique_stages = sorted(df['Stage'].dropna().unique().tolist())
            all_stages.extend(unique_stages)

    return html.Div([
        dbc.Container([
            # Enhanced Header - REMOVED SEARCH BOX
            # In get_x_sell_dashboard_layout() function, replace the header row with:

            # Enhanced Header with Export Dropdown
            dbc.Row([
                dbc.Col([
                    html.Div([
                        html.H1("X Sell", style={
                            "background": "linear-gradient(135deg, #C41E3A 0%, #8b1528 100%)",
                            "-webkit-background-clip": "text",
                            "-webkit-text-fill-color": "transparent",
                            "fontWeight": "800",
                            "fontSize": "26px",
                            "margin": "15px 0 8px 0",
                            "letterSpacing": "-0.5px",
                            "display": "inline-block"
                        }),
                        html.Span("📄", style={"fontSize": "24px", "marginLeft": "10px"})
                    ])
                ], width=6),
                dbc.Col([
                    html.Div([
                        # Export Dropdown Menu
                        dbc.DropdownMenu(
                            [
                                dbc.DropdownMenuItem([
                                    html.I(className="fas fa-file-excel", style={"marginRight": "8px"}),
                                    "Export as Excel"
                                ], id="xs-export-excel-btn", className="dropdown-item-custom"),
                                dbc.DropdownMenuItem([
                                    html.I(className="fas fa-file-pdf", style={"marginRight": "8px"}),
                                    "Export as PDF"
                                ], id="xs-export-pdf-btn", className="dropdown-item-custom"),
                                dbc.DropdownMenuItem([
                                    html.I(className="fas fa-file-powerpoint", style={"marginRight": "8px"}),
                                    "Export as PowerPoint"
                                ], id="xs-export-ppt-btn", className="dropdown-item-custom"),
                            ],
                            label=[
                                html.I(className="fas fa-download", style={"marginRight": "6px"}),
                                "Export"
                            ],
                            color="light",
                            size="sm",
                            className="me-2",
                            style={"display": "inline-block", "fontSize": "11px"}
                        ),

                        # Back Button
                        dbc.Button([
                            html.I(className="fas fa-arrow-left", style={"marginRight": "5px", "fontSize": "10px"}),
                            "Back"
                        ], id="back-to-dashboard-xsell", color="secondary", size="sm",
                            style={"fontSize": "11px", "padding": "6px 14px"})
                    ], style={"textAlign": "right"})
                ], width=6, style={"paddingTop": "18px"})
            ]),

            html.Hr(style={"margin": "8px 0", "borderColor": "#e0e0e0"}),

            dbc.Row([
                # Left Sidebar - CLEANED UP FILTERS
                dbc.Col([
                    html.Div([
                        # Filter Header
                        html.Div([
                            html.Div([
                                html.I(className="fas fa-filter", style={
                                    "color": "#C41E3A",
                                    "marginRight": "8px",
                                    "fontSize": "14px"
                                }),
                                html.H5("Filters", style={
                                    "display": "inline",
                                    "fontWeight": "700",
                                    "fontSize": "15px",
                                    "margin": "0"
                                })
                            ], style={"display": "inline-block"}),
                            html.Button([
                                html.I(id="xs-filter-collapse-icon", className="fas fa-chevron-up")
                            ], id="xs-filter-collapse-btn", style={
                                "border": "none",
                                "background": "transparent",
                                "float": "right",
                                "cursor": "pointer",
                                "color": "#666"
                            })
                        ], style={
                            "marginBottom": "12px",
                            "paddingBottom": "8px",
                            "borderBottom": "2px solid #e5e7eb"
                        }),

                        # Collapsible Filter Content
                        dbc.Collapse([
                            # Region Filter
                            html.Div([
                                html.Label([
                                    html.I(className="fas fa-map-marker-alt", style={
                                        "marginRight": "6px",
                                        "color": "#6b7280",
                                        "fontSize": "11px"
                                    }),
                                    "Region"
                                ], style={
                                    "fontSize": "12px",
                                    "fontWeight": "600",
                                    "marginBottom": "6px",
                                    "display": "block",
                                    "color": "#374151"
                                }),
                                dcc.Dropdown(
                                    id='xs-region-filter',
                                    options=[{'label': r, 'value': r} for r in regions],
                                    value='All Regions',
                                    clearable=False,
                                    style={"fontSize": "12px", "minHeight": "36px"},
                                    optionHeight=35
                                )
                            ], style={"marginBottom": "12px"}),

                            # Banker Filter
                            html.Div([
                                html.Label([
                                    html.I(className="fas fa-user-tie", style={
                                        "marginRight": "6px",
                                        "color": "#6b7280",
                                        "fontSize": "11px"
                                    }),
                                    "Banker"
                                ], style={
                                    "fontSize": "12px",
                                    "fontWeight": "600",
                                    "marginBottom": "6px",
                                    "display": "block",
                                    "color": "#374151"
                                }),
                                dcc.Dropdown(
                                    id='xs-banker-filter',
                                    options=[{'label': 'All Bankers', 'value': 'All Bankers'}],
                                    value='All Bankers',
                                    clearable=False,
                                    style={"fontSize": "12px", "minHeight": "36px"},
                                    optionHeight=35
                                )
                            ], style={"marginBottom": "12px"}),

                            # Product Filter
                            html.Div([
                                html.Label([
                                    html.I(className="fas fa-box", style={
                                        "marginRight": "6px",
                                        "color": "#6b7280",
                                        "fontSize": "11px"
                                    }),
                                    "Product"
                                ], style={
                                    "fontSize": "12px",
                                    "fontWeight": "600",
                                    "marginBottom": "6px",
                                    "display": "block",
                                    "color": "#374151"
                                }),
                                dcc.Dropdown(
                                    id='xs-product-filter',
                                    options=[{'label': p, 'value': p} for p in all_products],
                                    value='All Products',
                                    clearable=False,
                                    style={"fontSize": "12px", "minHeight": "36px"},
                                    optionHeight=35
                                )
                            ], style={"marginBottom": "12px"}),

                            # Stage Filter
                            html.Div([
                                html.Label([
                                    html.I(className="fas fa-tasks", style={
                                        "marginRight": "6px",
                                        "color": "#6b7280",
                                        "fontSize": "11px"
                                    }),
                                    "Stage"
                                ], style={
                                    "fontSize": "12px",
                                    "fontWeight": "600",
                                    "marginBottom": "6px",
                                    "display": "block",
                                    "color": "#374151"
                                }),
                                dcc.Dropdown(
                                    id='xs-stage-filter',
                                    options=[{'label': s, 'value': s} for s in all_stages],
                                    value='All Stages',
                                    clearable=False,
                                    style={"fontSize": "12px", "minHeight": "36px"},
                                    optionHeight=35
                                )
                            ], style={"marginBottom": "12px"}),

                            html.Div(style={
                                "height": "1px",
                                "background": "linear-gradient(90deg, transparent, #e5e7eb, transparent)",
                                "margin": "15px 0"
                            }),

                            # Strategist Filter
                            html.Div([
                                html.Label([
                                    html.I(className="fas fa-chart-line", style={
                                        "marginRight": "6px",
                                        "color": "#6b7280",
                                        "fontSize": "11px"
                                    }),
                                    "Investment Strategist ",
                                    html.Span("INDEPENDENT", style={
                                        "fontSize": "9px",
                                        "color": "#10b981",
                                        "background": "rgba(16, 185, 129, 0.1)",
                                        "padding": "2px 6px",
                                        "borderRadius": "8px",
                                        "fontWeight": "700",
                                        "letterSpacing": "0.5px"
                                    })
                                ], style={
                                    "fontSize": "12px",
                                    "fontWeight": "600",
                                    "marginBottom": "6px",
                                    "display": "block",
                                    "color": "#374151"
                                }),
                                dcc.Dropdown(
                                    id='xs-strategist-filter',
                                    options=[{'label': s, 'value': s} for s in all_strategists],
                                    value='All Strategists',
                                    clearable=False,
                                    style={"fontSize": "12px", "minHeight": "36px"},
                                    optionHeight=35
                                )
                            ], style={"marginBottom": "15px"}),
                        ], id="xs-filter-collapse", is_open=True),

                        html.Div(style={
                            "height": "1px",
                            "background": "linear-gradient(90deg, transparent, #e5e7eb, transparent)",
                            "margin": "15px 0"
                        }),

                        # Stage Progress Indicator
                        html.Div([
                            html.Div([
                                html.I(className="fas fa-chart-pie", style={
                                    "color": "#C41E3A",
                                    "marginRight": "8px",
                                    "fontSize": "14px"
                                }),
                                html.H5("Stage Progress", style={
                                    "display": "inline",
                                    "fontWeight": "700",
                                    "fontSize": "15px",
                                    "margin": "0"
                                })
                            ], style={"marginBottom": "12px"}),

                            html.Div(id="xs-stage-progress")
                        ]),

                        html.Div(style={
                            "height": "1px",
                            "background": "linear-gradient(90deg, transparent, #e5e7eb, transparent)",
                            "margin": "15px 0"
                        }),

                        # Comparison Section
                        html.Div([
                            html.Div([
                                html.I(className="fas fa-balance-scale", style={
                                    "color": "#C41E3A",
                                    "marginRight": "8px",
                                    "fontSize": "14px"
                                }),
                                html.H5("Compare", style={
                                    "display": "inline",
                                    "fontWeight": "700",
                                    "fontSize": "15px",
                                    "margin": "0"
                                }),
                                html.Button([
                                    html.I(id="xs-compare-collapse-icon", className="fas fa-chevron-down")
                                ], id="xs-compare-collapse-btn", style={
                                    "border": "none",
                                    "background": "transparent",
                                    "float": "right",
                                    "cursor": "pointer",
                                    "color": "#666"
                                })
                            ], style={"marginBottom": "12px"}),

                            dbc.Collapse([
                                html.Label("Compare Multiple", style={
                                    "fontSize": "12px",
                                    "fontWeight": "600",
                                    "color": "#374151",
                                    "display": "block",
                                    "marginBottom": "6px"
                                }),
                                dcc.Dropdown(
                                    id='xs-compare-type',
                                    options=[
                                        {'label': '🗺️ Compare Regions', 'value': 'region'},
                                        {'label': '👔 Compare Bankers', 'value': 'banker'},
                                        {'label': '📈 Compare Strategists', 'value': 'strategist'}
                                    ],
                                    placeholder='Select type',
                                    clearable=True,
                                    style={"fontSize": "12px", "marginBottom": "8px", "minHeight": "36px"},
                                    optionHeight=35
                                ),
                                dcc.Dropdown(
                                    id='xs-compare-values',
                                    options=[],
                                    placeholder='Select 2-5 items',
                                    multi=True,
                                    clearable=True,
                                    style={"fontSize": "12px", "minHeight": "36px"},
                                    optionHeight=35
                                )
                            ], id="xs-compare-collapse", is_open=False)
                        ])
                    ], className="glass-card", style={
                        "padding": "18px",
                        "borderRadius": "12px",
                        "position": "sticky",
                        "top": "20px",
                        "maxHeight": "calc(100vh - 150px)",
                        "overflowY": "auto"
                    })
                ], width=3),

                # Right Side - Content Area (NO TABS - DIRECT TABLE)
                dbc.Col([
                    html.Div(id='xs-content', style={
                        "padding": "15px",
                        "minHeight": "calc(100vh - 250px)"
                    })
                ], width=9)
            ], style={"alignItems": "stretch"})
        ], fluid=True),

        # Store components
        dcc.Store(id='xs-region-banker-map', data=region_banker_map),
        # Download component
        dcc.Download(id="xs-download-dataframe"),

        # Interval for auto-refresh (optional)
        dcc.Interval(id='xs-auto-refresh', interval=300000, n_intervals=0, disabled=True)

    ], style={"minHeight": "100vh", "padding": "20px"})
def get_new_mandates_layout():
    """Enhanced Mandates with interactive features"""
    df = load_mandate_detailed_data()

    regions = ['All Regions']
    region_banker_map = {'All Regions': ['All Bankers']}
    all_strategists = ['All Strategists']

    if not df.empty:
        if 'Investment_Strategist' in df.columns:
            unique_strategists = sorted(
                df[df['Investment_Strategist'].notna()]['Investment_Strategist'].unique().tolist())
            all_strategists.extend([s for s in unique_strategists if s != 'All Strategists'])

        if 'Region' in df.columns:
            unique_regions = sorted(df[df['Region'].notna()]['Region'].unique().tolist())
            regions.extend([r for r in unique_regions if r != 'All Regions'])

            for region in unique_regions:
                if 'Banker' in df.columns:
                    bankers = sorted(df[df['Region'] == region]['Banker'].dropna().unique().tolist())
                    region_banker_map[region] = ['All Bankers'] + [b for b in bankers if b != 'All Bankers']

    return html.Div([
        dbc.Container([
            # Enhanced Header
            dbc.Row([
                dbc.Col([
                    html.Div([
                        html.H1("New Mandates", style={
                            "background": "linear-gradient(135deg, #C41E3A 0%, #8b1528 100%)",
                            "-webkit-background-clip": "text",
                            "-webkit-text-fill-color": "transparent",
                            "fontWeight": "800",
                            "fontSize": "26px",
                            "margin": "15px 0 8px 0",
                            "letterSpacing": "-0.5px",
                            "display": "inline-block"
                        }),
                        html.Span("💼", style={"fontSize": "24px", "marginLeft": "10px"})
                    ])
                ], width=6),
                dbc.Col([
                    html.Div([
                        # Export Dropdown Menu
                        dbc.DropdownMenu(
                            [
                                dbc.DropdownMenuItem([
                                    html.I(className="fas fa-file-excel", style={"marginRight": "8px"}),
                                    "Export as Excel"
                                ], id="md-export-excel-btn", className="dropdown-item-custom"),
                                dbc.DropdownMenuItem([
                                    html.I(className="fas fa-file-pdf", style={"marginRight": "8px"}),
                                    "Export as PDF"
                                ], id="md-export-pdf-btn", className="dropdown-item-custom"),
                                dbc.DropdownMenuItem([
                                    html.I(className="fas fa-file-powerpoint", style={"marginRight": "8px"}),
                                    "Export as PowerPoint"
                                ], id="md-export-ppt-btn", className="dropdown-item-custom"),
                            ],
                            label=[
                                html.I(className="fas fa-download", style={"marginRight": "6px"}),
                                "Export"
                            ],
                            color="light",
                            size="sm",
                            className="me-2",
                            style={"display": "inline-block", "fontSize": "12px"}
                        ),

                        # Back Button
                        dbc.Button([
                            html.I(className="fas fa-arrow-left", style={"marginRight": "6px"}),
                            "Back"
                        ], id="back-to-dashboard-mandates", color="secondary", size="sm",
                            style={"fontSize": "12px"})
                    ], style={"textAlign": "right"})
                ], width=6, style={"paddingTop": "18px"})
            ]),

            html.Hr(style={"margin": "8px 0", "borderColor": "#e0e0e0"}),

            dbc.Row([
                # Left Sidebar - Enhanced
                dbc.Col([
                    html.Div([
                        # Filter Header
                        html.Div([
                            html.Div([
                                html.I(className="fas fa-filter", style={
                                    "color": "#C41E3A",
                                    "marginRight": "8px",
                                    "fontSize": "14px"
                                }),
                                html.H5("Filters", style={
                                    "display": "inline",
                                    "fontWeight": "700",
                                    "fontSize": "15px",
                                    "margin": "0"
                                })
                            ], style={"display": "inline-block"}),
                            html.Button([
                                html.I(id="md-filter-collapse-icon", className="fas fa-chevron-up")
                            ], id="md-filter-collapse-btn", style={
                                "border": "none",
                                "background": "transparent",
                                "float": "right",
                                "cursor": "pointer",
                                "color": "#666"
                            })
                        ], style={
                            "marginBottom": "12px",
                            "paddingBottom": "8px",
                            "borderBottom": "2px solid #e5e7eb"
                        }),

                        # Collapsible Filter Content
                        dbc.Collapse([
                            html.Div([
                                html.Label([
                                    html.I(className="fas fa-map-marker-alt", style={
                                        "marginRight": "6px",
                                        "color": "#6b7280",
                                        "fontSize": "11px"
                                    }),
                                    "Region"
                                ], style={
                                    "fontSize": "12px",
                                    "fontWeight": "600",
                                    "marginBottom": "6px",
                                    "display": "block",
                                    "color": "#374151"
                                }),
                                dcc.Dropdown(
                                    id='md-region-filter',
                                    options=[{'label': r, 'value': r} for r in regions],
                                    value='All Regions',
                                    clearable=False,
                                    style={"fontSize": "12px", "minHeight": "36px"},
                                    optionHeight=35
                                )
                            ], style={"marginBottom": "12px"}),

                            html.Div([
                                html.Label([
                                    html.I(className="fas fa-user-tie", style={
                                        "marginRight": "6px",
                                        "color": "#6b7280",
                                        "fontSize": "11px"
                                    }),
                                    "Banker"
                                ], style={
                                    "fontSize": "12px",
                                    "fontWeight": "600",
                                    "marginBottom": "6px",
                                    "display": "block",
                                    "color": "#374151"
                                }),
                                dcc.Dropdown(
                                    id='md-banker-filter',
                                    options=[{'label': 'All Bankers', 'value': 'All Bankers'}],
                                    value='All Bankers',
                                    clearable=False,
                                    style={"fontSize": "12px", "minHeight": "36px"},
                                    optionHeight=35
                                )
                            ], style={"marginBottom": "12px"}),

                            html.Div(style={
                                "height": "1px",
                                "background": "linear-gradient(90deg, transparent, #e5e7eb, transparent)",
                                "margin": "15px 0"
                            }),

                            html.Div([
                                html.Label([
                                    html.I(className="fas fa-chart-line", style={
                                        "marginRight": "6px",
                                        "color": "#6b7280",
                                        "fontSize": "11px"
                                    }),
                                    "Investment Strategist ",
                                    html.Span("INDEPENDENT", style={
                                        "fontSize": "9px",
                                        "color": "#10b981",
                                        "background": "rgba(16, 185, 129, 0.1)",
                                        "padding": "2px 6px",
                                        "borderRadius": "8px",
                                        "fontWeight": "700",
                                        "letterSpacing": "0.5px"
                                    })
                                ], style={
                                    "fontSize": "12px",
                                    "fontWeight": "600",
                                    "marginBottom": "6px",
                                    "display": "block",
                                    "color": "#374151"
                                }),
                                dcc.Dropdown(
                                    id='md-strategist-filter',
                                    options=[{'label': s, 'value': s} for s in all_strategists],
                                    value='All Strategists',
                                    clearable=False,
                                    style={"fontSize": "12px", "minHeight": "36px"},
                                    optionHeight=35
                                )
                            ], style={"marginBottom": "15px"}),
                        ], id="md-filter-collapse", is_open=True),

                        html.Div(style={
                            "height": "1px",
                            "background": "linear-gradient(90deg, transparent, #e5e7eb, transparent)",
                            "margin": "15px 0"
                        }),

                        # Conversion Funnel - NEW
                        html.Div([
                            html.Div([
                                html.I(className="fas fa-funnel-dollar", style={
                                    "color": "#C41E3A",
                                    "marginRight": "8px",
                                    "fontSize": "14px"
                                }),
                                html.H5("Conversion Funnel", style={
                                    "display": "inline",
                                    "fontWeight": "700",
                                    "fontSize": "15px",
                                    "margin": "0"
                                })
                            ], style={"marginBottom": "12px"}),

                            html.Div(id="md-conversion-funnel")
                        ]),

                        html.Div(style={
                            "height": "1px",
                            "background": "linear-gradient(90deg, transparent, #e5e7eb, transparent)",
                            "margin": "15px 0"
                        }),

                        # Win/Loss Analysis - NEW
                        html.Div([
                            html.Div([
                                html.I(className="fas fa-trophy", style={
                                    "color": "#C41E3A",
                                    "marginRight": "8px",
                                    "fontSize": "14px"
                                }),
                                html.H5("Win/Loss", style={
                                    "display": "inline",
                                    "fontWeight": "700",
                                    "fontSize": "15px",
                                    "margin": "0"
                                })
                            ], style={"marginBottom": "12px"}),

                            html.Div(id="md-win-loss")
                        ]),

                        html.Div(style={
                            "height": "1px",
                            "background": "linear-gradient(90deg, transparent, #e5e7eb, transparent)",
                            "margin": "15px 0"
                        }),

                        # Comparison Section
                        html.Div([
                            html.Div([
                                html.I(className="fas fa-balance-scale", style={
                                    "color": "#C41E3A",
                                    "marginRight": "8px",
                                    "fontSize": "14px"
                                }),
                                html.H5("Compare", style={
                                    "display": "inline",
                                    "fontWeight": "700",
                                    "fontSize": "15px",
                                    "margin": "0"
                                }),
                                html.Button([
                                    html.I(id="md-compare-collapse-icon", className="fas fa-chevron-down")
                                ], id="md-compare-collapse-btn", style={
                                    "border": "none",
                                    "background": "transparent",
                                    "float": "right",
                                    "cursor": "pointer",
                                    "color": "#666"
                                })
                            ], style={"marginBottom": "12px"}),

                            dbc.Collapse([
                                html.Label("Compare Multiple", style={
                                    "fontSize": "12px",
                                    "fontWeight": "600",
                                    "color": "#374151",
                                    "display": "block",
                                    "marginBottom": "6px"
                                }),
                                dcc.Dropdown(
                                    id='md-compare-type',
                                    options=[
                                        {'label': '📍 Compare Regions', 'value': 'region'},
                                        {'label': '👔 Compare Bankers', 'value': 'banker'},
                                        {'label': '📈 Compare Strategists', 'value': 'strategist'}
                                    ],
                                    placeholder='Select type',
                                    clearable=True,
                                    style={"fontSize": "12px", "marginBottom": "8px", "minHeight": "36px"},
                                    optionHeight=35
                                ),
                                dcc.Dropdown(
                                    id='md-compare-values',
                                    options=[],
                                    placeholder='Select 2-5 items',
                                    multi=True,
                                    clearable=True,
                                    style={"fontSize": "12px", "minHeight": "36px"},
                                    optionHeight=35
                                )
                            ], id="md-compare-collapse", is_open=False)
                        ])
                    ], className="glass-card", style={
                        "padding": "18px",
                        "borderRadius": "12px",
                        "position": "sticky",
                        "top": "20px",
                        "maxHeight": "calc(100vh - 150px)",
                        "overflowY": "auto"
                    })
                ], width=3),

                # Right Side - Enhanced with Tabs
                dbc.Col([
                    # View Toggle Tabs
                    dbc.Tabs([
                        dbc.Tab(label="📊 Timeline View", tab_id="md-tab-timeline",
                                label_style={"fontSize": "12px", "fontWeight": "600"}),
                        dbc.Tab(label="📈 Trend Analysis", tab_id="md-tab-trends",
                                label_style={"fontSize": "12px", "fontWeight": "600"}),
                        dbc.Tab(label="📝 Remarks & Notes", tab_id="md-tab-remarks",
                                label_style={"fontSize": "12px", "fontWeight": "600"}),
                    ], id="md-tabs", active_tab="md-tab-timeline", style={"marginBottom": "15px"}),

                    html.Div(id='md-content', style={
                        "padding": "15px",
                        "minHeight": "calc(100vh - 250px)"
                    })
                ], width=9)
            ], style={"alignItems": "stretch"})
        ], fluid=True),

        # Store components
        dcc.Store(id='md-region-banker-map', data=region_banker_map),
        dcc.Store(id='md-last-update', data=datetime.now().strftime('%Y-%m-%d %H:%M:%S')),

        # Download component
        dcc.Download(id="md-download-dataframe")

    ], style={"minHeight": "100vh", "padding": "20px"})
# Gray color palette for all visualizations
GRAY_PALETTE = {
    'dark': '#2d3436',
    'medium_dark': '#636e72',
    'medium': '#95a5a6',
    'medium_light': '#b2bec3',
    'light': '#dfe6e9'
}

GRAY_SHADES_LIST = ['#2d3436', '#636e72', '#7f8c8d', '#95a5a6', '#b2bec3', '#dfe6e9']

def get_gray_shades(n):
    """Get n gray shades for visualizations"""
    if n <= len(GRAY_SHADES_LIST):
        return GRAY_SHADES_LIST[:n]
    else:
        # Generate more shades if needed
        return GRAY_SHADES_LIST * (n // len(GRAY_SHADES_LIST) + 1)


def create_km_trend_view_v2(dashboard_df, key_matrix_df, region, banker, strategist):
    """Create trend analysis view with real data from both sheets"""

    # Filter data
    if strategist != 'All Strategists':
        chart_df = dashboard_df[dashboard_df['Investment_Strategist'] == strategist]
        km_filtered = key_matrix_df[key_matrix_df['Investment Strategist'] == strategist]
    elif region == 'All Regions':
        chart_df = dashboard_df[(dashboard_df['Region'] == 'All Regions') & (dashboard_df['Banker'] == 'All Bankers')]
        km_filtered = key_matrix_df[(key_matrix_df['Region'].isin(['All Region', 'All Regions'])) &
                                    (key_matrix_df['Banker'].isin(['All Banker', 'All Bankers']))]
    else:
        chart_df = dashboard_df[dashboard_df['Region'] == region]
        if banker != 'All Bankers':
            chart_df = chart_df[chart_df['Banker'] == banker]
        km_filtered = key_matrix_df[key_matrix_df['Region'] == region]
        if banker != 'All Bankers':
            km_filtered = km_filtered[km_filtered['Banker'] == banker]

    if chart_df.empty:
        return dbc.Alert("No data available for trend analysis", color="info")

    # Get summary values
    if not km_filtered.empty:
        km_row = km_filtered.iloc[0]
        families_growth = int(km_row['families_growth'])
        aum_growth = int(km_row['aum_growth'])
        revenue_growth = int(km_row['revenue_growth'])
    else:
        families_growth = aum_growth = revenue_growth = 0

    # Sort and prepare data
    chart_df = chart_df.sort_values(['Year', 'Quarter'])

    # Create quarter labels
    quarter_labels = []
    for _, row in chart_df.iterrows():
        year_short = str(row['Year'])[-2:]
        quarter_labels.append(f"FY{year_short}-{row['Quarter']}")

    # Create trend line charts
    fig = make_subplots(
        rows=2, cols=2,
        subplot_titles=('Families Trend', 'AUM Trend (Cr)', 'Billed Revenue Trend (Cr)', 'Total Revenue Trend (Cr)'),
        specs=[[{"type": "scatter"}, {"type": "scatter"}],
               [{"type": "scatter"}, {"type": "scatter"}]]
    )

    # Families trend
    fig.add_trace(
        go.Scatter(
            x=quarter_labels,
            y=chart_df['Families'].tolist(),
            mode='lines+markers+text',
            marker=dict(size=10, color='#636e72'),
            line=dict(color='#636e72', width=3),
            text=chart_df['Families'].tolist(),
            textposition='top center',
            textfont={'size': 10},
            name='Families',
            hovertemplate='<b>%{x}</b><br>Families: %{y}<extra></extra>'
        ),
        row=1, col=1
    )

    # AUM trend
    fig.add_trace(
        go.Scatter(
            x=quarter_labels,
            y=chart_df['AUM_Cr'].tolist(),
            mode='lines+markers+text',
            marker=dict(size=10, color='#7f8c8d'),
            line=dict(color='#7f8c8d', width=3),
            text=chart_df['AUM_Cr'].tolist(),
            textposition='top center',
            textfont={'size': 10},
            name='AUM',
            hovertemplate='<b>%{x}</b><br>AUM: ₹%{y} Cr<extra></extra>'
        ),
        row=1, col=2
    )

    # Billed Revenue trend
    fig.add_trace(
        go.Scatter(
            x=quarter_labels,
            y=chart_df['Billed_Revenue_Cr'].tolist(),
            mode='lines+markers+text',
            marker=dict(size=10, color='#95a5a6'),
            line=dict(color='#95a5a6', width=3),
            text=chart_df['Billed_Revenue_Cr'].round(1).tolist(),
            textposition='top center',
            textfont={'size': 10},
            name='Billed Revenue',
            hovertemplate='<b>%{x}</b><br>Revenue: ₹%{y} Cr<extra></extra>'
        ),
        row=2, col=1
    )

    # Total Revenue trend
    fig.add_trace(
        go.Scatter(
            x=quarter_labels,
            y=chart_df['Total_Revenue_Cr'].tolist(),
            mode='lines+markers+text',
            marker=dict(size=10, color='#b2bec3'),
            line=dict(color='#b2bec3', width=3),
            text=chart_df['Total_Revenue_Cr'].round(1).tolist(),
            textposition='top center',
            textfont={'size': 10},
            name='Total Revenue',
            hovertemplate='<b>%{x}</b><br>Revenue: ₹%{y} Cr<extra></extra>'
        ),
        row=2, col=2
    )

    fig.update_layout(
        height=700,
        showlegend=False,
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=11),
        hoverlabel=dict(
            bgcolor="white",
            font_size=14,
            font_family="Inter",
            font_color="#333",
            bordercolor="#C41E3A"
        )
    )

    fig.update_xaxes(showgrid=False, showline=True, linecolor='#ddd', tickangle=-45)
    fig.update_yaxes(showgrid=True, gridcolor='#f0f0f0')

    return html.Div([
        html.H4("Trend Analysis", style={"color": "#C41E3A", "marginBottom": "20px"}),

        # Growth Summary Cards
        dbc.Row([
            dbc.Col([
                html.Div([
                    html.I(className="fas fa-users", style={"fontSize": "24px", "color": "#636e72"}),
                    html.H3(f"{'+' if families_growth >= 0 else ''}{families_growth}%", style={
                        "color": get_percentage_color(families_growth),
                        "margin": "10px 0 5px 0",
                        "fontWeight": "bold"
                    }),
                    html.P("Families Growth", style={"fontSize": "12px", "color": "#666", "margin": "0"})
                ], style={
                    "padding": "20px",
                    "backgroundColor": "white",
                    "border": "1px solid #ddd",
                    "borderRadius": "10px",
                    "textAlign": "center"
                })
            ], width=4),
            dbc.Col([
                html.Div([
                    html.I(className="fas fa-chart-line", style={"fontSize": "24px", "color": "#7f8c8d"}),
                    html.H3(f"{'+' if aum_growth >= 0 else ''}{aum_growth}%", style={
                        "color": get_percentage_color(aum_growth),
                        "margin": "10px 0 5px 0",
                        "fontWeight": "bold"
                    }),
                    html.P("AUM Growth", style={"fontSize": "12px", "color": "#666", "margin": "0"})
                ], style={
                    "padding": "20px",
                    "backgroundColor": "white",
                    "border": "1px solid #ddd",
                    "borderRadius": "10px",
                    "textAlign": "center"
                })
            ], width=4),
            dbc.Col([
                html.Div([
                    html.I(className="fas fa-dollar-sign", style={"fontSize": "24px", "color": "#95a5a6"}),
                    html.H3(f"{'+' if revenue_growth >= 0 else ''}{revenue_growth}%", style={
                        "color": get_percentage_color(revenue_growth),
                        "margin": "10px 0 5px 0",
                        "fontWeight": "bold"
                    }),
                    html.P("Revenue Growth", style={"fontSize": "12px", "color": "#666", "margin": "0"})
                ], style={
                    "padding": "20px",
                    "backgroundColor": "white",
                    "border": "1px solid #ddd",
                    "borderRadius": "10px",
                    "textAlign": "center"
                })
            ], width=4)
        ], className="mb-4"),

        # Trend Chart
        dcc.Graph(figure=fig, config={'displayModeBar': False})
    ])


def create_km_table_view_v2(dashboard_df, key_matrix_df, region, banker, strategist):
    """Create data table view with real filtered data"""

    # Filter data
    if strategist != 'All Strategists':
        filtered_df = dashboard_df[dashboard_df['Investment_Strategist'] == strategist]
    elif region == 'All Regions':
        filtered_df = dashboard_df[
            (dashboard_df['Region'] == 'All Regions') & (dashboard_df['Banker'] == 'All Bankers')]
    else:
        filtered_df = dashboard_df[dashboard_df['Region'] == region]
        if banker != 'All Bankers':
            filtered_df = filtered_df[filtered_df['Banker'] == banker]

    if filtered_df.empty:
        return dbc.Alert("No data available", color="info")

    # Sort by Year and Quarter
    filtered_df = filtered_df.sort_values(['Year', 'Quarter'])

    # Select key columns for display
    display_columns = ['Quarter', 'Year', 'Region', 'Banker', 'Families', 'Families YoY',
                       'AUM_Cr', 'YoY(AUM+AUA)', 'Core_AUM', 'XSell_AUM',
                       'Billed_Revenue_Cr', 'Total_Revenue_Cr']

    display_df = filtered_df[display_columns].copy()

    # Create FY-Quarter column
    display_df['FY-Quarter'] = display_df.apply(
        lambda row: f"FY{str(row['Year'])[-2:]}-{row['Quarter']}", axis=1
    )

    # Reorder columns
    cols = ['FY-Quarter'] + [col for col in display_df.columns if col != 'FY-Quarter']
    display_df = display_df[cols]

    # Create table
    return html.Div([
        html.H4("Data Table", style={"color": "#C41E3A", "marginBottom": "20px"}),

        html.Div([
            html.Table([
                html.Thead([
                    html.Tr([
                        html.Th(col, style={
                            "padding": "12px",
                            "backgroundColor": "#f8f9fa",
                            "borderBottom": "2px solid #dee2e6",
                            "fontSize": "11px",
                            "fontWeight": "bold",
                            "position": "sticky",
                            "top": "0",
                            "zIndex": "10"
                        }) for col in display_df.columns
                    ])
                ]),
                html.Tbody([
                    html.Tr([
                        html.Td(str(display_df.iloc[i][col]), style={
                            "padding": "10px",
                            "borderBottom": "1px solid #dee2e6",
                            "fontSize": "11px"
                        }) for col in display_df.columns
                    ], style={
                        "backgroundColor": "#f8f9fa" if i % 2 == 0 else "white"
                    }) for i in range(len(display_df))
                ])
            ], style={
                "width": "100%",
                "borderCollapse": "collapse",
                "border": "1px solid #dee2e6"
            })
        ], style={"overflowX": "auto", "overflowY": "auto", "maxHeight": "600px"}),

        html.P(f"Showing {len(display_df)} rows",
               style={"marginTop": "10px", "fontSize": "11px", "color": "#666"})
    ])


# ============================================================================
# X-SELL EXPORTS - Add these new functions
# ============================================================================

def create_xs_excel_export(df, region, banker, strategist, product_filter, stage_filter, filter_desc):
    """Create comprehensive Excel export for X-Sell data"""

    # Filter data
    filtered_df = filter_data_by_hierarchy(df, region, banker, strategist)

    if product_filter and product_filter != 'All Products':
        filtered_df = filtered_df[filtered_df['Product'] == product_filter]

    if stage_filter and stage_filter != 'All Stages':
        filtered_df = filtered_df[filtered_df['Stage'] == stage_filter]

    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        workbook = writer.book

        # Define formats
        header_format = workbook.add_format({
            'bold': True,
            'bg_color': '#b04d59',
            'font_color': 'white',
            'align': 'center',
            'valign': 'vcenter',
            'border': 1
        })

        # Sheet 1: Summary
        if not filtered_df.empty:
            summary_row = filtered_df.iloc[0]

            summary_data = {
                'Metric': ['Filter Applied', 'Total (Cr)', 'YoY Growth (%)', 'HC Prop (Cr)',
                           'HC Prop YoY (%)', 'HC TP (Cr)', 'HC TP YoY (%)'],
                'Value': [
                    filter_desc,
                    round(filtered_df['Total'].sum() if 'Total' in filtered_df.columns else filtered_df[
                        'Amount_Cr'].sum(), 1),
                    int(summary_row.get('YoY_Growth', 0)),
                    round(filtered_df['HC_Prop'].sum() if 'HC_Prop' in filtered_df.columns else 0, 1),
                    int(summary_row.get('HC_Prop_YoY', 0)),
                    round(filtered_df['HC_TP'].sum() if 'HC_TP' in filtered_df.columns else 0, 1),
                    int(summary_row.get('HC_TP_YoY', 0))
                ]
            }
            summary_df = pd.DataFrame(summary_data)
            summary_df.to_excel(writer, sheet_name='Summary', index=False)

            worksheet = writer.sheets['Summary']
            worksheet.set_column('A:A', 25)
            worksheet.set_column('B:B', 20)

            for col_num, value in enumerate(summary_df.columns.values):
                worksheet.write(0, col_num, value, header_format)

        # Sheet 2: Detailed Data by Stage
        agg_data = filtered_df.groupby(['Product', 'Category', 'Stage'])['Amount_Cr'].sum().reset_index()
        stage_pivot = agg_data.pivot_table(
            index=['Product', 'Category'],
            columns='Stage',
            values='Amount_Cr',
            fill_value=0
        ).reset_index()

        stage_pivot.to_excel(writer, sheet_name='By Stage', index=False)
        worksheet = writer.sheets['By Stage']

        for col_num, value in enumerate(stage_pivot.columns.values):
            worksheet.write(0, col_num, value, header_format)
            worksheet.set_column(col_num, col_num, 18)

        # Sheet 3: Raw Data
        filtered_df.to_excel(writer, sheet_name='Raw Data', index=False)
        worksheet = writer.sheets['Raw Data']

        for col_num, value in enumerate(filtered_df.columns.values):
            worksheet.write(0, col_num, value, header_format)
            worksheet.set_column(col_num, col_num, 15)

    output.seek(0)
    return output


def create_xs_pdf_export(df, region, banker, strategist, product_filter, stage_filter, filter_desc):
    """Create comprehensive PDF for X-Sell data"""

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter,
                            rightMargin=50, leftMargin=50,
                            topMargin=50, bottomMargin=50)

    elements = []
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=28,
        textColor=colors.HexColor('#b04d59'),
        spaceAfter=30,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
    )

    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=18,
        textColor=colors.HexColor('#b04d59'),
        spaceAfter=15,
        spaceBefore=20,
        fontName='Helvetica-Bold'
    )

    # Title Page
    elements.append(Paragraph("TRUSTPLUTUS", title_style))
    elements.append(Paragraph("Cross-Sell Dashboard Report", heading_style))
    elements.append(Paragraph(f"Filter: {filter_desc}", styles['Normal']))
    elements.append(Paragraph(
        f"Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}",
        styles['Normal']
    ))
    elements.append(Spacer(1, 30))

    # Filter data
    filtered_df = filter_data_by_hierarchy(df, region, banker, strategist)

    if product_filter and product_filter != 'All Products':
        filtered_df = filtered_df[filtered_df['Product'] == product_filter]

    if stage_filter and stage_filter != 'All Stages':
        filtered_df = filtered_df[filtered_df['Stage'] == stage_filter]

    if filtered_df.empty:
        elements.append(Paragraph("No data available for selected filters", styles['Normal']))
        doc.build(elements)
        buffer.seek(0)
        return buffer

    # Summary Table
    elements.append(Paragraph("Executive Summary", heading_style))

    summary_row = filtered_df.iloc[0]
    total_value = round(
        filtered_df['Total'].sum() if 'Total' in filtered_df.columns else filtered_df['Amount_Cr'].sum(), 1)
    hc_prop_value = round(filtered_df['HC_Prop'].sum() if 'HC_Prop' in filtered_df.columns else 0, 1)
    hc_tp_value = round(filtered_df['HC_TP'].sum() if 'HC_TP' in filtered_df.columns else 0, 1)

    summary_data = [
        ['Metric', 'Value', 'YoY Growth'],
        ['Total (Cr)', str(total_value), f"{int(summary_row.get('YoY_Growth', 0))}%"],
        ['HC Prop (Cr)', str(hc_prop_value), f"{int(summary_row.get('HC_Prop_YoY', 0))}%"],
        ['HC TP (Cr)', str(hc_tp_value), f"{int(summary_row.get('HC_TP_YoY', 0))}%"]
    ]

    summary_table = Table(summary_data, colWidths=[3 * inch, 2 * inch, 1.5 * inch])
    summary_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#b04d59')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
    ]))

    elements.append(summary_table)
    elements.append(Spacer(1, 30))

    # Stage Distribution Chart
    elements.append(PageBreak())
    elements.append(Paragraph("Stage Distribution", heading_style))

    stage_totals = filtered_df.groupby('Stage')['Amount_Cr'].sum()

    fig = go.Figure()
    fig.add_trace(go.Bar(
        x=list(stage_totals.index),
        y=list(stage_totals.values),
        marker_color=['#636e72', '#95a5a6', '#b2bec3'],
        text=[f"{val:.1f}" for val in stage_totals.values],
        textposition='outside'
    ))

    fig.update_layout(
        title="Amount by Stage (Cr)",
        xaxis_title="Stage",
        yaxis_title="Amount (Cr)",
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=12),
        height=400,
        showlegend=False
    )

    img_bytes = save_plotly_as_image(fig, width=700, height=400)
    if img_bytes:
        img_stream = io.BytesIO(img_bytes)
        img = Image(img_stream, width=6 * inch, height=3.5 * inch)
        elements.append(img)

    # Product Distribution
    elements.append(PageBreak())
    elements.append(Paragraph("Product Distribution", heading_style))

    product_totals = filtered_df.groupby('Product')['Amount_Cr'].sum().sort_values(ascending=False).head(10)

    fig2 = go.Figure()
    fig2.add_trace(go.Bar(
        x=list(product_totals.index),
        y=list(product_totals.values),
        marker_color='#f7ce7c',
        text=[f"{val:.1f}" for val in product_totals.values],
        textposition='outside'
    ))

    fig2.update_layout(
        title="Top 10 Products by Amount (Cr)",
        xaxis_title="Product",
        yaxis_title="Amount (Cr)",
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=12),
        height=400,
        showlegend=False,
        xaxis_tickangle=-45
    )

    img_bytes = save_plotly_as_image(fig2, width=700, height=400)
    if img_bytes:
        img_stream = io.BytesIO(img_bytes)
        img = Image(img_stream, width=6 * inch, height=3.5 * inch)
        elements.append(img)

    # Footer
    elements.append(Spacer(1, 30))
    footer_text = Paragraph(
        "www.trustplutus.com | Private and Confidential",
        ParagraphStyle('Footer', parent=styles['Normal'],
                       fontSize=8, textColor=colors.grey,
                       alignment=TA_CENTER)
    )
    elements.append(footer_text)

    doc.build(elements)
    buffer.seek(0)
    return buffer


def create_xs_ppt_export(df, region, banker, strategist, product_filter, stage_filter, filter_desc):
    """Create comprehensive PowerPoint for X-Sell data"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title Slide
    prs = create_title_slide_ppt(prs, "Cross-Sell Dashboard", "Product Performance Analytics", filter_desc)

    # Filter data
    filtered_df = filter_data_by_hierarchy(df, region, banker, strategist)

    if product_filter and product_filter != 'All Products':
        filtered_df = filtered_df[filtered_df['Product'] == product_filter]

    if stage_filter and stage_filter != 'All Stages':
        filtered_df = filtered_df[filtered_df['Stage'] == stage_filter]

    if filtered_df.empty:
        return prs

    # Summary Slide
    MAROON = RGBColor(176, 77, 89)
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Executive Summary"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = MAROON

    summary_row = filtered_df.iloc[0]
    total_value = round(
        filtered_df['Total'].sum() if 'Total' in filtered_df.columns else filtered_df['Amount_Cr'].sum(), 1)
    hc_prop_value = round(filtered_df['HC_Prop'].sum() if 'HC_Prop' in filtered_df.columns else 0, 1)
    hc_tp_value = round(filtered_df['HC_TP'].sum() if 'HC_TP' in filtered_df.columns else 0, 1)

    # Create summary table
    rows = 4
    cols = 3
    left = Inches(2)
    top = Inches(2.5)
    width = Inches(6)
    height = Inches(3)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Header
    headers = ['Metric', 'Value (Cr)', 'YoY Growth']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = MAROON

    # Data
    data_rows = [
        ['Total', str(total_value), f"{int(summary_row.get('YoY_Growth', 0))}%"],
        ['HC Prop', str(hc_prop_value), f"{int(summary_row.get('HC_Prop_YoY', 0))}%"],
        ['HC TP', str(hc_tp_value), f"{int(summary_row.get('HC_TP_YoY', 0))}%"]
    ]

    for row_idx, row_data in enumerate(data_rows, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Stage Distribution Chart
    stage_totals = filtered_df.groupby('Stage')['Amount_Cr'].sum()

    fig_stages = go.Figure()
    fig_stages.add_trace(go.Bar(
        x=list(stage_totals.index),
        y=list(stage_totals.values),
        marker_color=['#636e72', '#95a5a6', '#b2bec3'],
        text=[f"{val:.1f}" for val in stage_totals.values],
        textposition='outside'
    ))

    fig_stages.update_layout(
        title="Amount Distribution by Stage (Cr)",
        xaxis_title="Stage",
        yaxis_title="Amount (Cr)",
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=14),
        height=500,
        showlegend=False
    )

    prs = add_chart_slide_ppt(prs, "Stage Distribution", fig_stages,
                              notes="Breakdown of cross-sell amounts across pipeline stages")

    # Product Distribution Chart
    product_totals = filtered_df.groupby('Product')['Amount_Cr'].sum().sort_values(ascending=False).head(10)

    fig_products = go.Figure()
    fig_products.add_trace(go.Bar(
        x=list(product_totals.index),
        y=list(product_totals.values),
        marker_color='#f7ce7c',
        text=[f"{val:.1f}" for val in product_totals.values],
        textposition='outside'
    ))

    fig_products.update_layout(
        title="Top 10 Products by Amount (Cr)",
        xaxis_title="Product",
        yaxis_title="Amount (Cr)",
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=14),
        height=500,
        showlegend=False,
        xaxis_tickangle=-45
    )

    prs = add_chart_slide_ppt(prs, "Top Products", fig_products,
                              notes="Highest performing products by cross-sell amount")

    # Category Distribution
    category_totals = filtered_df.groupby('Category')['Amount_Cr'].sum()

    fig_category = go.Figure()
    fig_category.add_trace(go.Pie(
        labels=list(category_totals.index),
        values=list(category_totals.values),
        marker_colors=['#b04d59', '#f7ce7c', '#636e72', '#95a5a6'],
        textinfo='label+percent',
        textposition='outside'
    ))

    fig_category.update_layout(
        title="Distribution by Category",
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=14),
        height=500,
        showlegend=True
    )

    prs = add_chart_slide_ppt(prs, "Category Distribution", fig_category,
                              notes="Product category breakdown")

    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


# ============================================================================
# MANDATES EXPORTS - Add these new functions
# ============================================================================

def create_md_excel_export(df, region, banker, strategist, filter_desc):
    """Create comprehensive Excel export for Mandates data"""

    filtered_df = filter_data_by_hierarchy(df, region, banker, strategist)

    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        workbook = writer.book

        header_format = workbook.add_format({
            'bold': True,
            'bg_color': '#b04d59',
            'font_color': 'white',
            'align': 'center',
            'valign': 'vcenter',
            'border': 1
        })

        # Sheet 1: Summary
        stage_counts = filtered_df.groupby('Stage')['Count'].sum()

        summary_data = {
            'Metric': ['Filter Applied', 'Total Pitched', 'Total Offered',
                       'Total Converted', 'Total Declined', 'Conversion Rate (%)'],
            'Value': [
                filter_desc,
                int(stage_counts.get('Pitched', 0)),
                int(stage_counts.get('Offered', 0)),
                int(stage_counts.get('Converted', 0)),
                int(stage_counts.get('Declined', 0)),
                round((stage_counts.get('Converted', 0) / stage_counts.get('Pitched', 1) * 100), 1) if stage_counts.get(
                    'Pitched', 0) > 0 else 0
            ]
        }

        summary_df = pd.DataFrame(summary_data)
        summary_df.to_excel(writer, sheet_name='Summary', index=False)

        worksheet = writer.sheets['Summary']
        worksheet.set_column('A:A', 25)
        worksheet.set_column('B:B', 20)

        for col_num, value in enumerate(summary_df.columns.values):
            worksheet.write(0, col_num, value, header_format)

        # Sheet 2: Quarterly Breakdown
        quarterly_data = filtered_df.groupby(['Quarter', 'Stage'])['Count'].sum().reset_index()
        quarterly_pivot = quarterly_data.pivot(index='Quarter', columns='Stage', values='Count').fillna(0).reset_index()

        quarterly_pivot.to_excel(writer, sheet_name='Quarterly', index=False)
        worksheet = writer.sheets['Quarterly']

        for col_num, value in enumerate(quarterly_pivot.columns.values):
            worksheet.write(0, col_num, value, header_format)
            worksheet.set_column(col_num, col_num, 15)

        # Sheet 3: Raw Data
        filtered_df.to_excel(writer, sheet_name='Raw Data', index=False)
        worksheet = writer.sheets['Raw Data']

        for col_num, value in enumerate(filtered_df.columns.values):
            worksheet.write(0, col_num, value, header_format)
            worksheet.set_column(col_num, col_num, 15)

    output.seek(0)
    return output


def create_md_pdf_export(df, region, banker, strategist, filter_desc):
    """Create comprehensive PDF for Mandates data"""

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter,
                            rightMargin=50, leftMargin=50,
                            topMargin=50, bottomMargin=50)

    elements = []
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=28,
        textColor=colors.HexColor('#b04d59'),
        spaceAfter=30,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
    )

    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=18,
        textColor=colors.HexColor('#b04d59'),
        spaceAfter=15,
        spaceBefore=20,
        fontName='Helvetica-Bold'
    )

    # Title Page
    elements.append(Paragraph("TRUSTPLUTUS", title_style))
    elements.append(Paragraph("New Mandates Dashboard Report", heading_style))
    elements.append(Paragraph(f"Filter: {filter_desc}", styles['Normal']))
    elements.append(Paragraph(
        f"Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}",
        styles['Normal']
    ))
    elements.append(Spacer(1, 30))

    # Filter data
    filtered_df = filter_data_by_hierarchy(df, region, banker, strategist)

    if filtered_df.empty:
        elements.append(Paragraph("No data available for selected filters", styles['Normal']))
        doc.build(elements)
        buffer.seek(0)
        return buffer

    # Summary Table
    elements.append(Paragraph("Executive Summary", heading_style))

    stage_counts = filtered_df.groupby('Stage')['Count'].sum()
    total_pitched = int(stage_counts.get('Pitched', 0))
    total_converted = int(stage_counts.get('Converted', 0))
    conversion_rate = round((total_converted / total_pitched * 100), 1) if total_pitched > 0 else 0

    summary_data = [
        ['Metric', 'Count'],
        ['Pitched', str(total_pitched)],
        ['Offered', str(int(stage_counts.get('Offered', 0)))],
        ['Converted', str(total_converted)],
        ['Declined', str(int(stage_counts.get('Declined', 0)))],
        ['Conversion Rate', f"{conversion_rate}%"]
    ]

    summary_table = Table(summary_data, colWidths=[3 * inch, 2 * inch])
    summary_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#b04d59')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
    ]))

    elements.append(summary_table)
    elements.append(Spacer(1, 30))

    # Pipeline Chart
    elements.append(PageBreak())
    elements.append(Paragraph("Pipeline Distribution", heading_style))

    fig = go.Figure()
    stages = ['Pitched', 'Offered', 'Converted', 'Declined']
    stage_values = [stage_counts.get(s, 0) for s in stages]
    colors_list = ['#b04d59', '#f7ce7c', '#28a745', '#6C757D']

    fig.add_trace(go.Bar(
        x=stages,
        y=stage_values,
        marker_color=colors_list,
        text=stage_values,
        textposition='outside'
    ))

    fig.update_layout(
        title="Mandate Pipeline by Stage",
        xaxis_title="Stage",
        yaxis_title="Count",
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=12),
        height=400,
        showlegend=False
    )

    img_bytes = save_plotly_as_image(fig, width=700, height=400)
    if img_bytes:
        img_stream = io.BytesIO(img_bytes)
        img = Image(img_stream, width=6 * inch, height=3.5 * inch)
        elements.append(img)

    # Quarterly Trend
    elements.append(PageBreak())
    elements.append(Paragraph("Quarterly Trend", heading_style))

    quarterly_data = filtered_df.groupby(['Quarter', 'Stage'])['Count'].sum().reset_index()

    fig2 = go.Figure()

    for stage, color in zip(['Pitched', 'Offered', 'Converted', 'Declined'], colors_list):
        stage_data = quarterly_data[quarterly_data['Stage'] == stage]
        fig2.add_trace(go.Scatter(
            x=stage_data['Quarter'],
            y=stage_data['Count'],
            mode='lines+markers',
            name=stage,
            line=dict(color=color, width=3),
            marker=dict(size=8, color=color)
        ))

    fig2.update_layout(title="Quarterly Mandate Trends",
        xaxis_title="Quarter",
        yaxis_title="Count",
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=12),
        height=400,
        showlegend=True,
        legend=dict(orientation="h", yanchor="bottom", y=-0.3, xanchor="center", x=0.5)
    )

    img_bytes = save_plotly_as_image(fig2, width=700, height=400)
    if img_bytes:
        img_stream = io.BytesIO(img_bytes)
        img = Image(img_stream, width=6*inch, height=3.5*inch)
        elements.append(img)

    # Footer
    elements.append(Spacer(1, 30))
    footer_text = Paragraph(
        "www.trustplutus.com | Private and Confidential",
        ParagraphStyle('Footer', parent=styles['Normal'],
                      fontSize=8, textColor=colors.grey,
                      alignment=TA_CENTER)
    )
    elements.append(footer_text)

    doc.build(elements)
    buffer.seek(0)
    return buffer


def create_md_ppt_export(df, region, banker, strategist, filter_desc):
    """Create comprehensive PowerPoint for Mandates data"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title Slide
    prs = create_title_slide_ppt(prs, "New Mandates Dashboard", "Pipeline Analytics", filter_desc)

    # Filter data
    filtered_df = filter_data_by_hierarchy(df, region, banker, strategist)

    if filtered_df.empty:
        return prs

    # Summary Slide
    MAROON = RGBColor(176, 77, 89)
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Executive Summary"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = MAROON

    stage_counts = filtered_df.groupby('Stage')['Count'].sum()
    total_pitched = int(stage_counts.get('Pitched', 0))
    total_converted = int(stage_counts.get('Converted', 0))
    conversion_rate = round((total_converted / total_pitched * 100), 1) if total_pitched > 0 else 0

    # Create summary table
    rows = 6
    cols = 2
    left = Inches(2.5)
    top = Inches(2.5)
    width = Inches(5)
    height = Inches(3.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Header
    headers = ['Metric', 'Count']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = MAROON

    # Data
    data_rows = [
        ['Pitched', str(total_pitched)],
        ['Offered', str(int(stage_counts.get('Offered', 0)))],
        ['Converted', str(total_converted)],
        ['Declined', str(int(stage_counts.get('Declined', 0)))],
        ['Conversion Rate', f"{conversion_rate}%"]
    ]

    for row_idx, row_data in enumerate(data_rows, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Pipeline Distribution Chart
    stages = ['Pitched', 'Offered', 'Converted', 'Declined']
    stage_values = [stage_counts.get(s, 0) for s in stages]
    colors_list = ['#b04d59', '#f7ce7c', '#28a745', '#6C757D']

    fig_pipeline = go.Figure()
    fig_pipeline.add_trace(go.Bar(
        x=stages,
        y=stage_values,
        marker_color=colors_list,
        text=stage_values,
        textposition='outside'
    ))

    fig_pipeline.update_layout(
        title="Mandate Pipeline Distribution",
        xaxis_title="Stage",
        yaxis_title="Count",
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=14),
        height=500,
        showlegend=False
    )

    prs = add_chart_slide_ppt(prs, "Pipeline Distribution", fig_pipeline,
                              notes=f"Overall conversion rate: {conversion_rate}%")

    # Quarterly Trend Chart
    quarterly_data = filtered_df.groupby(['Quarter', 'Stage'])['Count'].sum().reset_index()

    fig_trend = go.Figure()

    for stage, color in zip(['Pitched', 'Offered', 'Converted', 'Declined'], colors_list):
        stage_data = quarterly_data[quarterly_data['Stage'] == stage]
        fig_trend.add_trace(go.Scatter(
            x=stage_data['Quarter'],
            y=stage_data['Count'],
            mode='lines+markers',
            name=stage,
            line=dict(color=color, width=3),
            marker=dict(size=10, color=color)
        ))

    fig_trend.update_layout(
        title="Quarterly Mandate Trends",
        xaxis_title="Quarter",
        yaxis_title="Count",
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=14),
        height=500,
        showlegend=True,
        legend=dict(orientation="h", yanchor="bottom", y=-0.2, xanchor="center", x=0.5)
    )

    prs = add_chart_slide_ppt(prs, "Quarterly Trends", fig_trend,
                              notes="Pipeline progression across quarters")

    # Conversion Funnel
    funnel_data = {
        'Stage': ['Pitched', 'Offered', 'Converted'],
        'Count': [
            stage_counts.get('Pitched', 0),
            stage_counts.get('Offered', 0),
            stage_counts.get('Converted', 0)
        ]
    }

    fig_funnel = go.Figure()
    fig_funnel.add_trace(go.Funnel(
        y=funnel_data['Stage'],
        x=funnel_data['Count'],
        textposition="inside",
        textinfo="value+percent initial",
        marker={"color": ["#b04d59", "#f7ce7c", "#28a745"]},
        connector={"line": {"color": "#ddd", "dash": "dot", "width": 2}}
    ))

    fig_funnel.update_layout(
        title="Conversion Funnel",
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=14),
        height=500
    )

    prs = add_chart_slide_ppt(prs, "Conversion Funnel", fig_funnel,
                              notes="Visual representation of mandate conversion")

    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

@app.callback(
    Output("session-store", "data", allow_duplicate=True),
    Input("nav-logout-btn", "n_clicks"),
    prevent_initial_call=True
)
def handle_navbar_logout(n_clicks):
    """Handle logout from navbar"""
    global current_user
    if n_clicks:
        current_user = {}
        return {"page": "home", "logged_in": False}
    return dash.no_update
# Mandates Excel Export
@app.callback(
    Output("md-download-dataframe", "data", allow_duplicate=True),
    Input("md-export-excel-btn", "n_clicks"),
    State('md-region-filter', 'value'),
    State('md-banker-filter', 'value'),
    State('md-strategist-filter', 'value'),
    prevent_initial_call=True
)
def export_md_excel(n_clicks, region, banker, strategist):
    if n_clicks:
        try:
            df = pd.read_excel('data/mandate_data.xlsx')

            if strategist != 'All Strategists':
                filter_desc = f"Investment Strategist: {strategist}"
            elif region == 'All Regions':
                filter_desc = "All Regions | All Bankers"
            else:
                filter_desc = f"{region} | {banker}"

            excel_buffer = create_md_excel_export(df, region, banker, strategist, filter_desc)

            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"Mandates_{region}_{banker}_{timestamp}.xlsx"

            return dcc.send_bytes(excel_buffer.getvalue(), filename)
        except Exception as e:
            print(f"Excel Export Error: {e}")
            return None
    return None


# Mandates PDF Export
@app.callback(
    Output("md-download-dataframe", "data", allow_duplicate=True),
    Input("md-export-pdf-btn", "n_clicks"),
    State('md-region-filter', 'value'),
    State('md-banker-filter', 'value'),
    State('md-strategist-filter', 'value'),
    prevent_initial_call=True
)
def export_md_pdf(n_clicks, region, banker, strategist):
    if n_clicks:
        try:
            df = pd.read_excel('data/mandate_data.xlsx')

            if strategist != 'All Strategists':
                filter_desc = f"Investment Strategist: {strategist}"
            elif region == 'All Regions':
                filter_desc = "All Regions | All Bankers"
            else:
                filter_desc = f"{region} | {banker}"

            pdf_buffer = create_md_pdf_export(df, region, banker, strategist, filter_desc)

            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"Mandates_{region}_{banker}_{timestamp}.pdf"

            return dcc.send_bytes(pdf_buffer.getvalue(), filename)
        except Exception as e:
            print(f"PDF Export Error: {e}")
            return None
    return None


# Mandates PPT Export
@app.callback(
    Output("md-download-dataframe", "data", allow_duplicate=True),
    Input("md-export-ppt-btn", "n_clicks"),
    State('md-region-filter', 'value'),
    State('md-banker-filter', 'value'),
    State('md-strategist-filter', 'value'),
    prevent_initial_call=True
)
def export_md_ppt(n_clicks, region, banker, strategist):
    if n_clicks:
        try:
            df = pd.read_excel('data/mandate_data.xlsx')

            if strategist != 'All Strategists':
                filter_desc = f"Investment Strategist: {strategist}"
            elif region == 'All Regions':
                filter_desc = "All Regions | All Bankers"
            else:
                filter_desc = f"{region} | {banker}"

            ppt_buffer = create_md_ppt_export(df, region, banker, strategist, filter_desc)

            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"Mandates_{region}_{banker}_{timestamp}.pptx"

            return dcc.send_bytes(ppt_buffer.getvalue(), filename)
        except Exception as e:
            print(f"PPT Export Error: {e}")
            return None
    return None

# Key Matrix Collapse Callbacks
@app.callback(
    [Output("km-filter-collapse", "is_open"),
     Output("km-filter-collapse-icon", "className")],
    [Input("km-filter-collapse-btn", "n_clicks")],
    [State("km-filter-collapse", "is_open")]
)
def toggle_km_filters(n_clicks, is_open):
    if n_clicks:
        return not is_open, "fas fa-chevron-down" if is_open else "fas fa-chevron-up"
    return is_open, "fas fa-chevron-up"

@app.callback(
    [Output("km-compare-collapse", "is_open"),
     Output("km-compare-collapse-icon", "className")],
    [Input("km-compare-collapse-btn", "n_clicks")],
    [State("km-compare-collapse", "is_open")]
)
def toggle_km_compare(n_clicks, is_open):
    if n_clicks:
        return not is_open, "fas fa-chevron-up" if not is_open else "fas fa-chevron-down"
    return is_open, "fas fa-chevron-down"

# X-Sell Collapse Callbacks
@app.callback(
    [Output("xs-filter-collapse", "is_open"),
     Output("xs-filter-collapse-icon", "className")],
    [Input("xs-filter-collapse-btn", "n_clicks")],
    [State("xs-filter-collapse", "is_open")]
)
def toggle_xs_filters(n_clicks, is_open):
    if n_clicks:
        return not is_open, "fas fa-chevron-down" if is_open else "fas fa-chevron-up"
    return is_open, "fas fa-chevron-up"

@app.callback(
    [Output("xs-compare-collapse", "is_open"),
     Output("xs-compare-collapse-icon", "className")],
    [Input("xs-compare-collapse-btn", "n_clicks")],
    [State("xs-compare-collapse", "is_open")]
)
def toggle_xs_compare(n_clicks, is_open):
    if n_clicks:
        return not is_open, "fas fa-chevron-up" if not is_open else "fas fa-chevron-down"
    return is_open, "fas fa-chevron-down"

# Mandates Collapse Callbacks
@app.callback(
    [Output("md-filter-collapse", "is_open"),
     Output("md-filter-collapse-icon", "className")],
    [Input("md-filter-collapse-btn", "n_clicks")],
    [State("md-filter-collapse", "is_open")]
)
def toggle_md_filters(n_clicks, is_open):
    if n_clicks:
        return not is_open, "fas fa-chevron-down" if is_open else "fas fa-chevron-up"
    return is_open, "fas fa-chevron-up"

@app.callback(
    [Output("md-compare-collapse", "is_open"),
     Output("md-compare-collapse-icon", "className")],
    [Input("md-compare-collapse-btn", "n_clicks")],
    [State("md-compare-collapse", "is_open")]
)
def toggle_md_compare(n_clicks, is_open):
    if n_clicks:
        return not is_open, "fas fa-chevron-up" if not is_open else "fas fa-chevron-down"
    return is_open, "fas fa-chevron-down"


# ============================================================================
# EXPORT HELPER FUNCTIONS - Add after imports
# ============================================================================

def save_plotly_as_image(fig, width=800, height=600):
    """Convert Plotly figure to image bytes for embedding in PDF/PPT"""
    try:
        img_bytes = fig.to_image(format="png", width=width, height=height, engine="kaleido")
        return img_bytes
    except Exception as e:
        print(f"Error converting plot to image: {e}")
        return None


def create_title_slide_ppt(prs, title, subtitle, filter_desc):
    """Create title slide for PowerPoint with TRUSTPLUTUS branding"""
    MAROON = RGBColor(176, 77, 89)
    GOLD = RGBColor(247, 206, 124)

    title_slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(title_slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "TRUSTPLUTUS"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = MAROON
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = GOLD
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Filter info
    filter_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.5))
    filter_frame = filter_box.text_frame
    filter_frame.text = f"Filter: {filter_desc}"
    filter_para = filter_frame.paragraphs[0]
    filter_para.font.size = Pt(16)
    filter_para.font.color.rgb = RGBColor(102, 102, 102)
    filter_para.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = datetime.now().strftime('%B %d, %Y')
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(14)
    date_para.font.color.rgb = RGBColor(102, 102, 102)
    date_para.alignment = PP_ALIGN.CENTER

    return prs


def add_chart_slide_ppt(prs, title, fig, notes=""):
    """Add a slide with a Plotly chart to PowerPoint"""
    MAROON = RGBColor(176, 77, 89)

    slide_layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.color.rgb = MAROON

    # Convert Plotly to image
    img_bytes = save_plotly_as_image(fig, width=900, height=500)

    if img_bytes:
        img_stream = io.BytesIO(img_bytes)
        pic = slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.5), width=Inches(9))

    # Add notes if provided
    if notes:
        notes_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
        notes_frame = notes_box.text_frame
        notes_frame.text = notes
        notes_para = notes_frame.paragraphs[0]
        notes_para.font.size = Pt(12)
        notes_para.font.color.rgb = RGBColor(102, 102, 102)

    return prs


# ============================================================================
# KEY MATRIX EXPORT FUNCTIONS
# ============================================================================

def create_km_excel_export(dashboard_df, key_matrix_df, region, banker, strategist, filter_desc):
    """Create comprehensive Excel export with multiple sheets and formatting"""

    # Filter data
    if strategist != 'All Strategists':
        filtered_dashboard = dashboard_df[dashboard_df['Investment_Strategist'] == strategist]
        filtered_km = key_matrix_df[key_matrix_df['Investment Strategist'] == strategist]
    elif region == 'All Regions':
        filtered_dashboard = dashboard_df[
            (dashboard_df['Region'] == 'All Regions') & (dashboard_df['Banker'] == 'All Bankers')]
        filtered_km = key_matrix_df[(key_matrix_df['Region'].isin(['All Region', 'All Regions']))]
    else:
        filtered_dashboard = dashboard_df[dashboard_df['Region'] == region]
        filtered_km = key_matrix_df[key_matrix_df['Region'] == region]
        if banker != 'All Bankers':
            filtered_dashboard = filtered_dashboard[filtered_dashboard['Banker'] == banker]
            filtered_km = filtered_km[filtered_km['Banker'] == banker]

    # Create Excel with multiple sheets
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        workbook = writer.book

        # Define formats
        header_format = workbook.add_format({
            'bold': True,
            'bg_color': '#b04d59',
            'font_color': 'white',
            'align': 'center',
            'valign': 'vcenter',
            'border': 1
        })

        # Sheet 1: Summary
        if not filtered_km.empty:
            km_row = filtered_km.iloc[0]
            summary_data = {
                'Metric': ['Filter Applied', 'Total Families', 'Total AUM (Cr)', 'Core AUM (Cr)',
                           'X-Sell AUM (Cr)', 'Billed Revenue (Cr)', 'Total Revenue (Cr)'],
                'Value': [
                    filter_desc,
                    int(km_row['total_families']),
                    int(km_row['total_aum']),
                    int(km_row['core_aum']),
                    int(km_row['xsell_aum']),
                    round(float(km_row['billed_revenue']), 1),
                    round(float(km_row['total_revenue']), 1)
                ],
                'YoY Growth (%)': [
                    '',
                    int(km_row['families_growth']),
                    int(km_row['aum_growth']),
                    int(km_row['core_growth']),
                    int(km_row['xsell_growth']),
                    int(km_row['billed_growth']),
                    int(km_row['revenue_growth'])
                ]
            }
            summary_df = pd.DataFrame(summary_data)
            summary_df.to_excel(writer, sheet_name='Summary', index=False)

            worksheet = writer.sheets['Summary']
            worksheet.set_column('A:A', 25)
            worksheet.set_column('B:B', 20)
            worksheet.set_column('C:C', 18)

            for col_num, value in enumerate(summary_df.columns.values):
                worksheet.write(0, col_num, value, header_format)

        # Sheet 2: Quarterly Data
        filtered_dashboard.to_excel(writer, sheet_name='Quarterly Data', index=False)
        worksheet = writer.sheets['Quarterly Data']

        for col_num, value in enumerate(filtered_dashboard.columns.values):
            worksheet.write(0, col_num, value, header_format)
            worksheet.set_column(col_num, col_num, 15)

        # Sheet 3: Key Matrix Data
        filtered_km.to_excel(writer, sheet_name='Key Matrix', index=False)
        worksheet = writer.sheets['Key Matrix']

        for col_num, value in enumerate(filtered_km.columns.values):
            worksheet.write(0, col_num, value, header_format)
            worksheet.set_column(col_num, col_num, 18)

    output.seek(0)
    return output


def create_km_pdf_export(dashboard_df, key_matrix_df, region, banker, strategist, filter_desc):
    """Create comprehensive PDF with charts as images"""

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter,
                            rightMargin=50, leftMargin=50,
                            topMargin=50, bottomMargin=50)

    elements = []
    styles = getSampleStyleSheet()

    # Custom styles matching website colors
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=28,
        textColor=colors.HexColor('#b04d59'),
        spaceAfter=30,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
    )

    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=18,
        textColor=colors.HexColor('#b04d59'),
        spaceAfter=15,
        spaceBefore=20,
        fontName='Helvetica-Bold'
    )

    # Title Page
    elements.append(Paragraph("TRUSTPLUTUS", title_style))
    elements.append(Paragraph("Key Matrix Dashboard Report", heading_style))
    elements.append(Paragraph(f"Filter: {filter_desc}", styles['Normal']))
    elements.append(Paragraph(
        f"Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}",
        styles['Normal']
    ))
    elements.append(Spacer(1, 30))

    # Filter data
    if strategist != 'All Strategists':
        filtered_dashboard = dashboard_df[dashboard_df['Investment_Strategist'] == strategist]
        km_filtered = key_matrix_df[key_matrix_df['Investment Strategist'] == strategist]
    elif region == 'All Regions':
        filtered_dashboard = dashboard_df[
            (dashboard_df['Region'] == 'All Regions') & (dashboard_df['Banker'] == 'All Bankers')]
        km_filtered = key_matrix_df[(key_matrix_df['Region'].isin(['All Region', 'All Regions']))]
    else:
        filtered_dashboard = dashboard_df[dashboard_df['Region'] == region]
        km_filtered = key_matrix_df[key_matrix_df['Region'] == region]
        if banker != 'All Bankers':
            filtered_dashboard = filtered_dashboard[filtered_dashboard['Banker'] == banker]
            km_filtered = km_filtered[km_filtered['Banker'] == banker]

    # Executive Summary Table
    if not km_filtered.empty:
        km_row = km_filtered.iloc[0]

        elements.append(Paragraph("Executive Summary", heading_style))

        summary_data = [
            ['Metric', 'Value', 'YoY Growth'],
            ['Families', str(int(km_row['total_families'])), f"{int(km_row['families_growth'])}%"],
            ['Total AUM (Cr)', str(int(km_row['total_aum'])), f"{int(km_row['aum_growth'])}%"],
            ['Core AUM (Cr)', str(int(km_row['core_aum'])), f"{int(km_row['core_growth'])}%"],
            ['X-Sell AUM (Cr)', str(int(km_row['xsell_aum'])), f"{int(km_row['xsell_growth'])}%"],
            ['Billed Revenue (Cr)', f"{float(km_row['billed_revenue']):.1f}", f"{int(km_row['billed_growth'])}%"],
            ['Total Revenue (Cr)', f"{float(km_row['total_revenue']):.1f}", f"{int(km_row['revenue_growth'])}%"]
        ]

        summary_table = Table(summary_data, colWidths=[3 * inch, 2 * inch, 1.5 * inch])
        summary_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#b04d59')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))

        elements.append(summary_table)
        elements.append(Spacer(1, 20))

    # Create visualizations
    sorted_df = filtered_dashboard.sort_values(['Year', 'Quarter'])
    quarter_labels = []
    families_values = []
    aum_values = []

    for _, row in sorted_df.iterrows():
        year_short = str(row['Year'])[-2:]
        quarter_labels.append(f"FY{year_short}-{row['Quarter']}")
        families_values.append(int(row['Families']))
        aum_values.append(int(row['AUM_Cr']))

    # Chart 1: Families Trend
    elements.append(PageBreak())
    elements.append(Paragraph("Families Trend Analysis", heading_style))

    fig_families = go.Figure()
    fig_families.add_trace(go.Bar(
        x=quarter_labels,
        y=families_values,
        marker_color='#b04d59',
        text=families_values,
        textposition='outside'
    ))

    fig_families.update_layout(
        title="Number of Families - Quarterly Trend",
        xaxis_title="Quarter",
        yaxis_title="Families",
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=12),
        height=400,
        showlegend=False
    )

    # Add chart as image
    img_bytes = save_plotly_as_image(fig_families, width=700, height=400)
    if img_bytes:
        img_stream = io.BytesIO(img_bytes)
        img = Image(img_stream, width=6 * inch, height=3.5 * inch)
        elements.append(img)

    elements.append(Spacer(1, 20))

    # Chart 2: AUM Trend
    elements.append(PageBreak())
    elements.append(Paragraph("AUM Trend Analysis", heading_style))

    fig_aum = go.Figure()
    fig_aum.add_trace(go.Bar(
        x=quarter_labels,
        y=aum_values,
        marker_color='#f7ce7c',
        text=aum_values,
        textposition='outside'
    ))

    fig_aum.update_layout(
        title="AUM (Cr) - Quarterly Trend",
        xaxis_title="Quarter",
        yaxis_title="AUM (Cr)",
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=12),
        height=400,
        showlegend=False
    )

    img_bytes = save_plotly_as_image(fig_aum, width=700, height=400)
    if img_bytes:
        img_stream = io.BytesIO(img_bytes)
        img = Image(img_stream, width=6 * inch, height=3.5 * inch)
        elements.append(img)

    # Footer
    elements.append(Spacer(1, 30))
    footer_text = Paragraph(
        "www.trustplutus.com | Private and Confidential",
        ParagraphStyle('Footer', parent=styles['Normal'],
                       fontSize=8, textColor=colors.grey,
                       alignment=TA_CENTER)
    )
    elements.append(footer_text)

    doc.build(elements)
    buffer.seek(0)
    return buffer


def create_km_ppt_export(dashboard_df, key_matrix_df, region, banker, strategist, filter_desc):
    """Create comprehensive PowerPoint with all visualizations"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title Slide
    prs = create_title_slide_ppt(prs, "Key Matrix Dashboard", "Performance Analytics", filter_desc)

    # Filter data
    if strategist != 'All Strategists':
        filtered_dashboard = dashboard_df[dashboard_df['Investment_Strategist'] == strategist]
        km_filtered = key_matrix_df[key_matrix_df['Investment Strategist'] == strategist]
    elif region == 'All Regions':
        filtered_dashboard = dashboard_df[
            (dashboard_df['Region'] == 'All Regions') & (dashboard_df['Banker'] == 'All Bankers')]
        km_filtered = key_matrix_df[(key_matrix_df['Region'].isin(['All Region', 'All Regions']))]
    else:
        filtered_dashboard = dashboard_df[dashboard_df['Region'] == region]
        km_filtered = key_matrix_df[key_matrix_df['Region'] == region]
        if banker != 'All Bankers':
            filtered_dashboard = filtered_dashboard[filtered_dashboard['Banker'] == banker]
            km_filtered = km_filtered[km_filtered['Banker'] == banker]

    # Summary Slide with Table
    MAROON = RGBColor(176, 77, 89)
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Executive Summary"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = MAROON

    if not km_filtered.empty:
        km_row = km_filtered.iloc[0]

        # Create summary table
        rows = 7
        cols = 3
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(7)
        height = Inches(4)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set column widths
        table.columns[0].width = Inches(3)
        table.columns[1].width = Inches(2)
        table.columns[2].width = Inches(2)

        # Header
        headers = ['Metric', 'Value', 'YoY Growth']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = MAROON

        # Data
        data_rows = [
            ['Families', str(int(km_row['total_families'])), f"{int(km_row['families_growth'])}%"],
            ['Total AUM (Cr)', str(int(km_row['total_aum'])), f"{int(km_row['aum_growth'])}%"],
            ['Core AUM (Cr)', str(int(km_row['core_aum'])), f"{int(km_row['core_growth'])}%"],
            ['X-Sell AUM (Cr)', str(int(km_row['xsell_aum'])), f"{int(km_row['xsell_growth'])}%"],
            ['Billed Revenue (Cr)', f"{float(km_row['billed_revenue']):.1f}", f"{int(km_row['billed_growth'])}%"],
            ['Total Revenue (Cr)', f"{float(km_row['total_revenue']):.1f}", f"{int(km_row['revenue_growth'])}%"]
        ]

        for row_idx, row_data in enumerate(data_rows, start=1):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = cell_text
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Create visualizations
    sorted_df = filtered_dashboard.sort_values(['Year', 'Quarter'])
    quarter_labels = []
    families_values = []
    aum_values = []
    revenue_values = []

    for _, row in sorted_df.iterrows():
        year_short = str(row['Year'])[-2:]
        quarter_labels.append(f"FY{year_short}-{row['Quarter']}")
        families_values.append(int(row['Families']))
        aum_values.append(int(row['AUM_Cr']))
        revenue_values.append(round(float(row['Total_Revenue_Cr']), 1))

    # Families Chart Slide
    fig_families = go.Figure()
    fig_families.add_trace(go.Bar(
        x=quarter_labels,
        y=families_values,
        marker_color='#b04d59',
        text=families_values,
        textposition='outside'
    ))

    fig_families.update_layout(
        title="Families Trend - Quarterly Performance",
        xaxis_title="Quarter",
        yaxis_title="Number of Families",
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=14),
        height=500,
        showlegend=False
    )

    if not km_filtered.empty:
        prs = add_chart_slide_ppt(prs, "Families Trend Analysis", fig_families,
                                  notes=f"Total Growth: {int(km_filtered.iloc[0]['families_growth'])}% YoY")

    # AUM Chart Slide
    fig_aum = go.Figure()
    fig_aum.add_trace(go.Bar(
        x=quarter_labels,
        y=aum_values,
        marker_color='#f7ce7c',
        text=aum_values,
        textposition='outside'
    ))

    fig_aum.update_layout(
        title="AUM Trend - Quarterly Performance",
        xaxis_title="Quarter",
        yaxis_title="AUM (Cr)",
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=14),
        height=500,
        showlegend=False
    )

    if not km_filtered.empty:
        prs = add_chart_slide_ppt(prs, "AUM Trend Analysis", fig_aum,
                                  notes=f"Total Growth: {int(km_filtered.iloc[0]['aum_growth'])}% YoY")

    # Revenue Chart Slide
    fig_revenue = go.Figure()
    fig_revenue.add_trace(go.Bar(
        x=quarter_labels,
        y=revenue_values,
        marker_color='#636e72',
        text=revenue_values,
        textposition='outside'
    ))

    fig_revenue.update_layout(
        title="Total Revenue Trend - Quarterly Performance",
        xaxis_title="Quarter",
        yaxis_title="Revenue (Cr)",
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=14),
        height=500,
        showlegend=False
    )

    if not km_filtered.empty:
        prs = add_chart_slide_ppt(prs, "Revenue Trend Analysis", fig_revenue,
                                  notes=f"Total Growth: {int(km_filtered.iloc[0]['revenue_growth'])}% YoY")

    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer




# ============================================================================
# COMBINED KEY MATRIX EXPORT CALLBACK
# ============================================================================
@app.callback(
    Output("km-download-dataframe", "data"),
    Input("km-export-excel-btn", "n_clicks"),
    Input("km-export-pdf-btn", "n_clicks"),
    Input("km-export-ppt-btn", "n_clicks"),
    State('km-region-filter', 'value'),
    State('km-banker-filter', 'value'),
    State('km-strategist-filter', 'value'),
    prevent_initial_call=True
)
def export_km_combined(excel_clicks, pdf_clicks, ppt_clicks, region, banker, strategist):
    """Combined export callback for Key Matrix - handles Excel, PDF, and PPT"""

    ctx = dash.callback_context
    if not ctx.triggered:
        return None

    button_id = ctx.triggered[0]['prop_id'].split('.')[0]

    try:
        dashboard_df = pd.read_excel('data/key_matrix_data.xlsx', sheet_name='Dashboard')
        key_matrix_df = pd.read_excel('data/key_matrix_data.xlsx', sheet_name='Key Matrix Data')

        # Generate filter description
        if strategist != 'All Strategists':
            filter_desc = f"Investment Strategist: {strategist}"
        elif region == 'All Regions':
            filter_desc = "All Regions | All Bankers"
        else:
            filter_desc = f"{region} | {banker}"

        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

        # Excel Export
        if button_id == "km-export-excel-btn":
            print(f"Starting Excel export for {filter_desc}")
            excel_buffer = create_km_excel_export(dashboard_df, key_matrix_df, region, banker, strategist, filter_desc)
            filename = f"KeyMatrix_{region}_{banker}_{timestamp}.xlsx"
            print(f"Excel export successful: {filename}")
            return dcc.send_bytes(excel_buffer.getvalue(), filename)

        # PDF Export
        elif button_id == "km-export-pdf-btn":
            print(f"Starting PDF export for {filter_desc}")

            # Check kaleido first
            try:
                import kaleido
                print("Kaleido is installed")
            except ImportError:
                print("ERROR: Kaleido not installed!")
                return None

            pdf_buffer = create_km_pdf_export(dashboard_df, key_matrix_df, region, banker, strategist, filter_desc)
            filename = f"KeyMatrix_{region}_{banker}_{timestamp}.pdf"
            print(f"PDF export successful: {filename}")
            return dcc.send_bytes(pdf_buffer.getvalue(), filename)

        # PPT Export
        elif button_id == "km-export-ppt-btn":
            print(f"Starting PPT export for {filter_desc}")

            # Check kaleido first
            try:
                import kaleido
                print("Kaleido is installed")
            except ImportError:
                print("ERROR: Kaleido not installed!")
                return None

            ppt_buffer = create_km_ppt_export(dashboard_df, key_matrix_df, region, banker, strategist, filter_desc)
            filename = f"KeyMatrix_{region}_{banker}_{timestamp}.pptx"
            print(f"PPT export successful: {filename}")
            return dcc.send_bytes(ppt_buffer.getvalue(), filename)

    except Exception as e:
        print(f"Export Error: {e}")
        import traceback
        traceback.print_exc()
        return None

    return None


# ============================================================================
# COMBINED X-SELL EXPORT CALLBACK
# ============================================================================
@app.callback(
    Output("xs-download-dataframe", "data"),
    Input("xs-export-excel-btn", "n_clicks"),
    Input("xs-export-pdf-btn", "n_clicks"),
    Input("xs-export-ppt-btn", "n_clicks"),
    State('xs-region-filter', 'value'),
    State('xs-banker-filter', 'value'),
    State('xs-strategist-filter', 'value'),
    State('xs-product-filter', 'value'),
    State('xs-stage-filter', 'value'),
    prevent_initial_call=True
)
def export_xs_combined(excel_clicks, pdf_clicks, ppt_clicks, region, banker, strategist, product_filter, stage_filter):
    """Combined export callback for X-Sell - handles Excel, PDF, and PPT"""

    ctx = dash.callback_context
    if not ctx.triggered:
        return None

    button_id = ctx.triggered[0]['prop_id'].split('.')[0]

    try:
        df = pd.read_excel('data/crosssell_data.xlsx')

        # Generate filter description
        if strategist != 'All Strategists':
            filter_desc = f"Investment Strategist: {strategist}"
        elif region == 'All Regions':
            filter_desc = "All Regions | All Bankers"
        else:
            filter_desc = f"{region} | {banker}"

        if product_filter and product_filter != 'All Products':
            filter_desc += f" | Product: {product_filter}"
        if stage_filter and stage_filter != 'All Stages':
            filter_desc += f" | Stage: {stage_filter}"

        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

        # Excel Export
        if button_id == "xs-export-excel-btn":
            excel_buffer = create_xs_excel_export(df, region, banker, strategist,
                                                  product_filter, stage_filter, filter_desc)
            filename = f"CrossSell_{region}_{banker}_{timestamp}.xlsx"
            return dcc.send_bytes(excel_buffer.getvalue(), filename)

        # PDF Export
        elif button_id == "xs-export-pdf-btn":
            pdf_buffer = create_xs_pdf_export(df, region, banker, strategist,
                                              product_filter, stage_filter, filter_desc)
            filename = f"CrossSell_{region}_{banker}_{timestamp}.pdf"
            return dcc.send_bytes(pdf_buffer.getvalue(), filename)

        # PPT Export
        elif button_id == "xs-export-ppt-btn":
            ppt_buffer = create_xs_ppt_export(df, region, banker, strategist,
                                              product_filter, stage_filter, filter_desc)
            filename = f"CrossSell_{region}_{banker}_{timestamp}.pptx"
            return dcc.send_bytes(ppt_buffer.getvalue(), filename)

    except Exception as e:
        print(f"X-Sell Export Error: {e}")
        import traceback
        traceback.print_exc()
        return None

    return None


# ============================================================================
# COMBINED MANDATES EXPORT CALLBACK
# ============================================================================
@app.callback(
    Output("md-download-dataframe", "data"),
    Input("md-export-excel-btn", "n_clicks"),
    Input("md-export-pdf-btn", "n_clicks"),
    Input("md-export-ppt-btn", "n_clicks"),
    State('md-region-filter', 'value'),
    State('md-banker-filter', 'value'),
    State('md-strategist-filter', 'value'),
    prevent_initial_call=True
)
def export_md_combined(excel_clicks, pdf_clicks, ppt_clicks, region, banker, strategist):
    """Combined export callback for Mandates - handles Excel, PDF, and PPT"""

    ctx = dash.callback_context
    if not ctx.triggered:
        return None

    button_id = ctx.triggered[0]['prop_id'].split('.')[0]

    try:
        df = pd.read_excel('data/mandate_data.xlsx')

        # Generate filter description
        if strategist != 'All Strategists':
            filter_desc = f"Investment Strategist: {strategist}"
        elif region == 'All Regions':
            filter_desc = "All Regions | All Bankers"
        else:
            filter_desc = f"{region} | {banker}"

        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

        # Excel Export
        if button_id == "md-export-excel-btn":
            excel_buffer = create_md_excel_export(df, region, banker, strategist, filter_desc)
            filename = f"Mandates_{region}_{banker}_{timestamp}.xlsx"
            return dcc.send_bytes(excel_buffer.getvalue(), filename)

        # PDF Export
        elif button_id == "md-export-pdf-btn":
            pdf_buffer = create_md_pdf_export(df, region, banker, strategist, filter_desc)
            filename = f"Mandates_{region}_{banker}_{timestamp}.pdf"
            return dcc.send_bytes(pdf_buffer.getvalue(), filename)

        # PPT Export
        elif button_id == "md-export-ppt-btn":
            ppt_buffer = create_md_ppt_export(df, region, banker, strategist, filter_desc)
            filename = f"Mandates_{region}_{banker}_{timestamp}.pptx"
            return dcc.send_bytes(ppt_buffer.getvalue(), filename)

    except Exception as e:
        print(f"Mandates Export Error: {e}")
        import traceback
        traceback.print_exc()
        return None

    return None


# Key Matrix Tab Switching
@app.callback(
    Output('km-content', 'children', allow_duplicate=True),
    Input('km-tabs', 'active_tab'),
    State('km-region-filter', 'value'),
    State('km-banker-filter', 'value'),
    State('km-strategist-filter', 'value'),
    prevent_initial_call=True
)
def switch_km_tab(active_tab, region, banker, strategist):
    try:
        df = pd.read_excel('data/key_matrix_data.xlsx', sheet_name='Dashboard')
        filtered_df = filter_data_by_hierarchy(df, region, banker, strategist)

        if active_tab == "km-tab-dashboard":
            # Return the original dashboard view (your existing update_km_content logic)
            return update_km_content(region, banker, strategist, None, None)

        elif active_tab == "km-tab-trends":
            # Trend analysis with line charts
            return create_km_trend_view(filtered_df)

        elif active_tab == "km-tab-table":
            # Data table view
            return create_km_table_view(filtered_df)
    except:
        return dbc.Alert("Error loading data", color="warning")


def create_km_trend_view(df):
    """Create trend analysis view for Key Matrix - GRAY SHADES"""
    if df.empty:
        return dbc.Alert("No data available", color="info")

    # Create trend charts with gray shades
    fig = make_subplots(
        rows=2, cols=2,
        subplot_titles=('Families Trend', 'AUM Trend', 'Revenue Trend', 'Growth Rates'),
        specs=[[{"type": "scatter"}, {"type": "scatter"}],
               [{"type": "scatter"}, {"type": "bar"}]]
    )

    quarters = ['Q1', 'Q2', 'Q3', 'Q4']

    # Families trend - GRAY
    fig.add_trace(
        go.Scatter(x=quarters, y=[138, 165, 189, 204], mode='lines+markers',
                   name='Families', line=dict(color='#636e72', width=3),
                   marker=dict(size=8, color='#636e72')),
        row=1, col=1
    )

    # AUM trend - GRAY
    fig.add_trace(
        go.Scatter(x=quarters, y=[6500, 7200, 7900, 8525], mode='lines+markers',
                   name='AUM', line=dict(color='#7f8c8d', width=3),
                   marker=dict(size=8, color='#7f8c8d')),
        row=1, col=2
    )

    # Revenue trend - GRAY
    fig.add_trace(
        go.Scatter(x=quarters, y=[4.2, 5.8, 7.5, 10.6], mode='lines+markers',
                   name='Revenue', line=dict(color='#95a5a6', width=3),
                   marker=dict(size=8, color='#95a5a6')),
        row=2, col=1
    )

    # Growth rates - GRAY SHADES
    fig.add_trace(
        go.Bar(x=['Families', 'AUM', 'Revenue'], y=[27, 31, 56],
               marker_color=['#636e72', '#7f8c8d', '#95a5a6']),
        row=2, col=2
    )

    fig.update_layout(
        height=700,
        showlegend=False,
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=11)
    )

    fig.update_xaxes(showgrid=False, showline=True, linecolor='#ddd')
    fig.update_yaxes(showgrid=True, gridcolor='#f0f0f0')

    return html.Div([
        html.H4("Trend Analysis", style={"color": "#C41E3A", "marginBottom": "20px"}),
        dcc.Graph(figure=fig, config={'displayModeBar': False})
    ])


def create_km_table_view(df):
    """Create data table view for Key Matrix"""
    if df.empty:
        return dbc.Alert("No data available", color="info")

    # Create interactive data table
    return html.Div([
        html.H4("Data Table", style={"color": "#C41E3A", "marginBottom": "20px"}),
        html.Div([
            html.Table([
                html.Thead([
                    html.Tr([
                        html.Th(col, style={
                            "padding": "12px",
                            "backgroundColor": "#f8f9fa",
                            "borderBottom": "2px solid #dee2e6",
                            "fontSize": "12px",
                            "fontWeight": "bold"
                        }) for col in df.columns[:10]  # Show first 10 columns
                    ])
                ]),
                html.Tbody([
                    html.Tr([
                        html.Td(str(df.iloc[i][col]), style={
                            "padding": "10px",
                            "borderBottom": "1px solid #dee2e6",
                            "fontSize": "12px"
                        }) for col in df.columns[:10]
                    ], style={
                        "backgroundColor": "#f8f9fa" if i % 2 == 0 else "white"
                    }) for i in range(min(20, len(df)))  # Show first 20 rows
                ])
            ], style={
                "width": "100%",
                "borderCollapse": "collapse",
                "border": "1px solid #dee2e6"
            })
        ], style={"overflowX": "auto"}),
        html.P(f"Showing {min(20, len(df))} of {len(df)} rows",
               style={"marginTop": "10px", "fontSize": "11px", "color": "#666"})
    ])


# X-Sell Tab Switching
@app.callback(
    Output('xs-content', 'children', allow_duplicate=True),
    Input('xs-tabs', 'active_tab'),
    State('xs-region-filter', 'value'),
    State('xs-banker-filter', 'value'),
    State('xs-strategist-filter', 'value'),
    prevent_initial_call=True
)
def switch_xs_tab(active_tab, region, banker, strategist):
    try:
        df = pd.read_excel('data/crosssell_data.xlsx')
        filtered_df = filter_data_by_hierarchy(df, region, banker, strategist)

        if active_tab == "xs-tab-table":
            # Return the original table view
            return update_xs_content(region, banker, strategist, None, None)

        elif active_tab == "xs-tab-chart":
            # Chart view with interactive graphs
            return create_xs_chart_view(filtered_df)

        elif active_tab == "xs-tab-funnel":
            # Funnel analysis
            return create_xs_funnel_view(filtered_df)
    except:
        return dbc.Alert("Error loading data", color="warning")


def create_xs_chart_view(df):
    """Create chart view for X-Sell - GRAY SHADES"""
    if df.empty:
        return dbc.Alert("No data available", color="info")

    # Create stacked bar chart with gray shades
    stage_pivot = df.pivot_table(
        index='Product',
        columns='Stage',
        values='Amount_Cr',
        aggfunc='sum',
        fill_value=0
    ).reset_index()

    fig = go.Figure()

    # Gray shades for stages
    colors = {
        'Planned': '#636e72',
        'Proposed': '#95a5a6',
        'Transacted': '#b2bec3'
    }

    for stage in ['Planned', 'Proposed', 'Transacted']:
        if stage in stage_pivot.columns:
            fig.add_trace(go.Bar(
                name=stage,
                x=stage_pivot['Product'],
                y=stage_pivot[stage],
                marker_color=colors.get(stage, '#7f8c8d'),
                text=stage_pivot[stage].round(1),
                textposition='inside',
                textfont={'size': 10, 'color': 'white'}
            ))

    fig.update_layout(
        barmode='stack',
        title="Product-wise Stage Distribution",
        xaxis_title="Products",
        yaxis_title="Amount (Cr)",
        plot_bgcolor='white',
        paper_bgcolor='white',
        height=500,
        font=dict(family='Inter', size=11)
    )

    fig.update_xaxes(showgrid=False, tickangle=-45)
    fig.update_yaxes(showgrid=True, gridcolor='#f0f0f0')

    return html.Div([
        html.H4("Chart View", style={"color": "#C41E3A", "marginBottom": "20px"}),
        dcc.Graph(figure=fig, config={'displayModeBar': False})
    ])


def create_xs_funnel_view(df):
    """Create funnel analysis for X-Sell - GRAY SHADES"""
    if df.empty:
        return dbc.Alert("No data available", color="info")

    # Calculate funnel metrics
    stage_totals = df.groupby('Stage')['Amount_Cr'].sum()

    planned = stage_totals.get('Planned', 0)
    proposed = stage_totals.get('Proposed', 0)
    transacted = stage_totals.get('Transacted', 0)

    # Calculate conversion rates
    planned_to_proposed = (proposed / planned * 100) if planned > 0 else 0
    proposed_to_transacted = (transacted / proposed * 100) if proposed > 0 else 0
    overall_conversion = (transacted / planned * 100) if planned > 0 else 0
    # Create funnel chart
    # Create funnel chart - GRAY SHADES
    fig = go.Figure(go.Funnel(
        y=['Planned', 'Proposed', 'Transacted'],
        x=[planned, proposed, transacted],
        textposition="inside",
        textinfo="value+percent initial",
        marker={"color": ["#636e72", "#95a5a6", "#b2bec3"]},  # Gray shades
        connector={"line": {"color": "#ddd", "dash": "dot", "width": 2}}
    ))

    fig.update_layout(
        title="Sales Funnel Analysis",
        height=500,
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(family='Inter', size=12)
    )

    return html.Div([
        html.H4("Funnel Analysis", style={"color": "#C41E3A", "marginBottom": "20px"}),

        # Key Metrics Cards with conditional colors
        dbc.Row([
            dbc.Col([
                html.Div([
                    html.H6("Planned → Proposed", style={"fontSize": "12px", "color": "#666"}),
                    html.H3(f"{planned_to_proposed:.1f}%", style={
                        "color": get_percentage_color(planned_to_proposed),
                        "margin": "0",
                        "fontWeight": "bold"
                    })
                ], style={
                    "padding": "15px",
                    "backgroundColor": "white",
                    "border": "1px solid #ddd",
                    "borderRadius": "10px",
                    "textAlign": "center"
                })
            ], width=4),
            dbc.Col([
                html.Div([
                    html.H6("Proposed → Transacted", style={"fontSize": "12px", "color": "#666"}),
                    html.H3(f"{proposed_to_transacted:.1f}%", style={
                        "color": get_percentage_color(proposed_to_transacted),
                        "margin": "0",
                        "fontWeight": "bold"
                    })
                ], style={
                    "padding": "15px",
                    "backgroundColor": "white",
                    "border": "1px solid #ddd",
                    "borderRadius": "10px",
                    "textAlign": "center"
                })
            ], width=4),
            dbc.Col([
                html.Div([
                    html.H6("Overall Conversion", style={"fontSize": "12px", "color": "#666"}),
                    html.H3(f"{overall_conversion:.1f}%", style={
                        "color": get_percentage_color(overall_conversion),
                        "margin": "0",
                        "fontWeight": "bold"
                    })
                ], style={
                    "padding": "15px",
                    "backgroundColor": "white",
                    "border": "1px solid #ddd",
                    "borderRadius": "10px",
                    "textAlign": "center"
                })
            ], width=4)
        ], className="mb-4"),

        # Funnel Chart
        dcc.Graph(figure=fig, config={'displayModeBar': False})
    ])
# Mandates Tab Switching
@app.callback(
    Output('md-content', 'children', allow_duplicate=True),
    Input('md-tabs', 'active_tab'),
    State('md-region-filter', 'value'),
    State('md-banker-filter', 'value'),
    State('md-strategist-filter', 'value'),
    prevent_initial_call=True
)
def switch_md_tab(active_tab, region, banker, strategist):
    try:
        df = pd.read_excel('data/mandate_data.xlsx')
        filtered_df = filter_data_by_hierarchy(df, region, banker, strategist)

        if active_tab == "md-tab-timeline":
            # Return the original timeline view
            return update_md_content(region, banker, strategist, None, None)

        elif active_tab == "md-tab-trends":
            # Trend analysis
            return create_md_trend_view(filtered_df)

        elif active_tab == "md-tab-remarks":
            # Remarks view
            return create_md_remarks_view(filtered_df)
    except:
        return dbc.Alert("Error loading data", color="warning")


def create_md_trend_view(df):
    """Create trend analysis view for Mandates - GRAY SHADES"""
    if df.empty:
        return dbc.Alert("No data available", color="info")

    # Calculate stage counts over time
    stage_counts = df.groupby(['Quarter', 'Stage'])['Count'].sum().reset_index()

    # Create line chart with gray shades
    fig = go.Figure()

    # Gray shades for stages
    colors = {
        'Pitched': '#636e72',
        'Offered': '#7f8c8d',
        'Converted': '#95a5a6',
        'Declined': '#b2bec3'
    }

    for stage in ['Pitched', 'Offered', 'Converted', 'Declined']:
        stage_data = stage_counts[stage_counts['Stage'] == stage]
        if not stage_data.empty:
            fig.add_trace(go.Scatter(
                x=stage_data['Quarter'],
                y=stage_data['Count'],
                mode='lines+markers',
                name=stage,
                line=dict(color=colors.get(stage, '#7f8c8d'), width=3),
                marker=dict(size=8, color=colors.get(stage, '#7f8c8d'))
            ))

    fig.update_layout(
        title="Mandate Pipeline Trends",
        xaxis_title="Quarter",
        yaxis_title="Count",
        plot_bgcolor='white',
        paper_bgcolor='white',
        height=500,
        font=dict(family='Inter', size=11),
        hovermode='x unified',
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=-0.3,
            xanchor="center",
            x=0.5
        )
    )

    fig.update_xaxes(showgrid=False, showline=True, linecolor='#ddd')
    fig.update_yaxes(showgrid=True, gridcolor='#f0f0f0')

    # Calculate win rate
    total_pitched = df[df['Stage'] == 'Pitched']['Count'].sum()
    total_converted = df[df['Stage'] == 'Converted']['Count'].sum()
    total_declined = df[df['Stage'] == 'Declined']['Count'].sum()
    win_rate = (total_converted / total_pitched * 100) if total_pitched > 0 else 0
    loss_rate = (total_declined / total_pitched * 100) if total_pitched > 0 else 0

    return html.Div([
        html.H4("Trend Analysis", style={"color": "#C41E3A", "marginBottom": "20px"}),

        # Key Metrics with conditional colors
        dbc.Row([
            dbc.Col([
                html.Div([
                    html.I(className="fas fa-bullseye", style={"fontSize": "24px", "color": "#636e72"}),
                    html.H3(f"{win_rate:.1f}%", style={
                        "color": get_percentage_color(win_rate),
                        "margin": "10px 0 5px 0",
                        "fontWeight": "bold"
                    }),
                    html.P("Win Rate", style={"fontSize": "12px", "color": "#666", "margin": "0"})
                ], style={
                    "padding": "20px",
                    "backgroundColor": "white",
                    "border": "1px solid #ddd",
                    "borderRadius": "10px",
                    "textAlign": "center"
                })
            ], width=3),
            dbc.Col([
                html.Div([
                    html.I(className="fas fa-handshake", style={"fontSize": "24px", "color": "#636e72"}),
                    html.H3(str(total_converted), style={
                        "color": "#636e72",
                        "margin": "10px 0 5px 0",
                        "fontWeight": "bold"
                    }),
                    html.P("Converted", style={"fontSize": "12px", "color": "#666", "margin": "0"})
                ], style={
                    "padding": "20px",
                    "backgroundColor": "white",
                    "border": "1px solid #ddd",
                    "borderRadius": "10px",
                    "textAlign": "center"
                })
            ], width=3),
            dbc.Col([
                html.Div([
                    html.I(className="fas fa-hourglass-half", style={"fontSize": "24px", "color": "#636e72"}),
                    html.H3(str(df[df['Stage'] == 'Offered']['Count'].sum()),
                            style={"color": "#636e72", "margin": "10px 0 5px 0", "fontWeight": "bold"}),
                    html.P("In Pipeline", style={"fontSize": "12px", "color": "#666", "margin": "0"})
                ], style={
                    "padding": "20px",
                    "backgroundColor": "white",
                    "border": "1px solid #ddd",
                    "borderRadius": "10px",
                    "textAlign": "center"
                })
            ], width=3),
            dbc.Col([
                html.Div([
                    html.I(className="fas fa-times-circle", style={"fontSize": "24px", "color": "#636e72"}),
                    html.H3(str(total_declined), style={
                        "color": "#636e72",
                        "margin": "10px 0 5px 0",
                        "fontWeight": "bold"
                    }),
                    html.P([
                        "Declined ",
                        html.Span(f"({loss_rate:.1f}%)", style={
                            "color": get_percentage_color(-loss_rate),  # Negative for red
                            "fontWeight": "bold"
                        })
                    ], style={"fontSize": "12px", "color": "#666", "margin": "0"})
                ], style={
                    "padding": "20px",
                    "backgroundColor": "white",
                    "border": "1px solid #ddd",
                    "borderRadius": "10px",
                    "textAlign": "center"
                })
            ], width=3)
        ], className="mb-4"),

        # Trend Chart
        dcc.Graph(figure=fig, config={'displayModeBar': False})
    ])

def create_md_remarks_view(filtered_df):
    """Create remarks and notes view for Mandates"""
    user_id = current_user.get('UserID', '')
    remarks_file = f'data/user_data/{user_id}_mandate_remarks.xlsx'

    # Load existing remarks if available
    remarks_data = []
    if os.path.exists(remarks_file):
        remarks_df = pd.read_excel(remarks_file)
        remarks_data = remarks_df.to_dict('records')

    return html.Div([
        html.H4("Remarks & Notes", style={"color": "#C41E3A", "marginBottom": "20px"}),

        html.P("Track your progress and insights for each stage",
               style={"fontSize": "13px", "color": "#666", "marginBottom": "20px"}),

        # Remarks Form (same as before but enhanced styling)
        html.Div([
            html.Div([
                html.Div([
                    html.I(className="fas fa-lightbulb", style={"marginRight": "8px", "color": "#C41E3A"}),
                    html.H5("Pitched Stage", style={"display": "inline", "fontSize": "15px"})
                ], style={"marginBottom": "10px"}),
                dcc.Textarea(
                    id="pitched-remarks",
                    placeholder="Enter remarks for Pitched stage...",
                    value=remarks_data[-1]['Pitched_Remarks'] if remarks_data else "",
                    style={
                        "width": "100%",
                        "height": "80px",
                        "border": "1px solid #ddd",
                        "borderRadius": "8px",
                        "padding": "10px",
                        "fontSize": "12px",
                        "resize": "vertical"
                    }
                )
            ], style={
                "padding": "15px",
                "backgroundColor": "white",
                "border": "1px solid #ddd",
                "borderRadius": "10px",
                "marginBottom": "15px"
            }),

            html.Div([
                html.Div([
                    html.I(className="fas fa-file-invoice", style={"marginRight": "8px", "color": "#f0ad4e"}),
                    html.H5("Offered Stage", style={"display": "inline", "fontSize": "15px"})
                ], style={"marginBottom": "10px"}),
                dcc.Textarea(
                    id="offered-remarks",
                    placeholder="Enter remarks for Offered stage...",
                    value=remarks_data[-1]['Offered_Remarks'] if remarks_data else "",
                    style={
                        "width": "100%",
                        "height": "80px",
                        "border": "1px solid #ddd",
                        "borderRadius": "8px",
                        "padding": "10px",
                        "fontSize": "12px",
                        "resize": "vertical"
                    }
                )
            ], style={
                "padding": "15px",
                "backgroundColor": "white",
                "border": "1px solid #ddd",
                "borderRadius": "10px",
                "marginBottom": "15px"
            }),

            html.Div([
                html.Div([
                    html.I(className="fas fa-check-circle", style={"marginRight": "8px", "color": "#28a745"}),
                    html.H5("Converted Stage", style={"display": "inline", "fontSize": "15px"})
                ], style={"marginBottom": "10px"}),
                dcc.Textarea(
                    id="converted-remarks",
                    placeholder="Enter remarks for Converted stage...",
                    value=remarks_data[-1]['Converted_Remarks'] if remarks_data else "",
                    style={
                        "width": "100%",
                        "height": "80px",
                        "border": "1px solid #ddd",
                        "borderRadius": "8px",
                        "padding": "10px",
                        "fontSize": "12px",
                        "resize": "vertical"
                    }
                )
            ], style={
                "padding": "15px",
                "backgroundColor": "white",
                "border": "1px solid #ddd",
                "borderRadius": "10px",
                "marginBottom": "15px"
            }),

            html.Div([
                html.Div([
                    html.I(className="fas fa-times-circle", style={"marginRight": "8px", "color": "#dc3545"}),
                    html.H5("Declined Stage", style={"display": "inline", "fontSize": "15px"})
                ], style={"marginBottom": "10px"}),
                dcc.Textarea(
                    id="declined-remarks",
                    placeholder="Enter remarks for Declined stage...",
                    value=remarks_data[-1]['Declined_Remarks'] if remarks_data else "",
                    style={
                        "width": "100%",
                        "height": "80px",
                        "border": "1px solid #ddd",
                        "borderRadius": "8px",
                        "padding": "10px",
                        "fontSize": "12px",
                        "resize": "vertical"
                    }
                )
            ], style={
                "padding": "15px",
                "backgroundColor": "white",
                "border": "1px solid #ddd",
                "borderRadius": "10px",
                "marginBottom": "15px"
            }),

            # Save Button
            html.Div([
                dbc.Button([
                    html.I(className="fas fa-save", style={"marginRight": "8px"}),
                    "Save Remarks"
                ], id="save-mandate-remarks", color="primary", size="lg",
                    style={"fontSize": "14px", "padding": "10px 30px"})
            ], style={"textAlign": "center"})
        ])
    ])


# @app.callback(
#     Output('km-quick-stats', 'children'),
#     Input('km-region-filter', 'value'),
#     Input('km-banker-filter', 'value'),
#     Input('km-strategist-filter', 'value')
# )
# def update_km_quick_stats(region, banker, strategist):
#     try:
#         df = pd.read_excel('data/dashboard_data.xlsx', sheet_name='Dashboard')
#         filtered_df = filter_data_by_hierarchy(df, region, banker, strategist)
#
#         if filtered_df.empty:
#             return html.P("No data", style={"fontSize": "11px", "color": "#999"})
#
#         # Calculate quick stats
#         total_families = int(filtered_df['Families'].sum()) if 'Families' in filtered_df.columns else 0
#         total_aum = int(filtered_df['AUM_Cr'].sum()) if 'AUM_Cr' in filtered_df.columns else 0
#         avg_aum_per_family = int(total_aum / total_families) if total_families > 0 else 0
#
#         return html.Div([
#             html.Div([
#                 html.Span("👥", style={"fontSize": "16px", "marginRight": "6px"}),
#                 html.Span(f"{total_families} Families", style={"fontSize": "11px", "fontWeight": "600"})
#             ], style={"marginBottom": "8px"}),
#             html.Div([
#                 html.Span("💰", style={"fontSize": "16px", "marginRight": "6px"}),
#                 html.Span(f"₹{total_aum} Cr AUM", style={"fontSize": "11px", "fontWeight": "600"})
#             ], style={"marginBottom": "8px"}),
#             html.Div([
#                 html.Span("📊", style={"fontSize": "16px", "marginRight": "6px"}),
#                 html.Span(f"₹{avg_aum_per_family} Cr Avg", style={"fontSize": "11px", "fontWeight": "600"})
#             ])
#         ])
#     except:
#         return html.P("Error loading stats", style={"fontSize": "11px", "color": "#dc3545"})
#

@app.callback(
    Output('xs-stage-progress', 'children'),
    Input('xs-region-filter', 'value'),
    Input('xs-banker-filter', 'value'),
    Input('xs-strategist-filter', 'value')
)
def update_xs_stage_progress(region, banker, strategist):
    try:
        df = pd.read_excel('data/crosssell_data.xlsx')
        filtered_df = filter_data_by_hierarchy(df, region, banker, strategist)

        if filtered_df.empty:
            return html.P("No data", style={"fontSize": "11px", "color": "#999"})

        # Calculate stage percentages
        stage_totals = filtered_df.groupby('Stage')['Amount_Cr'].sum()
        total = stage_totals.sum()

        planned_pct = (stage_totals.get('Planned', 0) / total * 100) if total > 0 else 0
        proposed_pct = (stage_totals.get('Proposed', 0) / total * 100) if total > 0 else 0
        transacted_pct = (stage_totals.get('Transacted', 0) / total * 100) if total > 0 else 0

        return html.Div([
            # Planned
            html.Div([
                html.Div("Planned", style={"fontSize": "11px", "marginBottom": "4px"}),
                html.Div([
                    html.Div(style={
                        "width": f"{planned_pct}%",
                        "height": "8px",
                        "backgroundColor": "#C41E3A",
                        "borderRadius": "4px",
                        "transition": "width 0.3s ease"
                    })
                ], style={
                    "width": "100%",
                    "height": "8px",
                    "backgroundColor": "#f0f0f0",
                    "borderRadius": "4px",
                    "overflow": "hidden"
                }),
                html.Span(f"{planned_pct:.1f}%", style={"fontSize": "10px", "color": "#666"})
            ], style={"marginBottom": "12px"}),

            # Proposed
            html.Div([
                html.Div("Proposed", style={"fontSize": "11px", "marginBottom": "4px"}),
                html.Div([
                    html.Div(style={
                        "width": f"{proposed_pct}%",
                        "height": "8px",
                        "backgroundColor": "#f0ad4e",
                        "borderRadius": "4px",
                        "transition": "width 0.3s ease"
                    })
                ], style={
                    "width": "100%",
                    "height": "8px",
                    "backgroundColor": "#f0f0f0",
                    "borderRadius": "4px",
                    "overflow": "hidden"
                }),
                html.Span(f"{proposed_pct:.1f}%", style={"fontSize": "10px", "color": "#666"})
            ], style={"marginBottom": "12px"}),

            # Transacted
            html.Div([
                html.Div("Transacted", style={"fontSize": "11px", "marginBottom": "4px"}),
                html.Div([
                    html.Div(style={
                        "width": f"{transacted_pct}%",
                        "height": "8px",
                        "backgroundColor": "#28a745",
                        "borderRadius": "4px",
                        "transition": "width 0.3s ease"
                    })
                ], style={
                    "width": "100%",
                    "height": "8px",
                    "backgroundColor": "#f0f0f0",
                    "borderRadius": "4px",
                    "overflow": "hidden"
                }),
                html.Span(f"{transacted_pct:.1f}%", style={"fontSize": "10px", "color": "#666"})
            ])
        ])
    except:
        return html.P("Error loading progress", style={"fontSize": "11px", "color": "#dc3545"})


@app.callback(
    [Output('md-conversion-funnel', 'children'),
     Output('md-win-loss', 'children')],
    Input('md-region-filter', 'value'),
    Input('md-banker-filter', 'value'),
    Input('md-strategist-filter', 'value')
)
def update_md_sidebar_widgets(region, banker, strategist):
    try:
        df = pd.read_excel('data/mandate_data.xlsx')
        filtered_df = filter_data_by_hierarchy(df, region, banker, strategist)

        if filtered_df.empty:
            return html.P("No data", style={"fontSize": "11px", "color": "#999"}), \
                html.P("No data", style={"fontSize": "11px", "color": "#999"})

        # Calculate funnel metrics
        stage_counts = filtered_df.groupby('Stage')['Count'].sum()
        pitched = stage_counts.get('Pitched', 0)
        offered = stage_counts.get('Offered', 0)
        converted = stage_counts.get('Converted', 0)
        declined = stage_counts.get('Declined', 0)

        # Conversion funnel widget
        funnel_widget = html.Div([
            html.Div([
                html.Span("Pitched", style={"fontSize": "10px"}),
                html.Span(str(pitched), style={"fontSize": "12px", "fontWeight": "bold", "float": "right"})
            ], style={"marginBottom": "6px"}),
            html.Div([
                html.Span("Offered", style={"fontSize": "10px"}),
                html.Span(str(offered), style={"fontSize": "12px", "fontWeight": "bold", "float": "right"})
            ], style={"marginBottom": "6px", "paddingLeft": "10px"}),
            html.Div([
                html.Span("Converted", style={"fontSize": "10px", "color": "#28a745"}),
                html.Span(str(converted), style={
                    "fontSize": "12px",
                    "fontWeight": "bold",
                    "float": "right",
                    "color": "#28a745"
                })
            ], style={"paddingLeft": "20px"})
        ])

        # Win/Loss widget
        win_rate = (converted / pitched * 100) if pitched > 0 else 0
        loss_rate = (declined / pitched * 100) if pitched > 0 else 0

        win_loss_widget = html.Div([
            html.Div([
                html.Div([
                    html.Span("✓", style={"color": "#28a745", "fontSize": "16px", "marginRight": "6px"}),
                    html.Span(f"{win_rate:.1f}%", style={"fontSize": "14px", "fontWeight": "bold", "color": "#28a745"})
                ]),
                html.Span("Win Rate", style={"fontSize": "10px", "color": "#666"})
            ], style={
                "padding": "10px",
                "backgroundColor": "rgba(40, 167, 69, 0.1)",
                "borderRadius": "8px",
                "marginBottom": "8px"
            }),
            html.Div([
                html.Div([
                    html.Span("✗", style={"color": "#dc3545", "fontSize": "16px", "marginRight": "6px"}),
                    html.Span(f"{loss_rate:.1f}%", style={"fontSize": "14px", "fontWeight": "bold", "color": "#dc3545"})
                ]),
                html.Span("Loss Rate", style={"fontSize": "10px", "color": "#666"})
            ], style={
                "padding": "10px",
                "backgroundColor": "rgba(220, 53, 69, 0.1)",
                "borderRadius": "8px"
            })
        ])

        return funnel_widget, win_loss_widget
    except:
        return html.P("Error", style={"fontSize": "11px", "color": "#dc3545"}), \
            html.P("Error", style={"fontSize": "11px", "color": "#dc3545"})
# Add callback for saving remarks
@app.callback(
    Output("proposal-message", "children", allow_duplicate=True),
    [Input("save-mandate-remarks", "n_clicks")],
    [State("pitched-remarks", "value"),
     State("offered-remarks", "value"),
     State("converted-remarks", "value"),
     State("declined-remarks", "value")],
    prevent_initial_call=True
)
def save_mandate_remarks(n_clicks, pitched_remarks, offered_remarks, converted_remarks, declined_remarks):
    if not n_clicks:
        return ""

    user_id = current_user.get('UserID', '')
    if not user_id:
        return dbc.Alert("User session expired. Please login again.", color="warning")

    try:
        remarks_data = {
            'Timestamp': [datetime.now().strftime('%Y-%m-%d %H:%M:%S')],
            'Pitched_Remarks': [pitched_remarks or ''],
            'Offered_Remarks': [offered_remarks or ''],
            'Converted_Remarks': [converted_remarks or ''],
            'Declined_Remarks': [declined_remarks or '']
        }

        # Save mandate remarks
        filename = f'data/user_data/{user_id}_mandate_remarks.xlsx'
        if os.path.exists(filename):
            existing_df = pd.read_excel(filename)
            combined_df = pd.concat([existing_df, pd.DataFrame(remarks_data)], ignore_index=True)
        else:
            combined_df = pd.DataFrame(remarks_data)
        combined_df.to_excel(filename, index=False)

        return dbc.Alert([
            html.H6("Mandate Remarks Saved!", className="alert-heading"),
            html.P("Your remarks have been successfully saved.", className="mb-0")
        ], color="success")

    except Exception as e:
        return dbc.Alert(f"Error saving remarks: {str(e)}", color="danger")
def get_cross_sell_layout():
    return html.Div([
        html.Div([
            dbc.Container([
                dbc.Row([
                    dbc.Col([
                        html.Div([
                            html.H2("Cross Sell Management", className="text-center mb-2"),
                            html.P("Streamline your cross-selling workflow", className="text-center text-muted mb-5",
                                   style={"fontSize": "1.1rem", "fontWeight": "400"})
                        ], className="mb-5"),
                        dbc.Tabs([
                            dbc.Tab(label="Manual Entry", tab_id="manual-tab"),
                            dbc.Tab(label="Excel Upload", tab_id="upload-tab")
                        ], id="cross-sell-tabs", active_tab="manual-tab", className="mb-4"),
                        html.Div(id="manual-entry-content", children=[
                            dbc.Card([
                                dbc.CardHeader([html.H5("Manual Cross Sell Entry", className="mb-0")]),
                                dbc.CardBody([
                                    dbc.Form([
                                        dbc.Row([
                                            dbc.Label("Family Name", width=3, className="form-label"),
                                            dbc.Col([dbc.Input(id="family-name", type="text",
                                                               placeholder="Enter family/client name",
                                                               className="form-control")], width=9)
                                        ], className="mb-4"),
                                        dbc.Row([
                                            dbc.Label("Product", width=3, className="form-label"),
                                            dbc.Col([dbc.Input(id="product", type="text",
                                                               placeholder="Enter product/service name",
                                                               className="form-control")], width=9)
                                        ], className="mb-4"),
                                        dbc.Row([
                                            dbc.Label("Stage", width=3, className="form-label"),
                                            dbc.Col([
                                                dcc.Dropdown(
                                                    id="stage",
                                                    options=[
                                                        {"label": "Planned", "value": "Planned"},
                                                        {"label": "Proposed", "value": "Proposed"},
                                                        {"label": "Confirmed", "value": "Confirmed"},
                                                        {"label": "Executed", "value": "Executed"}
                                                    ],
                                                    placeholder="Select current stage"
                                                )
                                            ], width=9)
                                        ], className="mb-4"),
                                        dbc.Row([
                                            dbc.Label("New Liquidity", width=3, className="form-label"),
                                            dbc.Col([dbc.Input(id="new-liquidity", type="number",
                                                               placeholder="Enter amount (e.g., 50000)",
                                                               className="form-control")], width=9)
                                        ], className="mb-4"),
                                        dbc.Row([
                                            dbc.Col([
                                                dbc.Button("Submit Entry", id="submit-cross-sell",
                                                           className="btn-primary me-3",
                                                           style={"fontSize": "13px", "padding": "10px 20px"}),
                                                dbc.Button("Back to Dashboard", id="back-to-dashboard",
                                                           className="btn-secondary")
                                            ])
                                        ])
                                    ])
                                ])
                            ], className="card shadow-lg")
                        ]),
                        html.Div(id="upload-entry-content", style={"display": "none"}, children=[
                            dbc.Card([
                                dbc.CardHeader([html.H5("Bulk Excel Upload", className="mb-0")]),
                                dbc.CardBody([
                                    dbc.Alert([
                                        html.H6("Excel Format Requirements", className="alert-heading"),
                                        html.P("Your Excel file must contain exactly these columns:", className="mb-3"),
                                        dbc.Row([
                                            dbc.Col([html.Ul([html.Li("Family_Name"), html.Li("Product")])], width=6),
                                            dbc.Col([html.Ul([html.Li("Stage (Planned/Proposed/Confirmed/Executed)"),
                                                              html.Li("New_Liquidity (numeric value)")])], width=6)
                                        ]),
                                        html.P("Timestamp will be added automatically when you upload.",
                                               className="mb-3 text-muted")
                                    ], color="light", className="mb-4"),
                                    dcc.Upload(
                                        id="upload-excel",
                                        children=html.Div([
                                            html.I(className="fas fa-cloud-upload-alt",
                                                   style={"fontSize": "4rem", "marginBottom": "1rem",
                                                          "color": "var(--gold)"}),
                                            html.H5("Drag & Drop Your Excel File Here", className="mb-2"),
                                            html.P("or click to browse files", className="text-muted mb-2"),
                                            html.Small("Supports .xlsx and .xls files", className="text-muted")
                                        ], style={"textAlign": "center", "padding": "3rem 2rem"}),
                                        className="upload-area mb-4",
                                        multiple=False,
                                        accept=".xlsx,.xls"
                                    ),
                                    html.Div(id="upload-status"),
                                    dbc.Button("Back to Dashboard", id="back-to-dashboard-upload",
                                               className="btn-secondary")
                                ])
                            ], className="card shadow-lg")
                        ])
                    ], width=12)
                ], justify="center")
            ], fluid=True)
        ], className="content-section")
    ], className="fade-in-up", style={"paddingTop": "2rem"})


# def get_proposal_layout():
#     return html.Div([
#         html.Div([
#             dbc.Container([
#                 dbc.Row([
#                     dbc.Col([
#                         html.Div([
#                             html.H2("Proposal Management", className="text-center mb-2"),
#                             html.P("Create and manage client proposals", className="text-center text-muted mb-5",
#                                    style={"fontSize": "1.1rem", "fontWeight": "400"})
#                         ], className="mb-5"),
#                         dbc.Card([
#                             dbc.CardHeader([html.H5("Client Proposal Entry", className="mb-0")]),
#                             dbc.CardBody([
#                                 dbc.Form([
#                                     dbc.Row([
#                                         dbc.Label("Name of Family", width=3, className="form-label"),
#                                         dbc.Col([dbc.Input(id="prop-family-name", type="text",
#                                                            placeholder="Enter family/client name",
#                                                            className="form-control")], width=9)
#                                     ], className="mb-3"),
#                                     dbc.Row([
#                                         dbc.Label("Banker", width=3, className="form-label"),
#                                         dbc.Col([dbc.Input(id="prop-banker", type="text",
#                                                            placeholder="Enter banker name", className="form-control")],
#                                                 width=9)
#                                     ], className="mb-3"),
#                                     dbc.Row([
#                                         dbc.Label("Date", width=3, className="form-label"),
#                                         dbc.Col([dbc.Input(id="prop-date", type="date", className="form-control")],
#                                                 width=9)
#                                     ], className="mb-3"),
#                                     dbc.Row([
#                                         dbc.Label("Amount", width=3, className="form-label"),
#                                         dbc.Col([dbc.Input(id="prop-amount", type="number", placeholder="Enter amount",
#                                                            className="form-control")], width=9)
#                                     ], className="mb-3"),
#                                     dbc.Row([
#                                         dbc.Label("Risk Profile", width=3, className="form-label"),
#                                         dbc.Col([dbc.Input(id="prop-risk-profile", type="text",
#                                                            placeholder="Enter risk profile", className="form-control")],
#                                                 width=9)
#                                     ], className="mb-3"),
#                                     dbc.Row([
#                                         dbc.Label("Status", width=3, className="form-label"),
#                                         dbc.Col([
#                                             dcc.Dropdown(
#                                                 id="prop-status",
#                                                 options=[
#                                                     {"label": "Pitched", "value": "Pitched"},
#                                                     {"label": "Follow up", "value": "Follow up"},
#                                                     {"label": "Offer", "value": "Offer"},
#                                                     {"label": "Converted", "value": "Converted"}
#                                                 ],
#                                                 placeholder="Select status"
#                                             )
#                                         ], width=9)
#                                     ], className="mb-4"),
#                                     dbc.Row([
#                                         dbc.Col([
#                                             dbc.Button("Submit Proposal", id="submit-proposal",
#                                                        className="btn-primary me-3"),
#                                             dbc.Button("Back to Dashboard", id="back-to-dashboard-prop",
#                                                        className="btn-secondary")
#                                         ])
#                                     ])
#                                 ])
#                             ])
#                         ], className="card shadow-lg")
#                     ], width=12)
#                 ], justify="center")
#             ], fluid=True)
#         ], className="content-section")
#     ], className="fade-in-up", style={"paddingTop": "2rem"})
#
def get_proposal_layout():
    return html.Div([
        html.Div([
            html.H2("Proposal Management", className="text-center mb-2"),
            html.P("Create and manage client proposals", className="text-center text-muted mb-5",
                   style={"fontSize": "1.1rem", "fontWeight": "400"})
        ], className="mb-5"),

        dbc.Row([
            dbc.Col([
                dbc.Card([
                    dbc.CardHeader([html.H5("Client Proposal Entry", className="mb-0")]),
                    dbc.CardBody([
                        # Client Search Section
                        html.Div([
                            html.H6("Client Search", className="mb-3 text-primary"),
                            dbc.Row([
                                dbc.Label("Search Client", width=3, className="form-label"),
                                dbc.Col([
                                    dbc.InputGroup([
                                        dbc.Input(
                                            id="client-search",
                                            type="text",
                                            placeholder="Type client name to search...",
                                            className="form-control"
                                        ),
                                        dbc.Button("Search", id="search-btn", color="outline-primary")
                                    ])
                                ], width=9)
                            ], className="mb-3"),

                            # Search Results
                            html.Div(id="search-results", className="mb-4"),

                            # Option to create new client
                            dbc.Row([
                                dbc.Col([
                                    dbc.Button("Create New Client", id="new-client-btn",
                                               color="success", outline=True, className="me-2"),
                                    dbc.Button("Clear Form", id="clear-form-btn",
                                               color="secondary", outline=True)
                                ])
                            ], className="mb-4"),

                            html.Hr()
                        ]),

                        # Proposal Form
                        dbc.Form([
                            html.H6("Proposal Information", className="mb-4 text-primary"),

                            # Hidden field to track if this is an update
                            dcc.Store(id="is-update-store", data=False),
                            dcc.Store(id="selected-client-store", data=None),

                            dbc.Row([
                                dbc.Label("Name of Family", width=3, className="form-label"),
                                dbc.Col([
                                    dbc.Input(id="prop-family-name", type="text",
                                              placeholder="Enter family/client name",
                                              className="form-control"),
                                    html.Small(id="family-name-helper", className="text-muted")
                                ], width=9)
                            ], className="mb-3"),

                            dbc.Row([
                                dbc.Label("Banker", width=3, className="form-label"),
                                dbc.Col([dbc.Input(id="prop-banker", type="text",
                                                   placeholder="Enter banker name",
                                                   className="form-control")], width=9)
                            ], className="mb-3"),

                            dbc.Row([
                                dbc.Label("Date", width=3, className="form-label"),
                                dbc.Col([dbc.Input(id="prop-date", type="date",
                                                   className="form-control")], width=9)
                            ], className="mb-3"),

                            dbc.Row([
                                dbc.Label("Amount", width=3, className="form-label"),
                                dbc.Col([dbc.Input(id="prop-amount", type="number",
                                                   placeholder="Enter amount",
                                                   className="form-control")], width=9)
                            ], className="mb-3"),

                            dbc.Row([
                                dbc.Label("Risk Profile", width=3, className="form-label"),
                                dbc.Col([dbc.Input(id="prop-risk-profile", type="text",
                                                   placeholder="Enter risk profile",
                                                   className="form-control")], width=9)
                            ], className="mb-3"),

                            dbc.Row([
                                dbc.Label("Type of Mandate", width=3, className="form-label"),
                                dbc.Col([
                                    dcc.Dropdown(
                                        id="prop-mandate-type",
                                        options=[
                                            {"label": "Distribution", "value": "Distribution"},
                                            {"label": "Advisory", "value": "Advisory"},
                                            {"label": "Family Office", "value": "Family Office"},
                                            {"label": "Treasury-Product", "value": "Treasury-Product"},
                                            {"label": "Treasury-Family Office", "value": "Treasury-Family Office"},
                                            {"label": "Product", "value": "Product"},
                                            {"label": "Others", "value": "Others"}
                                        ],
                                        placeholder="Select mandate type"
                                    )
                                ], width=9)
                            ], className="mb-4"),

                            # Current Status Display (for existing clients)
                            html.Div(id="current-status-display", className="mb-4"),

                            html.H6("Upload Proposal", className="mb-3 text-primary"),
                            dbc.Row([
                                dbc.Col([
                                    dcc.Upload(
                                        id="upload-proposal-pdf",
                                        children=html.Div([
                                            html.Div("📄", style={"fontSize": "2rem", "marginBottom": "0.5rem",
                                                                 "color": "var(--primary-dark-gray)"}),
                                            html.P("Drag & Drop PDF or Click to Browse", className="mb-0"),
                                            html.Small("PDF files only", className="text-muted")
                                        ], style={"textAlign": "center", "padding": "2rem"}),
                                        className="upload-area mb-4",
                                        multiple=False,
                                        accept=".pdf"
                                    )
                                ], width=12)
                            ]),

                            html.H6("Status Update", className="mb-3 text-primary"),
                            dbc.Row([
                                dbc.Label("New Status", width=3, className="form-label"),
                                dbc.Col([
                                    dcc.Dropdown(
                                        id="prop-status",
                                        options=[
                                            {"label": "Pitched", "value": "Pitched"},
                                            {"label": "Follow up", "value": "Follow up"},
                                            {"label": "Offer", "value": "Offer"},
                                            {"label": "Converted", "value": "Converted"}
                                        ],
                                        placeholder="Select new status"
                                    )
                                ], width=9)
                            ], className="mb-4"),

                            html.Div(id="status-form-content"),

                            dbc.Row([
                                dbc.Col([
                                    dbc.Button("Update Proposal", id="submit-proposal",
                                               color="primary", className="me-3 btn-lg",
                                               style={"fontSize": "14px", "padding": "10px 24px"}),
                                    dbc.Button("Back to Dashboard", id="back-to-dashboard-prop",
                                               color="secondary", className="btn-lg")
                                ])
                            ])
                        ])
                    ])
                ], className="card shadow-lg")
            ], width=12)
        ])
    ], className="fade-in-up")

# Callbacks
# Enhanced client search callback
@app.callback(
    Output("search-results", "children"),
    [Input("search-btn", "n_clicks")],
    [State("client-search", "value")]
)
def search_clients(n_clicks, search_term):
    if not n_clicks or not search_term:
        return ""

    existing_clients = search_existing_clients(search_term)

    if not existing_clients:
        return dbc.Alert("No existing clients found with that name.", color="info")

    client_cards = []
    for client in existing_clients:
        card = dbc.Card([
            dbc.CardBody([
                html.H6(client.get('Family_Name', 'Unknown'), className="card-title"),
                html.P([
                    html.Strong("Current Status: "),
                    html.Span(client.get('Status', 'Unknown'),
                              className="badge bg-primary ms-2")
                ], className="mb-2"),
                html.P([
                    html.Strong("Banker: "), client.get('Banker', 'Unknown'), html.Br(),
                    html.Strong("Last Updated: "),
                    client.get('Timestamp', 'Unknown')[:10] if client.get('Timestamp') else 'Unknown'
                ], className="text-muted small mb-2"),
                dbc.Button("Select Client", id={"type": "select-client", "index": client.get('Family_Name')},
                           color="primary", size="sm")
            ])
        ], className="mb-2 client-card", style={"border": "1px solid #dee2e6"})
        client_cards.append(card)

    return html.Div([
        html.H6(f"Found {len(existing_clients)} existing client(s):", className="mb-3"),
        html.Div(client_cards)
    ])


# Callback 1: Update banker dropdown based on region
# Callback 1: Update banker dropdown based on region (unchanged)
@app.callback(
    Output('banker-filter', 'options'),
    Output('banker-filter', 'value'),
    Input('region-filter', 'value'),
    State('region-banker-map', 'data')
)
def update_banker_dropdown(selected_region, region_banker_map):
    """Update banker options when region changes"""
    if not selected_region or not region_banker_map:
        return [{'label': 'All Bankers', 'value': 'All Bankers'}], 'All Bankers'

    bankers = region_banker_map.get(selected_region, ['All Bankers'])
    banker_options = [{'label': b, 'value': b} for b in bankers]

    return banker_options, 'All Bankers'


# Callback 2: Main KPI update with independent strategist filter
# ============================================================================
# UPDATED MAIN DASHBOARD CALLBACK - Uses All New Excel Columns
# ============================================================================
# Replace the existing update_dashboard_kpis callback (around line 1800)

@app.callback(
    Output('dashboard-kpi-content', 'children'),
    Input('region-filter', 'value'),
    Input('banker-filter', 'value'),
    Input('strategist-filter', 'value')
)
def update_dashboard_kpis(selected_region, selected_banker, selected_strategist):
    """Filter and display KPIs based on filters - Uses REAL Excel Data"""

    # Helper functions
    def pick_from_row(row, candidates, default=None):
        for c in candidates:
            if c in row and pd.notna(row[c]):
                return row[c]
        return default

    def as_int(val, default=0):
        try:
            if pd.isna(val):
                return default
            return int(float(val))
        except Exception:
            return default

    def as_float(val, default=0.0):
        try:
            if pd.isna(val):
                return default
            return float(val)
        except Exception:
            return default

    def format_pct(val):
        """Format percentage values and return tuple (text, color)"""
        if val is None or (isinstance(val, float) and np.isnan(val)):
            return ("", "#666")
        s = str(val).strip()

        try:
            num_str = s.replace('%', '').replace('YoY', '').strip()
            f = float(num_str)
            color = "#28a745" if f >= 0 else "#dc3545"
            formatted = f"{int(f)}% YoY" if f == int(f) else f"{round(f, 1)}% YoY"
            return (formatted, color)
        except:
            return (s if s else "", "#666")

    # Load data from Dashboard sheet
    try:
        full_df = pd.read_excel('data/key_matrix_data.xlsx', sheet_name='Dashboard')
    except Exception as e:
        return html.Div([
            dbc.Alert(f"Error loading data: {str(e)}", color="danger"),
            html.P("Please ensure key_matrix_data.xlsx exists with Dashboard sheet.")
        ])

    # ============================================================================
    # FILTERING LOGIC (same as before - Investment Strategist is INDEPENDENT)
    # ============================================================================
    if selected_strategist != 'All Strategists':
        filtered_df = full_df[full_df['Investment_Strategist'] == selected_strategist].copy()
        filter_description = f"Investment Strategist: {selected_strategist}"

        if filtered_df.empty:
            data_row = pd.Series()
        else:
            if len(filtered_df) > 1:
                numeric_cols = filtered_df.select_dtypes(include=[np.number]).columns.tolist()
                if numeric_cols:
                    data_row = filtered_df.sort_values(['Year', 'Quarter'], ascending=False).iloc[0]
                    non_numeric_cols = filtered_df.select_dtypes(exclude=[np.number]).columns.tolist()
                    for col in non_numeric_cols:
                        data_row[col] = filtered_df[col].iloc[0]
                else:
                    data_row = filtered_df.iloc[0]
            else:
                data_row = filtered_df.iloc[0]

    elif selected_region == 'All Regions':
        # Get the latest quarter "All Regions" aggregated row
        all_regions_data = full_df[
            (full_df['Region'] == 'All Regions') &
            (full_df['Banker'] == 'All Bankers')
        ]
        if not all_regions_data.empty:
            # Get the most recent quarter
            data_row = all_regions_data.sort_values(['Year', 'Quarter'], ascending=False).iloc[0]
            filter_description = "All Regions | All Bankers | All Strategists"
        else:
            data_row = pd.Series()
            filter_description = "No data available"
    else:
        filtered_df = full_df[full_df['Region'].notna()].copy()

        if selected_region != 'All Regions':
            filtered_df = filtered_df[filtered_df['Region'] == selected_region]

        if selected_banker != 'All Bankers':
            filtered_df = filtered_df[filtered_df['Banker'] == selected_banker]

        filter_description = f"{selected_region} | {selected_banker}"

        if filtered_df.empty:
            data_row = pd.Series()
        else:
            # Get the latest quarter data
            data_row = filtered_df.sort_values(['Year', 'Quarter'], ascending=False).iloc[0]

    # ============================================================================
    # EXTRACT VALUES FROM EXCEL - ALL COLUMNS
    # ============================================================================
    families = as_int(pick_from_row(data_row, ['Families'], 0))
    families_yoy_text, families_yoy_color = format_pct(pick_from_row(data_row, ['Families YoY'], "0"))

    aum_aua = as_int(pick_from_row(data_row, ['AUM_Cr'], 0))
    aum_aua_yoy_text, aum_aua_yoy_color = format_pct(pick_from_row(data_row, ['YoY(AUM+AUA)'], "0"))

    # Calculate AUM per family
    aum_per_family = int(aum_aua / families) if families > 0 else 0

    # NEW: Extract HC Prop and HC TP from Excel
    hc_prop = as_int(pick_from_row(data_row, ['HC_Prop'], 0))
    hc_prop_pct = as_int(pick_from_row(data_row, ['HC_Prop_AUM_Pct'], 0))

    hc_tp = as_int(pick_from_row(data_row, ['HC_TP'], 0))
    hc_tp_pct = as_int(pick_from_row(data_row, ['HC_TP_AUM_Pct'], 0))

    total_revenue = as_float(pick_from_row(data_row, ['Total_Revenue_Cr'], 0))
    total_revenue_yoy_text, total_revenue_yoy_color = format_pct(
        pick_from_row(data_row, ['Total Revenue YoY'], "0"))

    billed_revenue = as_float(pick_from_row(data_row, ['Billed_Revenue_Cr'], 0))
    billed_revenue_yoy_text, billed_revenue_yoy_color = format_pct(
        pick_from_row(data_row, ['Billed Revenue YoY'], "0"))

    # NEW: Advisory metrics from Excel
    advisory_asset_pct = as_int(pick_from_row(data_row, ['Advisory_Assets_Pct'], 0))
    advisory_revenue_pct = as_int(pick_from_row(data_row, ['Advisory_Revenue_Pct'], 0))

    # NEW: Qualified families from Excel (note: column name has >5cr in it)
    qualified_families = as_int(pick_from_row(data_row, ['Qualified_Families>5cr'], 0))
    qualified_families_assets = as_int(pick_from_row(data_row, ['Qualified_Families_Assets'], 0))

    # NEW: Inhouse strategy from Excel
    inhouse_strategy_pct = as_int(pick_from_row(data_row, ['Inhouse_Strategy_Pct'], 0))

    # Styling - More compact and professional
    kpi_card_style = {
        "border": "none",
        "padding": "15px",
        "backgroundColor": "white",
        "height": "150px",
        "display": "flex",
        "flexDirection": "column",
        "justifyContent": "space-between",
        "alignItems": "center",
        "textAlign": "center",
        "borderRadius": "10px",
        "boxShadow": "0 2px 4px rgba(0,0,0,0.04), 0 1px 2px rgba(0,0,0,0.08)",
        "position": "relative",
        "overflow": "hidden",
        "transition": "all 0.3s ease"
    }

    main_value_style = {
        "fontSize": "32px",
        "fontWeight": "bold",
        "color": "#C41E3A",
        "lineHeight": "1",
        "margin": "0",
        "paddingTop": "5px"
    }

    yoy_style_base = {
        "fontSize": "12px",
        "margin": "5px 0",
        "fontWeight": "500",
        "minHeight": "20px"
    }

    label_style = {
        "fontSize": "11px",
        "color": "#333",
        "fontWeight": "normal",
        "borderTop": "1px solid #e0e0e0",
        "padding": "8px 5px 0 5px",
        "margin": "0",
        "width": "100%",
        "marginTop": "auto"
    }

    # Build the KPI layout
    return html.Div([
        # Display current selection
        html.Div([
            html.H6([
                html.Span("Viewing: ", style={"color": "#666", "fontWeight": "normal"}),
                html.Span(filter_description, style={"color": "#C41E3A", "fontWeight": "bold"})
            ], style={"marginBottom": "12px", "fontSize": "14px"})
        ]),

        # KPI Grid - All 7 boxes with REAL DATA
        html.Div([
            html.Div([
                html.Div(str(families), style=main_value_style),
                html.Div(families_yoy_text, style={**yoy_style_base, "color": families_yoy_color}),
                html.Div("Families", style=yoy_style_base)
            ], style={**kpi_card_style, "flex": "1", "minWidth": "110px", "margin": "0 4px"}),

            html.Div([
                html.Div(str(aum_aua), style=main_value_style),
                html.Div(aum_aua_yoy_text, style={**yoy_style_base, "color": aum_aua_yoy_color}),
                html.Div("AUM/AUA(Cr)", style=yoy_style_base)
            ], style={**kpi_card_style, "flex": "1", "minWidth": "110px", "margin": "0 4px"}),

            html.Div([
                html.Div(str(aum_per_family), style=main_value_style),
                html.Div("AUM/Family(Cr)", style=yoy_style_base)
            ], style={**kpi_card_style, "flex": "1", "minWidth": "110px", "margin": "0 4px"}),

            html.Div([
                html.Div(str(hc_prop), style=main_value_style),
                html.Div(f"{hc_prop_pct}% of AUM(Cr)",
                         style={**yoy_style_base, "fontSize": "11px", "color": total_revenue_yoy_color}),
                html.Div("HC Prop", style=yoy_style_base)
            ], style={**kpi_card_style, "flex": "1", "minWidth": "110px", "margin": "0 4px"}),

            html.Div([
                html.Div(str(hc_tp), style=main_value_style),
                html.Div(f"{hc_tp_pct}% of AUM(Cr)",
                         style={**yoy_style_base, "fontSize": "11px", "color": total_revenue_yoy_color}),
                html.Div("HC TP", style=yoy_style_base)
            ], style={**kpi_card_style, "flex": "1", "minWidth": "110px", "margin": "0 4px"}),

            html.Div([
                html.Div(str(total_revenue), style=main_value_style),
                html.Div(total_revenue_yoy_text, style={**yoy_style_base, "color": total_revenue_yoy_color}),
                html.Div("Total Revenue(Cr)", style=yoy_style_base)
            ], style={**kpi_card_style, "flex": "1", "minWidth": "110px", "margin": "0 4px"}),

            html.Div([
                html.Div(str(billed_revenue), style=main_value_style),
                html.Div(billed_revenue_yoy_text, style={**yoy_style_base, "color": billed_revenue_yoy_color}),
                html.Div("Billed Revenue(Cr)", style=yoy_style_base)
            ], style={**kpi_card_style, "flex": "1", "minWidth": "110px", "margin": "0 4px"})
        ], style={
            "display": "flex",
            "justifyContent": "space-between",
            "marginBottom": "20px",
            "gap": "4px",
            "flexWrap": "wrap"
        }),

        html.Hr(style={"margin": "20px 0", "borderColor": "#e0e0e0"}),

        html.Div([
            html.P("AUM / Total Revenue : Advisory Family assets across platform / products.",
                   style={"fontSize": "13px", "color": "#666", "margin": "15px 0", "paddingLeft": "15px"})
        ]),

        # Advisory Assets & Revenue - FROM EXCEL
        html.Div([
            html.P([
                html.Span(f"{advisory_asset_pct}%",
                          style={"color": "#C41E3A", "fontWeight": "bold", "fontSize": "16px"}),
                html.Span(" of Total Assets is advisory client assets. ", style={"fontSize": "15px"}),
                html.Span(f"{advisory_revenue_pct}%",
                          style={"color": "#C41E3A", "fontWeight": "bold", "fontSize": "16px"}),
                html.Span(" of total revenue earned in FY23 is from advisory clients.", style={"fontSize": "15px"})
            ], style={"margin": "15px 0 15px 15px", "lineHeight": "1.5"})
        ]),

        # Qualified Families - FROM EXCEL
        html.Div([
            html.P([
                html.Span("Qualified families (Asset > 5 Cr) : ", style={"fontSize": "15px"}),
                html.Span(str(qualified_families),
                          style={"color": "#C41E3A", "fontWeight": "bold", "fontSize": "16px"}),
                html.Span(" families having assets under advise / management of ", style={"fontSize": "15px"}),
                html.Span(str(qualified_families_assets),
                          style={"color": "#C41E3A", "fontWeight": "bold", "fontSize": "16px"}),
                html.Span(" Cr.", style={"fontSize": "15px"})
            ], style={"margin": "15px 0 15px 15px", "lineHeight": "1.5"})
        ]),

        # Inhouse Strategy - FROM EXCEL
        html.Div([
            html.P([
                html.Span(f"{inhouse_strategy_pct}%",
                          style={"color": "#C41E3A", "fontWeight": "bold", "fontSize": "16px"}),
                html.Span(" of total Advisory assets are advised / managed as per inhouse strategies (Prop products).",
                          style={"fontSize": "15px"})
            ], style={"margin": "15px 0 25px 15px", "lineHeight": "1.5"})
        ]),

        dbc.Row([
            dbc.Col([
                html.P("www.trustplutus.com", style={"fontSize": "11px", "color": "#999", "margin": "0"})
            ], width=6),
            dbc.Col([
                html.P("Private and Confidential | 2",
                       style={"fontSize": "11px", "color": "#999", "margin": "0", "textAlign": "right"})
            ], width=6)
        ], style={"marginTop": "25px", "paddingTop": "8px", "borderTop": "1px solid #eee"})
    ])


@app.callback(
    Output('km-banker-filter', 'options'),
    Output('km-banker-filter', 'value'),
    Input('km-region-filter', 'value'),
    State('km-region-banker-map', 'data')
)
def update_km_banker_dropdown(selected_region, region_banker_map):
    if not selected_region or not region_banker_map:
        return [{'label': 'All Bankers', 'value': 'All Bankers'}], 'All Bankers'
    bankers = region_banker_map.get(selected_region, ['All Bankers'])
    return [{'label': b, 'value': b} for b in bankers], 'All Bankers'


# 2. Update comparison options for Key Matrix
@app.callback(
    Output('km-compare-values', 'options'),
    Input('km-compare-type', 'value'),
    State('km-region-filter', 'value')
)
def update_km_comparison_options(compare_type, selected_region):
    if not compare_type:
        return []

    df = load_key_matrix_data()
    if df.empty:
        return []

    if compare_type == 'region':
        options = df['Region'].unique().tolist()
    elif compare_type == 'banker':
        if selected_region != 'All Regions':
            options = df[df['Region'] == selected_region]['Banker'].unique().tolist()
        else:
            options = df['Banker'].unique().tolist()
    elif compare_type == 'strategist':
        options = df['Investment_Strategist'].unique().tolist()
    else:
        return []

    options = [str(opt) for opt in options if pd.notna(opt) and str(opt) != 'nan' and not str(opt).startswith('All')]
    return [{'label': opt, 'value': opt} for opt in sorted(options)]


# 3. Main Key Matrix content update
# Helper function for percentage color
def get_percentage_color(value):
    """Return color based on percentage value"""
    try:
        # Handle different formats
        if isinstance(value, str):
            val_str = str(value).replace('%', '').replace('YoY', '').replace('+', '').strip()
        else:
            val_str = str(value)

        val_float = float(val_str)

        if val_float > 0:
            return "#28a745"  # Green
        elif val_float < 0:
            return "#dc3545"  # Red
        else:
            return "#666"  # Gray for zero
    except:
        return "#666"


# Updated Key Matrix Dashboard Callback - Uses Excel Data
# Replace the existing update_km_content callback with this

# ============================================================================
# COMPLETE KEY MATRIX DASHBOARD UPDATE - USES BOTH EXCEL SHEETS
# ============================================================================
# Replace the entire update_km_content callback with this version

# ============================================================================
# FINAL KEY MATRIX DASHBOARD - Complete Integration
# ============================================================================
# Features:
# 1. Left sidebar uses "Key Matrix Data" sheet (Region/Banker/Strategist specific)
# 2. Charts use "Dashboard" sheet with FY25-Q1, FY25-Q2... format
# 3. Fixed hover label positioning
# 4. Supports 6 quarters: Q1-Q4 2025, Q1-Q2 2026
# ============================================================================

@app.callback(
    Output('km-content', 'children'),
    Input('km-region-filter', 'value'),
    Input('km-banker-filter', 'value'),
    Input('km-strategist-filter', 'value'),
    Input('km-compare-type', 'value'),
    Input('km-compare-values', 'value'),
    Input('km-tabs', 'active_tab')
)
def update_km_content(region, banker, strategist, compare_type, compare_values, active_tab):
    """Key Matrix Dashboard - All Tabs with Real Data - 6 QUARTERS"""

    try:
        dashboard_df = pd.read_excel('data/key_matrix_data.xlsx', sheet_name='Dashboard')
        key_matrix_df = pd.read_excel('data/key_matrix_data.xlsx', sheet_name='Key Matrix Data')
    except Exception as e:
        return dbc.Alert(f"Error loading data: {str(e)}", color="danger")

    # ============================================================================
    # HANDLE TAB SWITCHING
    # ============================================================================
    if active_tab == "km-tab-trends":
        return create_km_trend_view_v2(dashboard_df, key_matrix_df, region, banker, strategist)

    elif active_tab == "km-tab-table":
        return create_km_table_view_v2(dashboard_df, key_matrix_df, region, banker, strategist)

    # ============================================================================
    # COMPARISON MODE (for Dashboard tab)
    # ============================================================================
    if active_tab == "km-tab-dashboard" and compare_type and compare_values and len(compare_values) >= 2:
        comparison_results = []
        for value in compare_values:
            if compare_type == 'region':
                filtered = dashboard_df[dashboard_df['Region'] == value]
            elif compare_type == 'banker':
                filtered = dashboard_df[dashboard_df['Banker'] == value]
            elif compare_type == 'strategist':
                filtered = dashboard_df[dashboard_df['Investment_Strategist'] == value]
            else:
                continue

            if not filtered.empty:
                latest = filtered.sort_values(['Year', 'Quarter'], ascending=False).iloc[0]
                comparison_results.append({
                    'Name': value,
                    'Families': int(latest['Families']),
                    'AUM_Cr': int(latest['AUM_Cr']),
                    'Billed_Revenue_Cr': round(float(latest['Billed_Revenue_Cr']), 1),
                    'Total_Revenue_Cr': round(float(latest['Total_Revenue_Cr']), 1)
                })

        comparison_df = pd.DataFrame(comparison_results)

        if comparison_df.empty:
            return dbc.Alert("No data found for selected comparison", color="warning")

        num_items = len(compare_values)
        gray_shades = ['#2d3436', '#636e72', '#95a5a6', '#b2bec3', '#dfe6e9'][:num_items]

        def create_gray_comparison_chart(data, metric, title):
            fig = go.Figure()
            fig.add_trace(go.Bar(
                x=data['Name'],
                y=data[metric],
                marker_color=gray_shades,
                text=data[metric].round(1) if metric != 'Families' else data[metric],
                textposition='inside',
                textfont={'size': 11,'color': '#fff'},

            ))
            fig.update_layout(
                title={'text': title, 'font': {'size': 16, 'color': '#C41E3A'}},
                plot_bgcolor='white',
                paper_bgcolor='white',
                font=dict(family='Inter', size=11),
                height=350,
                showlegend=False,
                xaxis=dict(showgrid=False, tickfont={'size': 11}),
                yaxis=dict(showgrid=True, gridcolor='#f0f0f0', tickfont={'size': 10}),
                margin={'l': 50, 'r': 50, 't': 50, 'b': 50},
                hoverlabel=dict(
                    bgcolor="white",
                    font_size=12,
                    font_family="Inter",
                    align="center"
                ),
                hovermode='closest'
            )
            return fig

        return html.Div([
            html.H5(f"Comparing {compare_type.title()}s", className="mb-3",
                    style={"color": "#C41E3A", "fontWeight": "bold", "fontSize": "18px"}),

            dbc.Row([
                dbc.Col([dcc.Graph(
                    figure=create_gray_comparison_chart(comparison_df, 'Families', 'Families Comparison'),
                    config={'displayModeBar': False})], width=6),
                dbc.Col([dcc.Graph(figure=create_gray_comparison_chart(comparison_df, 'AUM_Cr', 'AUM (Cr) Comparison'),
                                   config={'displayModeBar': False})], width=6)
            ], className="mb-3"),

            dbc.Row([
                dbc.Col([dcc.Graph(
                    figure=create_gray_comparison_chart(comparison_df, 'Billed_Revenue_Cr', 'Billed Revenue (Cr)'),
                    config={'displayModeBar': False})], width=6),
                dbc.Col([dcc.Graph(
                    figure=create_gray_comparison_chart(comparison_df, 'Total_Revenue_Cr', 'Total Revenue (Cr)'),
                    config={'displayModeBar': False})], width=6)
            ])
        ])

    # ============================================================================
    # NORMAL DASHBOARD VIEW WITH 4 CHARTS - 6 QUARTERS
    # ============================================================================

    # Filter data based on hierarchy
    if strategist != 'All Strategists':
        filtered_df = dashboard_df[dashboard_df['Investment_Strategist'] == strategist]
        km_filtered = key_matrix_df[key_matrix_df['Investment Strategist'] == strategist]
        filter_desc = f"Investment Strategist: {strategist}"
    elif region == 'All Regions':
        filtered_df = dashboard_df[
            (dashboard_df['Region'] == 'All Regions') & (dashboard_df['Banker'] == 'All Bankers')]
        km_filtered = key_matrix_df[(key_matrix_df['Region'].isin(['All Region', 'All Regions']))]
        filter_desc = "All Regions | All Bankers"
    else:
        filtered_df = dashboard_df[dashboard_df['Region'] == region]
        km_filtered = key_matrix_df[key_matrix_df['Region'] == region]
        if banker != 'All Bankers':
            filtered_df = filtered_df[filtered_df['Banker'] == banker]
            km_filtered = km_filtered[km_filtered['Banker'] == banker]
        filter_desc = f"{region} | {banker}"

    if filtered_df.empty or km_filtered.empty:
        return dbc.Alert("No data available for selected filters", color="info")

    # Get latest quarter from filtered data
    latest_quarter = filtered_df.sort_values(['Year', 'Quarter'], ascending=False).iloc[0]
    km_row = km_filtered.iloc[0]

    # Extract values from Key Matrix Data sheet (for left sidebar)
    total_families = int(km_row['total_families'])
    families_growth = int(km_row['families_growth'])
    total_aum = int(km_row['total_aum'])
    aum_growth = int(km_row['aum_growth'])
    core_aum = int(km_row['core_aum'])
    core_growth = int(km_row['core_growth'])
    xsell_aum = int(km_row['xsell_aum'])
    xsell_growth = int(km_row['xsell_growth'])
    billed_revenue = round(float(km_row['billed_revenue']), 1)
    billed_growth = int(km_row['billed_growth'])
    total_revenue = round(float(km_row['total_revenue']), 1)
    revenue_growth = int(km_row['revenue_growth'])

    # ============================================================================
    # GET ALL 6 QUARTERS DATA (2025 Q1-Q4, 2026 Q1-Q2)
    # ============================================================================

    # Sort by Year and Quarter to get chronological order
    sorted_df = filtered_df.sort_values(['Year', 'Quarter'])

    # Create quarter labels and extract values for ALL available quarters
    quarter_labels = []
    families_values = []
    aum_values = []
    billed_values = []
    total_rev_values = []

    for _, row in sorted_df.iterrows():
        year_short = str(row['Year'])[-2:]  # Get last 2 digits: 2025 -> 25
        quarter_label = f"FY{year_short}-{row['Quarter']}"

        quarter_labels.append(quarter_label)
        families_values.append(int(row['Families']))
        aum_values.append(int(row['AUM_Cr']))
        billed_values.append(round(float(row['Billed_Revenue_Cr']), 1))
        total_rev_values.append(round(float(row['Total_Revenue_Cr']), 1))

    # Limit to 6 quarters if more exist
    quarter_labels = quarter_labels[:6]
    families_values = families_values[:6]
    aum_values = aum_values[:6]
    billed_values = billed_values[:6]
    total_rev_values = total_rev_values[:6]

    # Generate gray shades for 6 quarters (lightest to darkest)
    gray_colors = ['#dfe6e9', '#b2bec3', '#95a5a6', '#7f8c8d', '#636e72', '#2d3436'][:len(quarter_labels)]

    # ============================================================================
    # LAYOUT: LEFT SIDEBAR + 4 CHARTS (6 QUARTERS EACH)
    # ============================================================================

    return html.Div([
        dbc.Row([
            # ========================================================================
            # LEFT SIDEBAR
            # ========================================================================
            dbc.Col([
                html.Div([
                    html.Small("Viewing:", style={"fontSize": "11px", "color": "#999"}),
                    html.Div(filter_desc, style={
                        "fontSize": "11px",
                        "color": "#C41E3A",
                        "fontWeight": "600",
                        "marginBottom": "20px"
                    })
                ]),

                # Families Section
                html.Div([
                    html.H3("Families", style={
                        "color": "#C41E3A",
                        "fontWeight": "bold",
                        "fontSize": "22px",
                        "margin": "0 0 8px 0"
                    }),
                    html.Div(str(total_families), style={
                        "fontSize": "32px",
                        "color": "#666",
                        "fontWeight": "bold",
                        "lineHeight": "1"
                    }),
                    html.Div(f"( {'+' if families_growth >= 0 else ''}{families_growth}%)", style={
                        "fontSize": "13px",
                        "color": get_percentage_color(families_growth),
                        "marginBottom": "15px",
                        "fontWeight": "600"
                    })
                ], style={
                    "border": "2px dashed #ddd",
                    "padding": "16px",
                    "marginBottom": "20px",
                    "borderRadius": "10px"
                }),

                # AUM Section
                html.Div([
                    html.H3("AUM (Cr)", style={
                        "color": "#C41E3A",
                        "fontWeight": "bold",
                        "fontSize": "22px",
                        "margin": "0 0 8px 0"
                    }),
                    html.Div([
                        html.Span("Total : ", style={"fontSize": "13px", "color": "#666"}),
                        html.Span(f"{total_aum}", style={"fontSize": "13px", "color": "#666"}),
                        html.Span(f"  ( {'+' if aum_growth >= 0 else ''}{aum_growth}%)", style={
                            "fontSize": "13px",
                            "color": get_percentage_color(aum_growth),
                            "fontWeight": "600"
                        })
                    ], style={"marginBottom": "4px"}),
                    html.Div([
                        html.Span("Core : ", style={"fontSize": "13px", "color": "#666"}),
                        html.Span(f"{core_aum}", style={"fontSize": "13px", "color": "#666"}),
                        html.Span(f" ({'+' if core_growth >= 0 else ''}{core_growth}%)", style={
                            "fontSize": "13px",
                            "color": get_percentage_color(core_growth),
                            "fontWeight": "600"
                        })
                    ], style={"marginBottom": "4px"}),
                    html.Div([
                        html.Span("X-Sell : ", style={"fontSize": "13px", "color": "#666"}),
                        html.Span(f"{xsell_aum}", style={"fontSize": "13px", "color": "#666"}),
                        html.Span(f" ({'+' if xsell_growth >= 0 else ''}{xsell_growth}%)", style={
                            "fontSize": "13px",
                            "color": get_percentage_color(xsell_growth),
                            "fontWeight": "600"
                        })
                    ], style={"marginBottom": "4px"})
                ], style={
                    "border": "2px dashed #ddd",
                    "padding": "16px",
                    "marginBottom": "20px",
                    "borderRadius": "10px"
                }),

                # Revenue Section
                html.Div([
                    html.H3("Revenue (Cr)", style={
                        "color": "#C41E3A",
                        "fontWeight": "bold",
                        "fontSize": "22px",
                        "margin": "0 0 8px 0"
                    }),
                    html.Div([
                        html.Span("Total : ", style={"fontSize": "13px", "color": "#666"}),
                        html.Span(f"{total_revenue}", style={"fontSize": "13px", "color": "#666"}),
                        html.Span(f"  ( {'+' if revenue_growth >= 0 else ''}{revenue_growth}%)", style={
                            "fontSize": "13px",
                            "color": get_percentage_color(revenue_growth),
                            "fontWeight": "600"
                        })
                    ], style={"marginBottom": "4px"}),
                    html.Div([
                        html.Span("Billed : ", style={"fontSize": "13px", "color": "#666"}),
                        html.Span(f"{billed_revenue}", style={"fontSize": "13px", "color": "#666"}),
                        html.Span(f" ({'+' if billed_growth >= 0 else ''}{billed_growth}%)", style={
                            "fontSize": "13px",
                            "color": get_percentage_color(billed_growth),
                            "fontWeight": "600"
                        })
                    ], style={"marginBottom": "4px"})
                ], style={
                    "border": "2px dashed #ddd",
                    "padding": "16px",
                    "marginBottom": "20px",
                    "borderRadius": "10px"
                }),

                html.Div("www.trustplutus.com", style={
                    "fontSize": "11px",
                    "color": "#999",
                    "marginTop": "40px"
                })
            ], width=3),

            # ========================================================================
            # RIGHT SIDE - FOUR CHARTS (6 QUARTERS EACH)
            # ========================================================================
            dbc.Col([
                dbc.Row([
                    # Top Left - Families Chart (6 QUARTERS)
                    dbc.Col([
                        html.Div([
                            html.H4("No. of Families", style={
                                "fontSize": "16px",
                                "fontWeight": "bold",
                                "color": "#333",
                                "textAlign": "center",
                                "marginBottom": "12px"
                            }),

                            html.Div([
                                html.Span(f"{abs(families_growth)}%", style={
                                    "fontSize": "22px",
                                    "fontWeight": "bold",
                                    "color": get_percentage_color(families_growth)
                                }),
                                html.Span("→", style={
                                    "fontSize": "18px",
                                    "color": "#666",
                                    "marginLeft": "8px"
                                }),
                                html.Span(str(total_families), style={
                                    "fontSize": "26px",
                                    "fontWeight": "bold",
                                    "color": "#333",
                                    "marginLeft": "12px"
                                })
                            ], style={"textAlign": "right", "marginBottom": "15px"}),

                            dcc.Graph(
                                figure={
                                    'data': [{
                                        'x': quarter_labels,
                                        'y': families_values,
                                        'type': 'bar',
                                        'marker': {'color': gray_colors},
                                        'text': families_values,
                                        'textposition': 'inside',
                                        'textfont': {'size': 11, 'color': '#fff'},
                                        'hovertemplate': '<b>%{x}</b><br>Families: %{y}<extra></extra>',
                                    }],
                                    'layout': {
                                        'height': 220,
                                        'margin': {'l': 35, 'r': 35, 't': 15, 'b': 45},
                                        'xaxis': {
                                            'showgrid': False,
                                            'showline': True,
                                            'linecolor': '#ddd',
                                            'tickangle': -45,
                                            'tickfont': {'size': 10}
                                        },
                                        'yaxis': {'showgrid': True, 'gridcolor': '#f0f0f0'},
                                        'plot_bgcolor': 'white',
                                        'paper_bgcolor': 'white',
                                        'hovermode': 'closest',
                                        'hoverlabel': {
                                            'bgcolor': 'white',
                                            'font_size': 14,
                                            'font_family': 'Inter',
                                            'font_color': '#333',
                                            'bordercolor': '#C41E3A',
                                            'align': 'left'
                                        }
                                    }
                                },
                                config={'displayModeBar': False}
                            )
                        ], style={
                            "border": "1px solid #ddd",
                            "padding": "16px",
                            "backgroundColor": "white",
                            "height": "320px",
                            "borderRadius": "8px"
                        })
                    ], width=6),
                    # Top Right - AUM Chart (6 QUARTERS)
                    dbc.Col([
                        html.Div([
                            html.H4("AUM. In Cr.", style={
                                "fontSize": "16px",
                                "fontWeight": "bold",
                                "color": "#333",
                                "textAlign": "center",
                                "marginBottom": "12px"
                            }),

                            html.Div([
                                html.Span(f"{abs(aum_growth)}%", style={
                                    "fontSize": "22px",
                                    "fontWeight": "bold",
                                    "color": get_percentage_color(aum_growth)
                                }),
                                html.Span("→", style={
                                    "fontSize": "18px",
                                    "color": "#666",
                                    "marginLeft": "8px"
                                }),
                                html.Span(str(total_aum), style={
                                    "fontSize": "26px",
                                    "fontWeight": "bold",
                                    "color": "#333",
                                    "marginLeft": "12px"
                                })
                            ], style={"textAlign": "right", "marginBottom": "15px"}),

                            dcc.Graph(
                                figure={
                                    'data': [{
                                        'x': quarter_labels,
                                        'y': aum_values,
                                        'type': 'bar',
                                        'marker': {'color': gray_colors},
                                        'text': aum_values,
                                        'textposition': 'inside',
                                        'textfont': {'size': 11, 'color': '#fff'},
                                        'hovertemplate': '<b>%{x}</b><br>AUM: ₹%{y} Cr<extra></extra>',
                                    }],
                                    'layout': {
                                        'height': 220,
                                        'margin': {'l': 35, 'r': 35, 't': 15, 'b': 45},
                                        'xaxis': {
                                            'showgrid': False,
                                            'showline': True,
                                            'linecolor': '#ddd',
                                            'tickangle': -45,
                                            'tickfont': {'size': 10}
                                        },
                                        'yaxis': {'showgrid': True, 'gridcolor': '#f0f0f0'},
                                        'plot_bgcolor': 'white',
                                        'paper_bgcolor': 'white',
                                        'hovermode': 'closest'
                                    }
                                },
                                config={'displayModeBar': False}
                            )
                        ], style={
                            "border": "1px solid #ddd",
                            "padding": "16px",
                            "backgroundColor": "white",
                            "height": "320px",
                            "borderRadius": "8px"
                        })
                    ], width=6)
                ], className="mb-3"),

                dbc.Row([
                    # Bottom Left - Billed Revenue (6 QUARTERS)
                    dbc.Col([
                        html.Div([
                            html.H4("Billed Revenue", style={
                                "fontSize": "16px",
                                "fontWeight": "bold",
                                "color": "#333",
                                "textAlign": "center",
                                "marginBottom": "12px"
                            }),

                            html.Div([
                                html.Span(f"{abs(billed_growth)}%", style={
                                    "fontSize": "22px",
                                    "fontWeight": "bold",
                                    "color": get_percentage_color(billed_growth)
                                }),
                                html.Span("→", style={
                                    "fontSize": "18px",
                                    "color": "#666",
                                    "marginLeft": "8px"
                                }),
                                html.Span(str(billed_revenue), style={
                                    "fontSize": "26px",
                                    "fontWeight": "bold",
                                    "color": "#333",
                                    "marginLeft": "12px"
                                })
                            ], style={"textAlign": "right", "marginBottom": "15px"}),

                            dcc.Graph(
                                figure={
                                    'data': [{
                                        'x': quarter_labels,
                                        'y': billed_values,
                                        'type': 'bar',
                                        'marker': {'color': gray_colors},
                                        'text': billed_values,
                                        'textposition': 'inside',
                                        'textfont': {'size': 11, 'color': '#fff'},
                                        'hovertemplate': '<b>%{x}</b><br>Revenue: ₹%{y} Cr<extra></extra>',
                                    }],
                                    'layout': {
                                        'height': 220,
                                        'margin': {'l': 35, 'r': 35, 't': 15, 'b': 45},
                                        'xaxis': {
                                            'showgrid': False,
                                            'showline': True,
                                            'linecolor': '#ddd',
                                            'tickangle': -45,
                                            'tickfont': {'size': 10}
                                        },
                                        'yaxis': {'showgrid': True, 'gridcolor': '#f0f0f0'},
                                        'plot_bgcolor': 'white',
                                        'paper_bgcolor': 'white',
                                        'hovermode': 'closest'
                                    }
                                },
                                config={'displayModeBar': False}
                            )
                        ], style={
                            "border": "1px solid #ddd",
                            "padding": "16px",
                            "backgroundColor": "white",
                            "height": "320px",
                            "borderRadius": "8px"
                        })
                    ], width=6),

                    # Bottom Right - Total Revenue (6 QUARTERS)
                    dbc.Col([
                        html.Div([
                            html.H4("Total Revenue", style={
                                "fontSize": "16px",
                                "fontWeight": "bold",
                                "color": "#333",
                                "textAlign": "center",
                                "marginBottom": "12px"
                            }),

                            html.Div([
                                html.Span(f"{abs(revenue_growth)}%", style={
                                    "fontSize": "22px",
                                    "fontWeight": "bold",
                                    "color": get_percentage_color(revenue_growth)
                                }),
                                html.Span("→", style={
                                    "fontSize": "18px",
                                    "color": "#666",
                                    "marginLeft": "8px"
                                }),
                                html.Span(str(total_revenue), style={
                                    "fontSize": "26px",
                                    "fontWeight": "bold",
                                    "color": "#333",
                                    "marginLeft": "12px"
                                })
                            ], style={"textAlign": "right", "marginBottom": "15px"}),

                            dcc.Graph(
                                figure={
                                    'data': [{
                                        'x': quarter_labels,
                                        'y': total_rev_values,
                                        'type': 'bar',
                                        'marker': {'color': gray_colors},
                                        'text': total_rev_values,
                                        'textposition': 'inside',
                                        'textfont': {'size': 11, 'color': '#fff'},
                                        'hovertemplate': '<b>%{x}</b><br>Revenue: ₹%{y} Cr<extra></extra>',
                                    }],
                                    'layout': {
                                        'height': 220,
                                        'margin': {'l': 35, 'r': 35, 't': 15, 'b': 45},
                                        'xaxis': {
                                            'showgrid': False,
                                            'showline': True,
                                            'linecolor': '#ddd',
                                            'tickangle': -45,
                                            'tickfont': {'size': 10}
                                        },
                                        'yaxis': {'showgrid': True, 'gridcolor': '#f0f0f0'},
                                        'plot_bgcolor': 'white',
                                        'paper_bgcolor': 'white',
                                        'hovermode': 'closest'
                                    }
                                },
                                config={'displayModeBar': False}
                            )
                        ], style={
                            "border": "1px solid #ddd",
                            "padding": "16px",
                            "backgroundColor": "white",
                            "height": "320px",
                            "borderRadius": "8px"
                        })
                    ], width=6)
                ])
            ], width=9)
        ])
    ])
# X-SELL CALLBACKS
# ============================================================================

@app.callback(
    Output('xs-banker-filter', 'options'),
    Output('xs-banker-filter', 'value'),
    Input('xs-region-filter', 'value'),
    State('xs-region-banker-map', 'data')
)
def update_xs_banker_dropdown(selected_region, region_banker_map):
    if not selected_region or not region_banker_map:
        return [{'label': 'All Bankers', 'value': 'All Bankers'}], 'All Bankers'
    bankers = region_banker_map.get(selected_region, ['All Bankers'])
    return [{'label': b, 'value': b} for b in bankers], 'All Bankers'


@app.callback(
    Output('xs-compare-values', 'options'),
    Input('xs-compare-type', 'value'),
    State('xs-region-filter', 'value')
)
def update_xs_comparison_options(compare_type, selected_region):
    if not compare_type:
        return []

    df = load_crosssell_detailed_data()
    if df.empty:
        return []

    if compare_type == 'region':
        options = df['Region'].unique().tolist()
    elif compare_type == 'banker':
        if selected_region != 'All Regions':
            options = df[df['Region'] == selected_region]['Banker'].unique().tolist()
        else:
            options = df['Banker'].unique().tolist()
    elif compare_type == 'strategist':
        options = df['Investment_Strategist'].unique().tolist()
    else:
        return []

    options = [str(opt) for opt in options if pd.notna(opt) and str(opt) != 'nan' and not str(opt).startswith('All')]
    return [{'label': opt, 'value': opt} for opt in sorted(options)]


@app.callback(
    Output('xs-content', 'children'),
    Input('xs-region-filter', 'value'),
    Input('xs-banker-filter', 'value'),
    Input('xs-strategist-filter', 'value'),
    Input('xs-compare-type', 'value'),
    Input('xs-compare-values', 'value'),
    Input('xs-product-filter', 'value'),
    Input('xs-stage-filter', 'value')
)
def update_xs_content(region, banker, strategist, compare_type, compare_values, product_filter, stage_filter):
    """X-Sell Dashboard - ALL 6 COLUMNS FROM EXCEL WITH COLOR-CODED YoY"""

    try:
        df = pd.read_excel('data/crosssell_data.xlsx')
    except:
        return dbc.Alert("Error loading crosssell_data.xlsx", color="danger")

    # ============================================================================
    # COMPARISON MODE
    # ============================================================================
    if compare_type and compare_values and len(compare_values) >= 2:
        comparison_results = []
        for value in compare_values:
            if compare_type == 'region':
                filtered = df[df['Region'] == value]
            elif compare_type == 'banker':
                filtered = df[df['Banker'] == value]
            elif compare_type == 'strategist':
                filtered = df[df['Investment_Strategist'] == value]
            else:
                continue

            if not filtered.empty:
                for stage in ['Planned', 'Proposed', 'Transacted']:
                    stage_data = filtered[filtered['Stage'] == stage]
                    comparison_results.append({
                        'Name': value,
                        'Stage': stage,
                        'Amount_Cr': stage_data['Amount_Cr'].sum(),
                        'Count': stage_data['Count'].sum()
                    })

        comp_df = pd.DataFrame(comparison_results)

        if comp_df.empty:
            return dbc.Alert("No data found for selected comparison", color="warning")

        fig = go.Figure()
        colors = {'Planned': '#636e72', 'Proposed': '#95a5a6', 'Transacted': '#b2bec3'}

        for stage in ['Planned', 'Proposed', 'Transacted']:
            stage_data = comp_df[comp_df['Stage'] == stage]
            fig.add_trace(go.Bar(
                name=stage,
                x=stage_data['Name'],
                y=stage_data['Amount_Cr'],
                marker_color=colors.get(stage, '#7f8c8d'),
                text=stage_data['Amount_Cr'].round(1),
                textposition='outside',
                textfont={'size': 11}
            ))

        fig.update_layout(
            title={
                'text': f'Cross-Sell Amount Comparison (Cr) - {compare_type.title()}s',
                'font': {'size': 16, 'color': '#C41E3A'}
            },
            barmode='group',
            plot_bgcolor='white',
            paper_bgcolor='white',
            height=450,
            font=dict(family='Inter', size=11),
            margin={'l': 40, 'r': 40, 't': 50, 'b': 40}
        )

        return html.Div([
            html.H5(f"Comparing {compare_type.title()}s", className="mb-3",
                    style={"color": "#C41E3A", "fontWeight": "bold", "fontSize": "18px"}),
            dcc.Graph(figure=fig, config={'displayModeBar': False})
        ])

    # ============================================================================
    # NORMAL VIEW - FILTER DATA
    # ============================================================================
    filtered_df = filter_data_by_hierarchy(df, region, banker, strategist)

    # Apply product filter
    if product_filter and product_filter != 'All Products':
        filtered_df = filtered_df[filtered_df['Product'] == product_filter]

    # Apply stage filter
    if stage_filter and stage_filter != 'All Stages':
        filtered_df = filtered_df[filtered_df['Stage'] == stage_filter]

    if filtered_df.empty:
        return dbc.Alert("No data available for selected filters", color="info")

    # ============================================================================
    # HELPER FUNCTION FOR YoY COLOR
    # ============================================================================
    def get_yoy_color(value):
        """Return green for positive, red for negative, gray for zero"""
        try:
            val = float(value)
            if val > 0:
                return "#28a745"  # Green
            elif val < 0:
                return "#dc3545"  # Red
            else:
                return "#6c757d"  # Gray
        except:
            return "#6c757d"

    # ============================================================================
    # EXTRACT ALL 6 VALUES FROM EXCEL
    # ============================================================================

    # Get the first row for aggregated metrics (assuming one summary row per filter)
    if not filtered_df.empty:
        summary_row = filtered_df.iloc[0]
    else:
        summary_row = pd.Series()

    # Extract Total
    if 'Total' in filtered_df.columns:
        total_value = round(filtered_df['Total'].sum(), 1)
    else:
        total_value = round(filtered_df['Amount_Cr'].sum(), 1)
        print("WARNING: 'Total' column missing, calculating from Amount_Cr")

    # Extract Total YoY Growth
    if 'YoY_Growth' in summary_row:
        total_yoy = summary_row['YoY_Growth']
        try:
            total_yoy = int(total_yoy) if not pd.isna(total_yoy) else 0
        except:
            total_yoy = 0
    else:
        total_yoy = 0
        print("WARNING: 'YoY_Growth' column missing in crosssell_data.xlsx")

    # Extract HC_Prop
    if 'HC_Prop' in filtered_df.columns:
        hc_prop_value = round(filtered_df['HC_Prop'].sum(), 1)
    else:
        hc_prop_value = 0
        print("WARNING: 'HC_Prop' column missing in crosssell_data.xlsx")

    # Extract HC_Prop YoY Growth
    if 'HC_Prop_YoY' in summary_row:
        hc_prop_yoy = summary_row['HC_Prop_YoY']
        try:
            hc_prop_yoy = int(hc_prop_yoy) if not pd.isna(hc_prop_yoy) else 0
        except:
            hc_prop_yoy = 0
    else:
        hc_prop_yoy = 0
        print("WARNING: 'HC_Prop_YoY' column missing in crosssell_data.xlsx")

    # Extract HC_TP
    if 'HC_TP' in filtered_df.columns:
        hc_tp_value = round(filtered_df['HC_TP'].sum(), 1)
    else:
        hc_tp_value = 0
        print("WARNING: 'HC_TP' column missing in crosssell_data.xlsx")

    # Extract HC_TP YoY Growth
    if 'HC_TP_YoY' in summary_row:
        hc_tp_yoy = summary_row['HC_TP_YoY']
        try:
            hc_tp_yoy = int(hc_tp_yoy) if not pd.isna(hc_tp_yoy) else 0
        except:
            hc_tp_yoy = 0
    else:
        hc_tp_yoy = 0
        print("WARNING: 'HC_TP_YoY' column missing in crosssell_data.xlsx")

    # Generate filter description
    if strategist != 'All Strategists':
        filter_desc = f"Investment Strategist: {strategist}"
    elif region == 'All Regions':
        filter_desc = "All Regions | All Bankers"
    else:
        filter_desc = f"{region} | {banker}"

    # Add active filters
    active_filters = []
    if product_filter and product_filter != 'All Products':
        active_filters.append(f"Product: {product_filter}")
    if stage_filter and stage_filter != 'All Stages':
        active_filters.append(f"Stage: {stage_filter}")

    if active_filters:
        filter_desc += " | " + " | ".join(active_filters)

    # ============================================================================
    # CREATE AGGREGATED TABLE DATA
    # ============================================================================
    agg_data = filtered_df.groupby(['Product', 'Category', 'Stage'])['Amount_Cr'].sum().reset_index()
    stage_pivot = agg_data.pivot_table(
        index=['Product', 'Category'],
        columns='Stage',
        values='Amount_Cr',
        fill_value=0
    ).reset_index()

    for stage in ['Planned', 'Proposed', 'Transacted']:
        if stage not in stage_pivot.columns:
            stage_pivot[stage] = 0

    stage_pivot['Planned'] = stage_pivot['Planned'].round(1)
    stage_pivot['Proposed'] = stage_pivot['Proposed'].round(1)
    stage_pivot['Transacted'] = stage_pivot['Transacted'].round(1)

    # ============================================================================
    # TABLE ROWS
    # ============================================================================
    table_rows = []
    for row_idx, row in stage_pivot.iterrows():
        product_name = row['Product']
        category = row['Category']
        planned_val = row['Planned']
        proposed_val = row['Proposed']
        transacted_val = row['Transacted']

        planned_display = str(planned_val) if planned_val > 0 else ""
        proposed_display = str(proposed_val) if proposed_val > 0 else ""
        transacted_display = str(transacted_val) if transacted_val > 0 else ""

        table_rows.append(
            html.Tr([
                html.Td(product_name, style={
                    "fontSize": "12px",
                    "color": "#333",
                    "padding": "12px",
                    "borderRight": "1px solid #e9ecef",
                    "width": "25%"
                }),
                html.Td(category, style={
                    "fontSize": "12px",
                    "color": "#333",
                    "padding": "12px",
                    "borderRight": "1px solid #e9ecef",
                    "width": "20%"
                }),
                html.Td(planned_display, style={
                    "fontSize": "12px",
                    "color": "#333",
                    "padding": "12px",
                    "borderRight": "1px solid #e9ecef",
                    "textAlign": "center",
                    "width": "15%"
                }),
                html.Td(proposed_display, style={
                    "fontSize": "12px",
                    "color": "#333",
                    "padding": "12px",
                    "borderRight": "1px solid #e9ecef",
                    "textAlign": "center",
                    "width": "15%"
                }),
                html.Td(transacted_display, style={
                    "fontSize": "12px",
                    "color": "#333",
                    "padding": "12px",
                    "borderRight": "1px solid #e9ecef",
                    "textAlign": "center",
                    "width": "15%"
                }),
                html.Td("", style={
                    "fontSize": "12px",
                    "padding": "12px",
                    "width": "10%"
                })
            ], style={
                "backgroundColor": "#f8f9fa" if row_idx % 2 == 0 else "white",
                "borderBottom": "1px solid #e9ecef"
            })
        )

    # ============================================================================
    # FINAL LAYOUT - 3 SECTIONS WITH COLOR-CODED YoY
    # ============================================================================
    return html.Div([
        dbc.Row([
            # ========================================================================
            # LEFT SIDEBAR - ALL 6 VALUES FROM EXCEL WITH COLOR-CODED YoY
            # ========================================================================
            dbc.Col([
                html.Div([
                    html.Small("Viewing:", style={"fontSize": "11px", "color": "#999"}),
                    html.Div(filter_desc, style={
                        "fontSize": "11px",
                        "color": "#C41E3A",
                        "fontWeight": "600",
                        "marginBottom": "20px",
                        "lineHeight": "1.4"
                    })
                ]),

                # ================================================================
                # TOTAL SECTION - WITH COLOR-CODED YoY
                # ================================================================
                html.Div([
                    html.H3("Total", style={
                        "color": "#C41E3A",
                        "fontWeight": "bold",
                        "fontSize": "22px",
                        "margin": "0 0 8px 0"
                    }),
                    html.Div(str(total_value), style={
                        "fontSize": "32px",
                        "color": "#636e72",
                        "fontWeight": "bold",
                        "lineHeight": "1",
                        "marginBottom": "5px"
                    }),
                    html.Div([
                        html.Span("( ", style={"fontSize": "13px", "color": "#666"}),
                        html.Span(
                            f"{'+' if total_yoy > 0 else ''}{total_yoy}%",
                            style={
                                "fontSize": "13px",
                                "color": get_yoy_color(total_yoy),
                                "fontWeight": "600"
                            }
                        ),
                        html.Span(" YoY )", style={"fontSize": "13px", "color": "#666"})
                    ])
                ], style={
                    "border": "2px dashed #ddd",
                    "padding": "16px",
                    "marginBottom": "20px",
                    "borderRadius": "10px"
                }),

                # ================================================================
                # HC - PROP SECTION - WITH COLOR-CODED YoY
                # ================================================================
                html.Div([
                    html.H3("HC - Prop", style={
                        "color": "#C41E3A",
                        "fontWeight": "bold",
                        "fontSize": "22px",
                        "margin": "0 0 8px 0"
                    }),
                    html.Div(str(hc_prop_value), style={
                        "fontSize": "32px",
                        "color": "#636e72",
                        "fontWeight": "bold",
                        "lineHeight": "1",
                        "marginBottom": "5px"
                    }),
                    html.Div([
                        html.Span("( ", style={"fontSize": "13px", "color": "#666"}),
                        html.Span(
                            f"{'+' if hc_prop_yoy > 0 else ''}{hc_prop_yoy}%",
                            style={
                                "fontSize": "13px",
                                "color": get_yoy_color(hc_prop_yoy),
                                "fontWeight": "600"
                            }
                        ),
                        html.Span(" YoY )", style={"fontSize": "13px", "color": "#666"})
                    ])
                ], style={
                    "border": "2px dashed #ddd",
                    "padding": "16px",
                    "marginBottom": "20px",
                    "borderRadius": "10px"
                }),

                # ================================================================
                # HC - TP SECTION - WITH COLOR-CODED YoY
                # ================================================================
                html.Div([
                    html.H3("HC - TP", style={
                        "color": "#C41E3A",
                        "fontWeight": "bold",
                        "fontSize": "22px",
                        "margin": "0 0 8px 0"
                    }),
                    html.Div(str(hc_tp_value), style={
                        "fontSize": "32px",
                        "color": "#636e72",
                        "fontWeight": "bold",
                        "lineHeight": "1",
                        "marginBottom": "5px"
                    }),
                    html.Div([
                        html.Span("( ", style={"fontSize": "13px", "color": "#666"}),
                        html.Span(
                            f"{'+' if hc_tp_yoy > 0 else ''}{hc_tp_yoy}%",
                            style={
                                "fontSize": "13px",
                                "color": get_yoy_color(hc_tp_yoy),
                                "fontWeight": "600"
                            }
                        ),
                        html.Span(" YoY )", style={"fontSize": "13px", "color": "#666"})
                    ])
                ], style={
                    "border": "2px dashed #ddd",
                    "padding": "16px",
                    "marginBottom": "20px",
                    "borderRadius": "10px"
                }),

                html.Div("www.trustplutus.com", style={
                    "fontSize": "11px",
                    "color": "#999",
                    "marginTop": "40px"
                })
            ], width=3),

            # ========================================================================
            # RIGHT SIDE - CLEAN TABLE
            # ========================================================================
            dbc.Col([
                html.Table([
                    html.Thead([
                        html.Tr([
                            html.Th("Products", style={
                                "fontWeight": "bold",
                                "fontSize": "13px",
                                "color": "#333",
                                "padding": "12px",
                                "backgroundColor": "#f8f9fa",
                                "borderBottom": "2px solid #ddd",
                                "borderRight": "1px solid #e9ecef",
                                "width": "22%",
                                "textAlign": "left"
                            }),
                            html.Th("Category", style={
                                "fontWeight": "bold",
                                "fontSize": "13px",
                                "color": "#333",
                                "padding": "12px",
                                "backgroundColor": "#f8f9fa",
                                "borderBottom": "2px solid #ddd",
                                "borderRight": "1px solid #e9ecef",
                                "width": "15%",
                                "textAlign": "left"
                            }),
                            html.Th("Planned", style={
                                "fontWeight": "bold",
                                "fontSize": "13px",
                                "color": "#333",
                                "padding": "12px",
                                "backgroundColor": "#e9ecef",
                                "borderBottom": "2px solid #ddd",
                                "borderRight": "1px solid #e9ecef",
                                "width": "15%",
                                "textAlign": "center"
                            }),
                            html.Th("Proposed", style={
                                "fontWeight": "bold",
                                "fontSize": "13px",
                                "color": "#333",
                                "padding": "12px",
                                "backgroundColor": "#e9ecef",
                                "borderBottom": "2px solid #ddd",
                                "borderRight": "1px solid #e9ecef",
                                "width": "15%",
                                "textAlign": "center"
                            }),
                            html.Th("Transacted", style={
                                "fontWeight": "bold",
                                "fontSize": "13px",
                                "color": "#333",
                                "padding": "12px",
                                "backgroundColor": "#e9ecef",
                                "borderBottom": "2px solid #ddd",
                                "borderRight": "1px solid #e9ecef",
                                "width": "13%",
                                "textAlign": "center"
                            }),
                            html.Th("Remarks", style={
                                "fontWeight": "bold",
                                "fontSize": "13px",
                                "color": "#333",
                                "padding": "12px",
                                "backgroundColor": "#f8f9fa",
                                "borderBottom": "2px solid #ddd",
                                "width": "15%",
                                "textAlign": "left"
                            })
                        ])
                    ]),
                    html.Tbody(table_rows)
                ], style={
                    "width": "100%",
                    "border": "1px solid #ddd",
                    "borderCollapse": "collapse",
                    "borderRadius": "5px",
                    "backgroundColor": "white",
                    "tableLayout": "fixed"
                })
            ], width=9)
        ], style={"marginTop": "10px"})
    ])
# ============================================================================
# MANDATES CALLBACKS
# ============================================================================

@app.callback(
    Output('md-banker-filter', 'options'),
    Output('md-banker-filter', 'value'),
    Input('md-region-filter', 'value'),
    State('md-region-banker-map', 'data')
)
def update_md_banker_dropdown(selected_region, region_banker_map):
    if not selected_region or not region_banker_map:
        return [{'label': 'All Bankers', 'value': 'All Bankers'}], 'All Bankers'
    bankers = region_banker_map.get(selected_region, ['All Bankers'])
    return [{'label': b, 'value': b} for b in bankers], 'All Bankers'


@app.callback(
    Output('md-compare-values', 'options'),
    Input('md-compare-type', 'value'),
    State('md-region-filter', 'value')
)
def update_md_comparison_options(compare_type, selected_region):
    if not compare_type:
        return []

    df = load_mandate_detailed_data()
    if df.empty:
        return []

    if compare_type == 'region':
        options = df['Region'].unique().tolist()
    elif compare_type == 'banker':
        if selected_region != 'All Regions':
            options = df[df['Region'] == selected_region]['Banker'].unique().tolist()
        else:
            options = df['Banker'].unique().tolist()
    elif compare_type == 'strategist':
        options = df['Investment_Strategist'].unique().tolist()
    else:
        return []

    options = [str(opt) for opt in options if pd.notna(opt) and str(opt) != 'nan' and not str(opt).startswith('All')]
    return [{'label': opt, 'value': opt} for opt in sorted(options)]


@app.callback(
    Output('md-content', 'children'),
    Input('md-region-filter', 'value'),
    Input('md-banker-filter', 'value'),
    Input('md-strategist-filter', 'value'),
    Input('md-compare-type', 'value'),
    Input('md-compare-values', 'value')
)
def update_md_content(region, banker, strategist, compare_type, compare_values):
    """Mandates with OPTIMIZED bar charts and remarks"""

    try:
        df = pd.read_excel('data/mandate_data.xlsx')
    except:
        return dbc.Alert("No data available", color="warning")

    # COMPARISON MODE
    if compare_type and compare_values and len(compare_values) >= 2:
        comparison_results = []
        for value in compare_values:
            if compare_type == 'region':
                filtered = df[df['Region'] == value]
            elif compare_type == 'banker':
                filtered = df[df['Banker'] == value]
            elif compare_type == 'strategist':
                filtered = df[df['Investment_Strategist'] == value]
            else:
                continue

            if not filtered.empty:
                for stage in ['Pitched', 'Offered', 'Converted', 'Declined']:
                    stage_data = filtered[filtered['Stage'] == stage]
                    comparison_results.append({
                        'Name': value,
                        'Stage': stage,
                        'Count': stage_data['Count'].sum(),
                        'Amount_Cr': stage_data['Amount_Cr'].sum()
                    })

        comp_df = pd.DataFrame(comparison_results)

        if comp_df.empty:
            return dbc.Alert("No data found for selected comparison", color="warning")

        fig = go.Figure()
        colors = {'Pitched': '#C41E3A', 'Offered': '#A0203A', 'Converted': '#28a745', 'Declined': '#6C757D'}
        for stage in ['Pitched', 'Offered', 'Converted', 'Declined']:
            stage_data = comp_df[comp_df['Stage'] == stage]
            fig.add_trace(go.Bar(
                name=stage,
                x=stage_data['Name'],
                y=stage_data['Count'],
                marker_color=colors.get(stage, '#C41E3A'),
                text=stage_data['Count'],
                textposition='outside',
                textfont={'size': 11}
            ))

        fig.update_layout(
            title={
                'text': f'Mandate Pipeline Comparison - {compare_type.title()}s',
                'font': {'size': 16, 'color': '#C41E3A'}
            },
            barmode='group',
            plot_bgcolor='white',
            paper_bgcolor='white',
            height=450,
            font=dict(family='Inter', size=11),
            margin={'l': 40, 'r': 40, 't': 50, 'b': 40}
        )

        return html.Div([
            html.H5(f"Comparing {compare_type.title()}s", className="mb-3",
                    style={"color": "#C41E3A", "fontWeight": "bold", "fontSize": "18px"}),
            dcc.Graph(figure=fig, config={'displayModeBar': False})
        ])

        # NORMAL VIEW
    filtered_df = filter_data_by_hierarchy(df, region, banker, strategist)

    if filtered_df.empty:
        return dbc.Alert("No data available for selected filters", color="info")

    # Calculate stage totals
    stage_totals = filtered_df.groupby('Stage')['Count'].sum()

    total_pitched = int(stage_totals.get('Pitched', 0))
    total_offered = int(stage_totals.get('Offered', 0))
    total_converted = int(stage_totals.get('Converted', 0))
    total_declined = int(stage_totals.get('Declined', 0))

    # Generate filter description
    if strategist != 'All Strategists':
        filter_desc = f"Investment Strategist: {strategist}"
    elif region == 'All Regions':
        filter_desc = "All Regions | All Bankers"
    else:
        filter_desc = f"{region} | {banker}"

    # OPTIMIZED LAYOUT
    return html.Div([
        dbc.Row([
            # Left Sidebar - OPTIMIZED
            dbc.Col([
                # Current filter
                html.Div([
                    html.Small("Viewing:", style={"fontSize": "11px", "color": "#999"}),
                    html.Div(filter_desc, style={
                        "fontSize": "11px",
                        "color": "#C41E3A",
                        "fontWeight": "600",
                        "marginBottom": "20px"
                    })
                ]),

                # Pitched Section - OPTIMIZED
                html.Div([
                    html.H3("Pitched", style={
                        "color": "#C41E3A",
                        "fontWeight": "bold",
                        "fontSize": "22px",
                        "margin": "0 0 8px 0"
                    }),
                    html.Div(str(total_pitched), style={
                        "fontSize": "32px",
                        "color": "#999",
                        "fontWeight": "bold",
                        "lineHeight": "1"
                    }),
                    html.Div("( + 27 %)", style={
                        "fontSize": "13px",
                        "color": "#28a745",
                        "marginBottom": "15px"
                    })
                ], style={
                    "border": "2px dashed #ddd",
                    "padding": "16px",
                    "marginBottom": "20px",
                    "borderRadius": "10px"
                }),

                # Offered Section - OPTIMIZED
                html.Div([
                    html.H3("Offered", style={
                        "color": "#C41E3A",
                        "fontWeight": "bold",
                        "fontSize": "22px",
                        "margin": "0 0 8px 0"
                    }),
                    html.Div(str(total_offered), style={
                        "fontSize": "32px",
                        "color": "#999",
                        "fontWeight": "bold",
                        "lineHeight": "1"
                    })
                ], style={
                    "border": "2px dashed #ddd",
                    "padding": "16px",
                    "marginBottom": "20px",
                    "borderRadius": "10px"
                }),

                # Converted Section - OPTIMIZED
                html.Div([
                    html.H3("Converted", style={
                        "color": "#C41E3A",
                        "fontWeight": "bold",
                        "fontSize": "22px",
                        "margin": "0 0 8px 0"
                    }),
                    html.Div(str(total_converted), style={
                        "fontSize": "32px",
                        "color": "#28a745",
                        "fontWeight": "bold",
                        "lineHeight": "1"
                    })
                ], style={
                    "border": "2px dashed #ddd",
                    "padding": "16px",
                    "marginBottom": "20px",
                    "borderRadius": "10px"
                }),

                # Declined Section - OPTIMIZED
                html.Div([
                    html.H3("Declined", style={
                        "color": "#C41E3A",
                        "fontWeight": "bold",
                        "fontSize": "22px",
                        "margin": "0 0 8px 0"
                    }),
                    html.Div(str(total_declined), style={
                        "fontSize": "32px",
                        "color": "#999",
                        "fontWeight": "bold",
                        "lineHeight": "1"
                    })
                ], style={
                    "border": "2px dashed #ddd",
                    "padding": "16px",
                    "marginBottom": "20px",
                    "borderRadius": "10px"
                }),

                # Website footer
                html.Div("www.trustplutus.com", style={
                    "fontSize": "11px",
                    "color": "#999",
                    "marginTop": "30px"
                })
            ], width=3),

            # Right side - OPTIMIZED BAR CHARTS WITH REMARKS
            dbc.Col([
                html.Div([
                    # Remarks Header - OPTIMIZED
                    html.Div([
                        html.H4("Remarks", style={
                            "fontSize": "16px",
                            "fontWeight": "bold",
                            "color": "#999",
                            "textAlign": "right",
                            "fontStyle": "italic",
                            "marginBottom": "15px"
                        })
                    ]),

                    # Pitched Row - OPTIMIZED
                    html.Div([
                        dbc.Row([
                            dbc.Col([
                                html.H4("Pitched", style={
                                    "fontSize": "16px",
                                    "fontWeight": "normal",
                                    "color": "#999",
                                    "margin": "0",
                                    "paddingTop": "25px"
                                })
                            ], width=2),
                            dbc.Col([
                                # Bar chart for Pitched - OPTIMIZED
                                html.Div([
                                    html.Div([
                                        html.Div(style={
                                            "width": "55px",
                                            "height": "35px",
                                            "backgroundColor": "#666",
                                            "margin": "0 auto 5px auto"
                                        }),
                                        html.Div("Q1", style={
                                            "textAlign": "center",
                                            "fontSize": "11px",
                                            "color": "#999"
                                        })
                                    ], style={"display": "inline-block", "margin": "0 12px"}),

                                    html.Div([
                                        html.Div(style={
                                            "width": "55px",
                                            "height": "48px",
                                            "backgroundColor": "#666",
                                            "margin": "0 auto 5px auto"
                                        }),
                                        html.Div("Q2", style={
                                            "textAlign": "center",
                                            "fontSize": "11px",
                                            "color": "#999"
                                        })
                                    ], style={"display": "inline-block", "margin": "0 12px"}),

                                    html.Div([
                                        html.Div(style={
                                            "width": "55px",
                                            "height": "62px",
                                            "backgroundColor": "#666",
                                            "margin": "0 auto 5px auto"
                                        }),
                                        html.Div("Q3", style={
                                            "textAlign": "center",
                                            "fontSize": "11px",
                                            "color": "#999"
                                        })
                                    ], style={"display": "inline-block", "margin": "0 12px"}),

                                    html.Div([
                                        html.Div(style={
                                            "width": "55px",
                                            "height": "75px",
                                            "backgroundColor": "#666",
                                            "margin": "0 auto 5px auto"
                                        }),
                                        html.Div("Q4", style={
                                            "textAlign": "center",
                                            "fontSize": "11px",
                                            "color": "#999"
                                        })
                                    ], style={"display": "inline-block", "margin": "0 12px"})
                                ], style={
                                    "display": "flex",
                                    "alignItems": "flex-end",
                                    "height": "90px",
                                    "justifyContent": "center"
                                })
                            ], width=8),
                            dbc.Col([
                                dcc.Textarea(
                                    id="pitched-remarks",
                                    placeholder="Enter remarks for Pitched...",
                                    style={
                                        "width": "100%",
                                        "height": "70px",
                                        "border": "1px solid #ddd",
                                        "borderRadius": "5px",
                                        "padding": "8px",
                                        "fontSize": "11px",
                                        "resize": "none"
                                    }
                                )
                            ], width=2)
                        ])
                    ], style={
                        "border": "1px solid #ddd",
                        "padding": "16px",
                        "marginBottom": "0",
                        "backgroundColor": "#f8f9fa"
                    }),

                    # Offered Row - OPTIMIZED
                    html.Div([
                        dbc.Row([
                            dbc.Col([
                                html.H4("Offered", style={
                                    "fontSize": "16px",
                                    "fontWeight": "normal",
                                    "color": "#999",
                                    "margin": "0",
                                    "paddingTop": "25px"
                                })
                            ], width=2),
                            dbc.Col([
                                html.Div([
                                    html.Div([
                                        html.Div(style={
                                            "width": "55px",
                                            "height": "30px",
                                            "backgroundColor": "#666",
                                            "margin": "0 auto 5px auto"
                                        }),
                                        html.Div("Q1", style={
                                            "textAlign": "center",
                                            "fontSize": "11px",
                                            "color": "#999"
                                        })
                                    ], style={"display": "inline-block", "margin": "0 12px"}),

                                    html.Div([
                                        html.Div(style={
                                            "width": "55px",
                                            "height": "43px",
                                            "backgroundColor": "#666",
                                            "margin": "0 auto 5px auto"
                                        }),
                                        html.Div("Q2", style={
                                            "textAlign": "center",
                                            "fontSize": "11px",
                                            "color": "#999"
                                        })
                                    ], style={"display": "inline-block", "margin": "0 12px"}),

                                    html.Div([
                                        html.Div(style={
                                            "width": "55px",
                                            "height": "57px",
                                            "backgroundColor": "#666",
                                            "margin": "0 auto 5px auto"
                                        }),
                                        html.Div("Q3", style={
                                            "textAlign": "center",
                                            "fontSize": "11px",
                                            "color": "#999"
                                        })
                                    ], style={"display": "inline-block", "margin": "0 12px"}),

                                    html.Div([
                                        html.Div(style={
                                            "width": "55px",
                                            "height": "70px",
                                            "backgroundColor": "#666",
                                            "margin": "0 auto 5px auto"
                                        }),
                                        html.Div("Q4", style={
                                            "textAlign": "center",
                                            "fontSize": "11px",
                                            "color": "#999"
                                        })
                                    ], style={"display": "inline-block", "margin": "0 12px"})
                                ], style={
                                    "display": "flex",
                                    "alignItems": "flex-end",
                                    "height": "90px",
                                    "justifyContent": "center"
                                })
                            ], width=8),
                            dbc.Col([
                                dcc.Textarea(
                                    id="offered-remarks",
                                    placeholder="Enter remarks for Offered...",
                                    style={
                                        "width": "100%",
                                        "height": "70px",
                                        "border": "1px solid #ddd",
                                        "borderRadius": "5px",
                                        "padding": "8px",
                                        "fontSize": "11px",
                                        "resize": "none"
                                    }
                                )
                            ], width=2)
                        ])
                    ], style={
                        "border": "1px solid #ddd",
                        "borderTop": "0",
                        "padding": "16px",
                        "marginBottom": "0",
                        "backgroundColor": "white"
                    }),

                    # Converted Row - OPTIMIZED
                    html.Div([
                        dbc.Row([
                            dbc.Col([
                                html.H4("Converted", style={
                                    "fontSize": "16px",
                                    "fontWeight": "normal",
                                    "color": "#999",
                                    "margin": "0",
                                    "paddingTop": "25px"
                                })
                            ], width=2),
                            dbc.Col([
                                html.Div([
                                    html.Div([
                                        html.Div(style={
                                            "width": "55px",
                                            "height": "26px",
                                            "backgroundColor": "#666",
                                            "margin": "0 auto 5px auto"
                                        }),
                                        html.Div("Q1", style={
                                            "textAlign": "center",
                                            "fontSize": "11px",
                                            "color": "#999"
                                        })
                                    ], style={"display": "inline-block", "margin": "0 12px"}),

                                    html.Div([
                                        html.Div(style={
                                            "width": "55px",
                                            "height": "39px",
                                            "backgroundColor": "#666",
                                            "margin": "0 auto 5px auto"
                                        }),
                                        html.Div("Q2", style={
                                            "textAlign": "center",
                                            "fontSize": "11px",
                                            "color": "#999"
                                        })
                                    ], style={"display": "inline-block", "margin": "0 12px"}),

                                    html.Div([
                                        html.Div(style={
                                            "width": "55px",
                                            "height": "52px",
                                            "backgroundColor": "#666",
                                            "margin": "0 auto 5px auto"
                                        }),
                                        html.Div("Q3", style={
                                            "textAlign": "center",
                                            "fontSize": "11px",
                                            "color": "#999"
                                        })
                                    ], style={"display": "inline-block", "margin": "0 12px"}),

                                    html.Div([
                                        html.Div(style={
                                            "width": "55px",
                                            "height": "66px",
                                            "backgroundColor": "#666",
                                            "margin": "0 auto 5px auto"
                                        }),
                                        html.Div("Q4", style={
                                            "textAlign": "center",
                                            "fontSize": "11px",
                                            "color": "#999"
                                        })
                                    ], style={"display": "inline-block", "margin": "0 12px"})
                                ], style={
                                    "display": "flex",
                                    "alignItems": "flex-end",
                                    "height": "90px",
                                    "justifyContent": "center"
                                })
                            ], width=8),
                            dbc.Col([
                                dcc.Textarea(
                                    id="converted-remarks",
                                    placeholder="Enter remarks for Converted...",
                                    style={
                                        "width": "100%",
                                        "height": "70px",
                                        "border": "1px solid #ddd",
                                        "borderRadius": "5px",
                                        "padding": "8px",
                                        "fontSize": "11px",
                                        "resize": "none"
                                    }
                                )
                            ], width=2)
                        ])
                    ], style={
                        "border": "1px solid #ddd",
                        "borderTop": "0",
                        "padding": "16px",
                        "marginBottom": "0",
                        "backgroundColor": "#f8f9fa"
                    }),

                    # Declined Row - OPTIMIZED
                    html.Div([
                        dbc.Row([
                            dbc.Col([
                                html.H4("Declined", style={
                                    "fontSize": "16px",
                                    "fontWeight": "normal",
                                    "color": "#999",
                                    "margin": "0",
                                    "paddingTop": "25px"
                                })
                            ], width=2),
                            dbc.Col([
                                html.Div([
                                    html.Div([
                                        html.Div(style={
                                            "width": "55px",
                                            "height": "22px",
                                            "backgroundColor": "#666",
                                            "margin": "0 auto 5px auto"
                                        }),
                                        html.Div("Q1", style={
                                            "textAlign": "center",
                                            "fontSize": "11px",
                                            "color": "#999"
                                        })
                                    ], style={"display": "inline-block", "margin": "0 12px"}),

                                    html.Div([
                                        html.Div(style={
                                            "width": "55px",
                                            "height": "35px",
                                            "backgroundColor": "#666",
                                            "margin": "0 auto 5px auto"
                                        }),
                                        html.Div("Q2", style={
                                            "textAlign": "center",
                                            "fontSize": "11px",
                                            "color": "#999"
                                        })
                                    ], style={"display": "inline-block", "margin": "0 12px"}),

                                    html.Div([
                                        html.Div(style={
                                            "width": "55px",
                                            "height": "48px",
                                            "backgroundColor": "#666",
                                            "margin": "0 auto 5px auto"
                                        }),
                                        html.Div("Q3", style={
                                            "textAlign": "center",
                                            "fontSize": "11px",
                                            "color": "#999"
                                        })
                                    ], style={"display": "inline-block", "margin": "0 12px"}),

                                    html.Div([
                                        html.Div(style={
                                            "width": "55px",
                                            "height": "62px",
                                            "backgroundColor": "#666",
                                            "margin": "0 auto 5px auto"
                                        }),
                                        html.Div("Q4", style={
                                            "textAlign": "center",
                                            "fontSize": "11px",
                                            "color": "#999"
                                        })
                                    ], style={"display": "inline-block", "margin": "0 12px"})
                                ], style={
                                    "display": "flex",
                                    "alignItems": "flex-end",
                                    "height": "90px",
                                    "justifyContent": "center"
                                })
                            ], width=8),
                            dbc.Col([
                                dcc.Textarea(
                                    id="declined-remarks",
                                    placeholder="Enter remarks for Declined...",
                                    style={
                                        "width": "100%",
                                        "height": "70px",
                                        "border": "1px solid #ddd",
                                        "borderRadius": "5px",
                                        "padding": "8px",
                                        "fontSize": "11px",
                                        "resize": "none"
                                    }
                                )
                            ], width=2)
                        ])
                    ], style={
                        "border": "1px solid #ddd",
                        "borderTop": "0",
                        "padding": "16px",
                        "backgroundColor": "white"
                    }),

                    # Save Button - OPTIMIZED
                    html.Div([
                        dbc.Button("Save Remarks", id="save-mandate-remarks",
                                   className="btn-primary",
                                   style={"marginTop": "15px", "fontSize": "13px", "padding": "8px 20px"})
                    ], style={"textAlign": "center"})
                ])
            ], width=9)
        ], style={"marginTop": "10px"})
    ])


# Enhanced client selection callback
@app.callback(
    [Output("prop-family-name", "value"), Output("prop-banker", "value"),
     Output("prop-date", "value"), Output("prop-amount", "value"),
     Output("prop-risk-profile", "value"), Output("prop-mandate-type", "value"),
     Output("current-status-display", "children"), Output("family-name-helper", "children"),
     Output("is-update-store", "data"), Output("selected-client-store", "data")],
    [Input({"type": "select-client", "index": ALL}, "n_clicks"),
     Input("new-client-btn", "n_clicks"), Input("clear-form-btn", "n_clicks")],
    [State("prop-family-name", "value")]
)
def handle_client_selection(select_clicks, new_clicks, clear_clicks, current_family_name):
    ctx = callback_context

    if not ctx.triggered:
        return "", "", "", "", "", None, "", "", False, None

    button_id = ctx.triggered[0]['prop_id']

    # Clear form
    if "clear-form-btn" in button_id:
        return "", "", "", "", "", None, "", "", False, None

    # New client
    if "new-client-btn" in button_id:
        return "", "", "", "", "", None, "", "Creating new client entry", False, None

    # Select existing client
    if "select-client" in button_id:
        try:
            button_data = json.loads(button_id.split('.')[0])
            family_name = button_data['index']

            client_data = load_client_data(family_name)
            if client_data:
                current_status_display = dbc.Alert([
                    html.H6("Existing Client Selected", className="alert-heading"),
                    html.P([
                        html.Strong("Current Status: "),
                        html.Span(client_data.get('Status', 'Unknown'),
                                  className="badge bg-info ms-2")
                    ]),
                    html.P(
                        f"Last updated: {client_data.get('Timestamp', 'Unknown')[:10] if client_data.get('Timestamp') else 'Unknown'}",
                        className="mb-0 small")
                ], color="info")

                return (
                    client_data.get('Family_Name', ''),
                    client_data.get('Banker', ''),
                    client_data.get('Date', ''),
                    client_data.get('Amount', ''),
                    client_data.get('Risk_Profile', ''),
                    client_data.get('Mandate_Type', ''),
                    current_status_display,
                    "Updating existing client",
                    True,
                    family_name
                )
        except Exception as e:
            print(f"Error in client selection: {e}")

    return "", "", "", "", "", None, "", "", False, None
# Main app layout
app.layout = html.Div([
    dcc.Store(id="session-store", data={"page": "home", "logged_in": False}),  # Make sure this is "home"
    html.Div(id="navbar-content"),
    html.Div(id="page-content", style={"paddingTop": "80px"}),
    html.Div(id="app-messages", children=[
        html.Div(id="registration-message"),
        html.Div(id="login-message"),
        html.Div(id="cross-sell-message"),
        html.Div(id="upload-message"),
        html.Div(id="proposal-message")
    ])
], className="premium-container")
@app.callback(
    Output("navbar-content", "children"),
    Input("session-store", "data")
)
def update_navbar(session_data):
    """Update navbar based on login status"""
    is_logged_in = session_data.get("logged_in", False)
    user_name = current_user.get('Name', '') if is_logged_in else ""
    return get_navigation(is_logged_in, user_name)
# Callbacks
@app.callback(
    Output("page-content", "children"),
    Input("session-store", "data")
)
def display_page(session_data):
    """Display page based on session data"""
    if not session_data:
        return get_home_layout()

    page = session_data.get("page", "home")

    if page == "home":
        return get_home_layout()
    elif page == "register":
        return get_registration_layout()
    elif page == "login":
        return get_login_layout()
    elif page == "main-dashboard":
        return get_main_dashboard_layout()
    elif page == "dashboard":
        return get_dashboard_layout()
    elif page == "key-matrix":
        return get_key_matrix_layout()
    elif page == "x-sell-dash":
        return get_x_sell_dashboard_layout()
    elif page == "new-mandates":
        return get_new_mandates_layout()
    elif page == "cross-sell":
        return get_cross_sell_layout()
    elif page == "proposal":
        return get_proposal_layout()
    else:
        return get_home_layout()

@app.callback(
    [Output("manual-entry-content", "style"), Output("upload-entry-content", "style")],
    [Input("cross-sell-tabs", "active_tab")]
)
def toggle_cross_sell_tabs(active_tab):
    if active_tab == "manual-tab":
        return {"display": "block"}, {"display": "none"}
    else:
        return {"display": "none"}, {"display": "block"}

@app.callback(
    Output("status-form-content", "children"),
    Input("prop-status", "value")
)
def update_status_form(status):
    if not status:
        return html.Div()

    if status in ["Pitched", "Follow up"]:
        return html.Div([
            html.Hr(className="my-4"),
            html.H6(f"{status} Details", className="mb-3 text-secondary"),

            dbc.Row([
                dbc.Label("Summary of Meeting", width=3, className="form-label"),
                dbc.Col([dbc.Textarea(id="status-summary", placeholder="Enter meeting summary",
                                      className="form-control", rows=3)], width=9)
            ], className="mb-3"),

            dbc.Row([
                dbc.Label("Feedback", width=3, className="form-label"),
                dbc.Col([dbc.Textarea(id="status-feedback", placeholder="Enter feedback received",
                                      className="form-control", rows=3)], width=9)
            ], className="mb-3"),

            dbc.Row([
                dbc.Label("Next Step", width=3, className="form-label"),
                dbc.Col([dbc.Input(id="status-next-step", type="text", placeholder="Enter next step",
                                   className="form-control")], width=9)
            ], className="mb-3"),

            dbc.Row([
                dbc.Label("Platform", width=3, className="form-label"),
                dbc.Col([dbc.Input(id="status-platform", type="text", placeholder="Enter platform",
                                   className="form-control")], width=9)
            ], className="mb-3"),

            dbc.Row([
                dbc.Label("Support Required", width=3, className="form-label"),
                dbc.Col([dbc.Textarea(id="status-support", placeholder="Enter support required",
                                      className="form-control", rows=2)], width=9)
            ], className="mb-3"),

            dbc.Row([
                dbc.Label("Follow up Date", width=3, className="form-label"),
                dbc.Col([dbc.Input(id="status-followup-date", type="date", className="form-control")], width=9)
            ], className="mb-3")
        ])

    elif status == "Offer":
        return html.Div([
            html.Hr(className="my-4"),
            html.H6("Offer Details", className="mb-3 text-secondary"),

            dbc.Row([
                dbc.Label("Platform", width=3, className="form-label"),
                dbc.Col([dbc.Input(id="status-platform", type="text", placeholder="Enter platform",
                                   className="form-control")], width=9)
            ], className="mb-3"),

            dbc.Row([
                dbc.Label("Amount", width=3, className="form-label"),
                dbc.Col([dbc.Input(id="status-amount", type="number", placeholder="Enter amount",
                                   className="form-control")], width=9)
            ], className="mb-3"),

            dbc.Row([
                dbc.Label("Pricing", width=3, className="form-label"),
                dbc.Col([dbc.Input(id="status-pricing", type="text", placeholder="Enter pricing details",
                                   className="form-control")], width=9)
            ], className="mb-3"),

            dbc.Row([
                dbc.Label("Date", width=3, className="form-label"),
                dbc.Col([dbc.Input(id="status-date", type="date", className="form-control")], width=9)
            ], className="mb-3"),

            dbc.Row([
                dbc.Label("Next Step", width=3, className="form-label"),
                dbc.Col([dbc.Input(id="status-next-step", type="text", placeholder="Enter next step",
                                   className="form-control")], width=9)
            ], className="mb-3"),

            dbc.Row([
                dbc.Label("Support Required", width=3, className="form-label"),
                dbc.Col([dbc.Textarea(id="status-support", placeholder="Enter support required",
                                      className="form-control", rows=2)], width=9)
            ], className="mb-3")
        ])

    elif status == "Converted":
        return html.Div([
            html.Hr(className="my-4"),
            html.H6("Conversion Details", className="mb-3 text-secondary"),

            dbc.Row([
                dbc.Label("Amount", width=3, className="form-label"),
                dbc.Col([dbc.Input(id="status-amount", type="number", placeholder="Enter amount",
                                   className="form-control")], width=9)
            ], className="mb-3"),

            dbc.Row([
                dbc.Label("Pricing", width=3, className="form-label"),
                dbc.Col([dbc.Input(id="status-pricing", type="text", placeholder="Enter pricing details",
                                   className="form-control")], width=9)
            ], className="mb-3"),

            dbc.Row([
                dbc.Label("Platform", width=3, className="form-label"),
                dbc.Col([dbc.Input(id="status-platform", type="text", placeholder="Enter platform",
                                   className="form-control")], width=9)
            ], className="mb-3"),

            dbc.Row([
                dbc.Label("AUM", width=3, className="form-label"),
                dbc.Col([dbc.Input(id="status-aum", type="number", placeholder="Enter AUM amount",
                                   className="form-control")], width=9)
            ], className="mb-3"),

            dbc.Row([
                dbc.Label("Support Required", width=3, className="form-label"),
                dbc.Col([dbc.Textarea(id="status-support", placeholder="Enter support required",
                                      className="form-control", rows=2)], width=9)
            ], className="mb-3")
        ])

    return html.Div()

@app.callback(
    [Output("otp-section", "style"), Output("registration-message", "children")],
    [Input("send-otp-btn", "n_clicks")],
    [State("reg-name", "value"), State("reg-phone", "value"), State("reg-email", "value"), State("reg-role", "value")]
)
def send_otp(n_clicks, name, phone, email, role):
    if not n_clicks:
        return {"display": "none"}, ""
    if not all([name, phone, email, role]):
        return {"display": "none"}, dbc.Alert("Please fill all fields including role selection!", color="warning")
    if not validate_email_domain(email):
        return {"display": "none"}, dbc.Alert("Email must end with a valid company domain!", color="danger")

    users_df = load_users()
    if not users_df.empty and not users_df[users_df['Email'] == email].empty:
        return {"display": "none"}, dbc.Alert("User with this email already exists!", color="warning")

    otp = generate_otp()
    otp_storage[email] = otp
    send_otp_email(email, otp)
    return {"display": "block"}, dbc.Alert(f"OTP sent to {email}! Check console for demo OTP: {otp}", color="success")


@app.callback(
    [Output("session-store", "data", allow_duplicate=True),
     Output("registration-message", "children", allow_duplicate=True)],
    [Input("register-btn", "n_clicks")],
    [State("reg-name", "value"), State("reg-phone", "value"), State("reg-email", "value"), State("reg-role", "value"),
     State("otp-input", "value"), State("reg-password", "value")],
    prevent_initial_call=True
)
def register_user(n_clicks, name, phone, email, role, otp, password):
    if not n_clicks:
        return {"page": "register", "logged_in": False}, ""
    if not all([name, phone, email, role, otp, password]):
        return {"page": "register", "logged_in": False}, dbc.Alert("Please fill all fields!", color="warning")

    if email in otp_storage and otp_storage[email] == otp:
        user_id = generate_user_id(name, role)
        save_user(name, phone, email, role, user_id, password)
        del otp_storage[email]
        success_alert = dbc.Alert([
            html.H5("Registration Successful!", className="alert-heading"),
            html.P(f"Your User ID is: {user_id}", className="mb-0")
        ], color="success")
        return {"page": "login", "logged_in": False}, success_alert
    else:
        return {"page": "register", "logged_in": False}, dbc.Alert("Invalid OTP!", color="danger")


@app.callback(
    Output("session-store", "data", allow_duplicate=True),
    [Input("nav-login-btn", "n_clicks"), Input("nav-register-btn", "n_clicks")],
    prevent_initial_call=True
)
def handle_main_navigation(nav_login, nav_register):
    ctx = callback_context
    if not ctx.triggered:
        return dash.no_update  # Changed from returning home

    # Check if the trigger actually has a click
    button_id = ctx.triggered[0]['prop_id'].split('.')[0]

    # Only proceed if there was an actual click (not just initialization)
    if nav_login is None and nav_register is None:
        return dash.no_update

    if button_id == "nav-login-btn" and nav_login:
        return {"page": "login", "logged_in": False}
    elif button_id == "nav-register-btn" and nav_register:
        return {"page": "register", "logged_in": False}

    return dash.no_update  # Changed from returning home
@app.callback(
    Output("session-store", "data", allow_duplicate=True),
    [Input("hero-get-started", "n_clicks")],
    prevent_initial_call=True
)
def handle_hero_navigation(hero_start):
    if hero_start:
        return {"page": "register", "logged_in": False}
    return {"page": "home", "logged_in": False}


@app.callback(
    Output("session-store", "data", allow_duplicate=True),
    [Input("go-to-login", "n_clicks")],
    prevent_initial_call=True
)
def go_to_login_page(login_clicks):
    if login_clicks:
        return {"page": "login", "logged_in": False}
    return {"page": "register", "logged_in": False}


@app.callback(
    Output("session-store", "data", allow_duplicate=True),
    [Input("go-to-register", "n_clicks")],
    prevent_initial_call=True
)
def go_to_register_page(register_clicks):
    if register_clicks:
        return {"page": "register", "logged_in": False}
    return {"page": "login", "logged_in": False}


@app.callback(
    [Output("session-store", "data", allow_duplicate=True), Output("login-message", "children")],
    [Input("login-btn", "n_clicks")],
    [State("login-userid", "value"), State("login-password", "value")],
    prevent_initial_call=True
)
def login_user(n_clicks, user_id, password):
    global current_user
    if not n_clicks:
        return {"page": "login", "logged_in": False}, ""
    if not all([user_id, password]):
        return {"page": "login", "logged_in": False}, dbc.Alert("Please enter both User ID and password!",
                                                                color="warning")

    user = authenticate_user(user_id, password)
    if user:
        current_user = user
        return {"page": "main-dashboard", "logged_in": True}, ""
    else:
        return {"page": "login", "logged_in": False}, dbc.Alert("Invalid credentials!", color="danger")


@app.callback(
    Output("session-store", "data", allow_duplicate=True),
    [Input("dashboard-btn", "n_clicks"),
     Input("cross-sell-btn", "n_clicks"),
     Input("proposal-btn", "n_clicks")],
    prevent_initial_call=True
)
def handle_main_dashboard_navigation(dashboard_clicks, cross_sell_clicks, proposal_clicks):
    ctx = callback_context
    if not ctx.triggered:
        return {"page": "main-dashboard", "logged_in": True}

    button_id = ctx.triggered[0]['prop_id'].split('.')[0]

    if button_id == "dashboard-btn":
        return {"page": "dashboard", "logged_in": True}
    elif button_id == "cross-sell-btn":
        return {"page": "cross-sell", "logged_in": True}
    elif button_id == "proposal-btn":
        return {"page": "proposal", "logged_in": True}

    return {"page": "main-dashboard", "logged_in": True}

# 2. Dashboard sub-navigation (only elements that exist on dashboard page)
@app.callback(
    Output("session-store", "data", allow_duplicate=True),
    [Input("key-matrix-btn", "n_clicks"),
     Input("x-sell-dash-btn", "n_clicks"),
     Input("new-mandates-btn", "n_clicks")],
    prevent_initial_call=True
)
def handle_dashboard_sub_navigation(matrix_clicks, xsell_clicks, mandates_clicks):
    ctx = callback_context
    if not ctx.triggered:
        return {"page": "dashboard", "logged_in": True}

    button_id = ctx.triggered[0]['prop_id'].split('.')[0]

    if button_id == "key-matrix-btn":
        return {"page": "key-matrix", "logged_in": True}
    elif button_id == "x-sell-dash-btn":
        return {"page": "x-sell-dash", "logged_in": True}
    elif button_id == "new-mandates-btn":
        return {"page": "new-mandates", "logged_in": True}

    return {"page": "dashboard", "logged_in": True}
# Add this new callback after the one above
@app.callback(
    Output("session-store", "data", allow_duplicate=True),
    Input("back-to-main-dash", "n_clicks"),
    State("session-store", "data"),
    prevent_initial_call=True
)
def handle_dashboard_back_button(n_clicks, session_data):
    if n_clicks:
        return {"page": "main-dashboard", "logged_in": True}
    return session_data
# 3. Separate callback for key-matrix page navigation
@app.callback(
    Output("session-store", "data", allow_duplicate=True),
    [Input("back-to-dashboard-matrix", "n_clicks")],
    [State("session-store", "data")],
    prevent_initial_call=True
)
def handle_key_matrix_navigation(back_clicks, session_data):
    if back_clicks and session_data.get("page") == "key-matrix":
        return {"page": "dashboard", "logged_in": True}
    return session_data

# 4. Separate callback for x-sell dashboard page navigation
@app.callback(
    Output("session-store", "data", allow_duplicate=True),
    [Input("back-to-dashboard-xsell", "n_clicks")],
    [State("session-store", "data")],
    prevent_initial_call=True
)
def handle_xsell_dashboard_navigation(back_clicks, session_data):
    if back_clicks and session_data.get("page") == "x-sell-dash":
        return {"page": "dashboard", "logged_in": True}
    return session_data

# 5. Separate callback for mandates page navigation
@app.callback(
    Output("session-store", "data", allow_duplicate=True),
    [Input("back-to-dashboard-mandates", "n_clicks")],
    [State("session-store", "data")],
    prevent_initial_call=True
)
def handle_mandates_navigation(back_clicks, session_data):
    if back_clicks and session_data.get("page") == "new-mandates":
        return {"page": "dashboard", "logged_in": True}
    return session_data

# 6. Cross-sell page navigation
@app.callback(
    Output("session-store", "data", allow_duplicate=True),
    [Input("back-to-dashboard", "n_clicks"),
     Input("back-to-dashboard-upload", "n_clicks")],
    [State("session-store", "data")],
    prevent_initial_call=True
)
def handle_cross_sell_navigation(back1, back2, session_data):
    ctx = callback_context
    if ctx.triggered and session_data.get("page") == "cross-sell":
        return {"page": "main-dashboard", "logged_in": True}
    return session_data

# 7. Proposal page navigation
@app.callback(
    Output("session-store", "data", allow_duplicate=True),
    [Input("back-to-dashboard-prop", "n_clicks")],
    [State("session-store", "data")],
    prevent_initial_call=True
)
def handle_proposal_navigation(back_prop, session_data):
    if back_prop and session_data.get("page") == "proposal":
        return {"page": "main-dashboard", "logged_in": True}
    return session_data



@app.callback(
    Output("cross-sell-message", "children"),
    [Input("submit-cross-sell", "n_clicks")],
    [State("family-name", "value"), State("product", "value"), State("stage", "value"), State("new-liquidity", "value")]
)
def submit_cross_sell(n_clicks, family_name, product, stage, new_liquidity):
    if not n_clicks:
        return ""
    if not all([family_name, product, stage, new_liquidity]):
        return dbc.Alert("Please fill all fields!", color="warning")

    user_id = current_user.get('UserID', '')
    if user_id:
        single_record = pd.DataFrame({
            'Timestamp': [datetime.now().strftime('%Y-%m-%d %H:%M:%S')],
            'Family_Name': [family_name],
            'Product': [product],
            'Stage': [stage],
            'New_Liquidity': [new_liquidity]
        })
        save_bulk_cross_sell_data(user_id, single_record)
        return dbc.Alert([
            html.H6("Success!", className="alert-heading"),
            html.P("Cross sell data submitted successfully!", className="mb-0")
        ], color="success")
    else:
        return dbc.Alert("User session expired. Please login again.", color="warning")


@app.callback(
    [Output("upload-status", "children"), Output("upload-message", "children")],
    [Input("upload-excel", "contents")],
    [State("upload-excel", "filename")]
)
def handle_excel_upload(contents, filename):
    if contents is None:
        return "", ""

    df, error = parse_excel_upload(contents, filename)
    if error:
        return dbc.Alert(f"Upload Failed: {error}", color="danger"), ""

    if df is not None and not df.empty:
        user_id = current_user.get('UserID', '')
        if user_id:
            try:
                records_count = save_bulk_cross_sell_data(user_id, df)
                preview_html = html.Div([
                    dbc.Alert([
                        html.H6("Upload Successful!", className="alert-heading"),
                        html.P(f"Successfully uploaded {records_count} records from {filename}", className="mb-0")
                    ], color="success"),
                    html.H6("Data Preview", className="mt-4 mb-3"),
                    html.Div([
                        html.Table([
                            html.Thead([html.Tr([html.Th(col, style={"color": "#ffffff"}) for col in df.columns])]),
                            html.Tbody([html.Tr([html.Td(str(df.iloc[i][col])) for col in df.columns]) for i in
                                        range(min(3, len(df)))])
                        ], className="table table-striped"),
                        html.Small(f"Showing first 3 of {len(df)} records" if len(
                            df) > 3 else f"Showing all {len(df)} records", className="text-muted")
                    ], style={"maxHeight": "300px", "overflowY": "auto"})
                ])
                return preview_html, ""
            except Exception as e:
                return dbc.Alert(f"Error saving data: {str(e)}", color="danger"), ""
        else:
            return dbc.Alert("User session expired. Please login again.", color="warning"), ""

    return dbc.Alert("No valid data found in the uploaded file.", color="warning"), ""


@app.callback(
    Output("proposal-message", "children"),
    [Input("submit-proposal", "n_clicks")],
    [State("prop-family-name", "value"), State("prop-banker", "value"), State("prop-date", "value"),
     State("prop-amount", "value"), State("prop-risk-profile", "value"), State("prop-status", "value")]
)
def submit_proposal(n_clicks, family_name, banker, date, amount, risk_profile, status):
    if not n_clicks:
        return ""

    basic_fields = [family_name, banker, date, amount, risk_profile, status]
    if not all(basic_fields):
        return dbc.Alert("Please fill all proposal fields!", color="warning")

    user_id = current_user.get('UserID', '')
    if not user_id:
        return dbc.Alert("User session expired. Please login again.", color="warning")

    try:
        proposal_data = {
            'Timestamp': [datetime.now().strftime('%Y-%m-%d %H:%M:%S')],
            'Family_Name': [family_name],
            'Banker': [banker],
            'Date': [date],
            'Amount': [amount],
            'Risk_Profile': [risk_profile],
            'Status': [status]
        }

        # Save proposal data
        filename = f'data/user_data/{user_id}_proposals.xlsx'
        if os.path.exists(filename):
            existing_df = pd.read_excel(filename)
            combined_df = pd.concat([existing_df, pd.DataFrame(proposal_data)], ignore_index=True)
        else:
            combined_df = pd.DataFrame(proposal_data)
        combined_df.to_excel(filename, index=False)

        success_message = dbc.Alert([
            html.H6("Proposal Submitted Successfully!", className="alert-heading"),
            html.P(f"Proposal for {family_name} has been created.", className="mb-0"),
            html.P(f"Status: {status}", className="mb-0")
        ], color="success")

        return success_message

    except Exception as e:
        return dbc.Alert(f"Error saving proposal: {str(e)}", color="danger")


@app.callback(
    Output("crosssell-table-container", "children"),
    [Input("page-content", "children")],
    prevent_initial_call=True
)
def update_crosssell_table(page_content):
    _, crosssell_df, _ = load_dashboard_data()

    if crosssell_df.empty:
        return html.P("No cross-sell data available. Please add sample data files.", className="text-muted")

    table_data = []
    for _, row in crosssell_df.iterrows():
        table_data.append(
            html.Tr([
                html.Td(row.get('Quarter', 'N/A')),
                html.Td(row.get('Family_Name', 'N/A')),
                html.Td(row.get('Product', 'N/A')),
                html.Td(row.get('Stage', 'N/A')),
                html.Td(f"â‚¹{row.get('Amount_Cr', 0)} Cr"),
                html.Td(row.get('Banker', 'N/A'))
            ])
        )

    table = html.Table([
        html.Thead([
            html.Tr([
                html.Th("Quarter", style={"color": "#ffffff"}),
                html.Th("Family Name", style={"color": "#ffffff"}),
                html.Th("Product", style={"color": "#ffffff"}),
                html.Th("Stage", style={"color": "#ffffff"}),
                html.Th("Amount", style={"color": "#ffffff"}),
                html.Th("Banker", style={"color": "#ffffff"})
            ])
        ]),
        html.Tbody(table_data)
    ], className="table table-striped")

    return table


@app.callback(
    Output("mandate-summary-stats", "children"),
    [Input("page-content", "children")],
    prevent_initial_call=True
)
def update_mandate_stats(page_content):
    _, _, mandate_df = load_dashboard_data()

    if mandate_df.empty:
        return html.P("No mandate data available. Please add sample data files.", className="text-muted")

    total_pitched = mandate_df[mandate_df['Stage'] == 'Pitched']['Count'].sum()
    total_converted = mandate_df[mandate_df['Stage'] == 'Converted']['Count'].sum()
    conversion_rate = round((total_converted / total_pitched * 100), 1) if total_pitched > 0 else 0

    return html.Div([
        dbc.Row([
            dbc.Col([
                html.H5(f"Total Pitched: {total_pitched}", style={"color": "var(--maroon)"})
            ], lg=3),
            dbc.Col([
                html.H5(f"Total Converted: {total_converted}", style={"color": "var(--gold-dark)"})
            ], lg=3),
            dbc.Col([
                html.H5(f"Conversion Rate: {conversion_rate}%", style={"color": "var(--charcoal-gray)"})
            ], lg=3)
        ])
    ])


if __name__ == "__main__":

    app.run(debug=True, port=8050)
